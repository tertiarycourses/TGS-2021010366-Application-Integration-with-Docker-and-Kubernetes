#!/usr/bin/env python3
"""
Application Integration with Docker and Kubernetes
Day 1 — Docker Fundamentals  (target ~62 slides)

Design: matches "Agentic AI Applications with Claude Code - v12.pptx"
Content: KillerCoda labs (GitHub tertiarycourses/TGS-2021010366-...)
Trainer: Mohan Pothula

Run:  python build_courseware_v3.py
Out:  Docker_Kubernetes_Day1_v3.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette (Cloud Native Pro) ───────────────────────────────────────────────
C_DARK    = RGBColor(0x0F, 0x17, 0x2A)  # deep navy   — dark slide bg + body text
C_CREAM   = RGBColor(0xF8, 0xFA, 0xFC)  # near white  — content slide bg
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_MUTED   = RGBColor(0x64, 0x74, 0x8B)  # slate 500   — secondary text
C_MUTED2  = RGBColor(0x94, 0xA3, 0xB8)  # slate 400   — subtext on dark bg
C_FOOTER  = RGBColor(0x71, 0x85, 0x96)  # blue-gray   — footer on both bg types
C_NAVY    = RGBColor(0x0F, 0x17, 0x2A)  # code text   — same as dark
C_BLUE    = RGBColor(0x24, 0x96, 0xED)  # Docker blue — links / KillerCoda
C_CODEBG  = RGBColor(0xF1, 0xF5, 0xF9)  # slate 100   — code block bg
C_AMBER   = RGBColor(0xFF, 0xF7, 0xED)  # knowledge check bg
C_LTBLUE  = RGBColor(0xEF, 0xF6, 0xFF)  # assessment bg

# Brand accent colours (C_ORANGE alias retained for diagram-function compat)
C_ORANGE  = RGBColor(0x24, 0x96, 0xED)  # Docker blue (C_ORANGE kept for compat)
C_HEADER  = RGBColor(0x24, 0x96, 0xED)  # Docker content header band
C_K8S_HDR = RGBColor(0x32, 0x6C, 0xE5)  # Kubernetes content header band
C_K8S     = RGBColor(0x32, 0x6C, 0xE5)  # Kubernetes accent
C_LAB_CODE_BG = RGBColor(0x0D, 0x1B, 0x3E)  # dark navy — lab command blocks
C_TEAL_STEP   = RGBColor(0x00, 0xBC, 0xD4)   # teal — lab step number indicators
C_SUMMARY_BAR = RGBColor(0x21, 0x96, 0xF3)   # blue — outline summary bar

# ── Dimensions ─────────────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.50)

COURSE   = "Application Integration with Docker and Kubernetes"
TGS      = "TGS-2021010366"
KC_BASE  = "https://killercoda.com/tertiary-labs/course/killercoda"

# ── Layout constants ───────────────────────────────────────────────────────────
LABEL_Y   = Inches(0.24)   # section label top (inside 1.50" header band)
TITLE_Y   = Inches(0.60)   # title top (inside header band)
BODY_Y    = Inches(1.72)   # body start (below header band)
CONTENT_MAX = Inches(6.40) # body must not go below this
KC_NOTE_Y = Inches(6.55)   # KillerCoda URL row
FOOTER_Y  = Inches(7.08)   # footer top

# ── Per-line height estimates (used to track cur_y in add_body_lines) ──────────
LH_T  = Inches(0.40)   # body text
LH_B  = Inches(0.40)   # bold
LH_I  = Inches(0.38)   # italic
LH_S  = Inches(0.34)   # small / muted
LH_H  = Inches(0.50)   # mini-heading
LH_BL = Inches(0.22)   # blank spacer
CODE_LINE_H = Inches(0.225)  # per code line (9pt Courier)
CODE_PAD    = Inches(0.16)   # top+bottom padding inside code block


# ══════════════════════════════════════════════════════════════════════════════
#  Low-level helpers
# ══════════════════════════════════════════════════════════════════════════════

def _rgb_to_hex(c):
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _no_bullet(para):
    pPr = para._p.get_or_add_pPr()
    etree.SubElement(pPr, qn("a:buNone"))


def add_textbox(slide, text, x, y, w, h,
                font="Calibri", size=16, bold=False, italic=False,
                color=C_DARK, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    _no_bullet(p)
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    if color:
        f.color.rgb = color
    return txBox


def add_rect(slide, x, y, w, h, fill, no_line=True):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    sp = shape._element
    sp_pr = sp.find(qn("p:spPr"))
    solid = etree.SubElement(sp_pr, qn("a:solidFill"))
    srgb = etree.SubElement(solid, qn("a:srgbClr"))
    srgb.set("val", _rgb_to_hex(fill))
    if no_line:
        ln = etree.SubElement(sp_pr, qn("a:ln"))
        etree.SubElement(ln, qn("a:noFill"))
    shape.text = ""
    return shape


def add_circle(slide, cx_in, cy_in, r_in, fill, alpha=15):
    shape = slide.shapes.add_shape(
        1, Inches(cx_in - r_in), Inches(cy_in - r_in),
        Inches(r_in * 2), Inches(r_in * 2))
    sp = shape._element
    sp_pr = sp.find(qn("p:spPr"))
    pg = sp_pr.find(qn("a:prstGeom"))
    if pg is not None:
        pg.set("prst", "ellipse")
    fill_el = sp_pr.find(qn("a:solidFill"))
    if fill_el is None:
        fill_el = etree.SubElement(sp_pr, qn("a:solidFill"))
    fill_el.clear()
    srgb = etree.SubElement(fill_el, qn("a:srgbClr"))
    srgb.set("val", _rgb_to_hex(fill))
    etree.SubElement(srgb, qn("a:alpha")).set("val", str(alpha * 1000))
    ln = sp_pr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(sp_pr, qn("a:ln"))
    etree.SubElement(ln, qn("a:noFill"))
    shape.text = ""


# ── Page-level layout helpers ──────────────────────────────────────────────────

def blank_slide(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    return slide


def add_footer(slide, page_num):
    add_textbox(slide, "© Tertiary Infotech Pte. Ltd.",
                Inches(0.50), FOOTER_Y, Inches(4.50), Inches(0.30),
                font="Calibri", size=9, color=C_FOOTER)
    add_textbox(slide, COURSE,
                Inches(4.00), FOOTER_Y, Inches(5.50), Inches(0.30),
                font="Calibri", size=9, color=C_FOOTER, align=PP_ALIGN.CENTER)
    add_textbox(slide, str(page_num),
                Inches(12.43), FOOTER_Y, Inches(0.50), Inches(0.30),
                font="Calibri", size=9, color=C_FOOTER, align=PP_ALIGN.RIGHT)


def add_section_label(slide, label, color=C_ORANGE):
    add_textbox(slide, label, Inches(0.60), LABEL_Y, Inches(11.0), Inches(0.32),
                font="Calibri", size=12, bold=True, color=color)


def add_main_title(slide, text, color=C_DARK, y=None, h=Inches(0.85)):
    y = y or TITLE_Y
    add_textbox(slide, text, Inches(0.60), y, Inches(12.10), h,
                font="Cambria", size=30, bold=True, color=color)


def kc_note(slide, scenario_path, y=None):
    y = y or KC_NOTE_Y
    add_textbox(slide, f"KillerCoda  ->  {KC_BASE}/{scenario_path}",
                Inches(0.70), y, Inches(12.0), Inches(0.35),
                font="Calibri", size=11, italic=True, color=C_BLUE)


# ── Body content renderer ──────────────────────────────────────────────────────

def add_body_lines(slide, lines,
                   x=Inches(0.70), y=BODY_Y, w=Inches(12.0)):
    """
    Render a list of (type, text) items onto a slide.
    Types:
        'T'   Calibri 16pt regular
        'B'   Calibri 16pt bold
        'I'   Calibri 14pt italic orange
        'S'   Calibri 12pt muted
        'H'   Cambria 20pt bold (mini-heading)
        'C'   Courier New 9pt on code-bg — text may contain \n
        'BL'  blank spacer
    Stops adding content once cur_y reaches CONTENT_MAX.
    """
    cur_y = y
    for item in lines:
        if cur_y >= CONTENT_MAX:
            break
        typ = item[0] if isinstance(item, (list, tuple)) else 'T'
        text = item[1] if isinstance(item, (list, tuple)) else item

        if typ == 'BL':
            cur_y += LH_BL
            continue

        if typ == 'C':
            n = text.count('\n') + 1
            block_h = CODE_LINE_H * n + CODE_PAD
            # Clamp to available space
            avail = CONTENT_MAX - cur_y
            if avail < Inches(0.35):
                break
            block_h = min(block_h, avail)
            add_rect(slide, x - Inches(0.08), cur_y - Inches(0.05),
                     w + Inches(0.16), block_h, C_CODEBG)
            add_textbox(slide, text, x, cur_y, w, block_h,
                        font="Courier New", size=9, color=C_NAVY, wrap=True)
            cur_y += block_h + Inches(0.10)
        elif typ == 'H':
            h_h = LH_H
            add_textbox(slide, text, x, cur_y, w, h_h,
                        font="Cambria", size=20, bold=True, color=C_DARK)
            cur_y += h_h
        elif typ == 'B':
            add_textbox(slide, text, x, cur_y, w, LH_B,
                        font="Calibri", size=16, bold=True, color=C_DARK)
            cur_y += LH_B
        elif typ == 'I':
            add_textbox(slide, text, x, cur_y, w, LH_I,
                        font="Calibri", size=14, italic=True, color=C_ORANGE)
            cur_y += LH_I
        elif typ == 'S':
            add_textbox(slide, text, x, cur_y, w, LH_S,
                        font="Calibri", size=12, color=C_MUTED)
            cur_y += LH_S
        else:  # 'T'
            add_textbox(slide, text, x, cur_y, w, LH_T,
                        font="Calibri", size=16, color=C_DARK, wrap=True)
            cur_y += LH_T
    return cur_y  # return final y so caller can add kc_note below


def add_lab_body_lines(slide, lines, x=Inches(0.70), y=BODY_Y, w=Inches(12.0)):
    """Like add_body_lines but renders 'C' blocks as dark navy with teal step indicators."""
    step_num = 0
    cur_y = y
    for item in lines:
        if cur_y >= CONTENT_MAX:
            break
        typ = item[0] if isinstance(item, (list, tuple)) else 'T'
        text = item[1] if isinstance(item, (list, tuple)) else item

        if typ == 'BL':
            cur_y += LH_BL
            continue

        if typ == 'C':
            step_num += 1
            n = text.count('\n') + 1
            block_h = CODE_LINE_H * n + CODE_PAD
            avail = CONTENT_MAX - cur_y
            if avail < Inches(0.35):
                break
            block_h = min(block_h, avail)
            # Dark navy background
            add_rect(slide, x - Inches(0.08), cur_y - Inches(0.05),
                     w + Inches(0.16), block_h, C_LAB_CODE_BG)
            # Teal step-indicator strip on left edge
            add_rect(slide, x - Inches(0.08), cur_y - Inches(0.05),
                     Inches(0.28), block_h, C_TEAL_STEP)
            # Step number in white on teal strip
            add_textbox(slide, str(step_num),
                        x - Inches(0.08), cur_y - Inches(0.02),
                        Inches(0.28), block_h,
                        font="Calibri", size=11, bold=True,
                        color=C_WHITE, align=PP_ALIGN.CENTER)
            # Code text in white, indented past the teal strip
            add_textbox(slide, text, x + Inches(0.22), cur_y,
                        w - Inches(0.22), block_h,
                        font="Courier New", size=9, color=C_WHITE, wrap=True)
            cur_y += block_h + Inches(0.10)
        elif typ == 'H':
            add_textbox(slide, text, x, cur_y, w, LH_H,
                        font="Cambria", size=20, bold=True, color=C_DARK)
            cur_y += LH_H
        elif typ == 'B':
            add_textbox(slide, text, x, cur_y, w, LH_B,
                        font="Calibri", size=16, bold=True, color=C_DARK)
            cur_y += LH_B
        elif typ == 'I':
            add_textbox(slide, text, x, cur_y, w, LH_I,
                        font="Calibri", size=14, italic=True, color=C_ORANGE)
            cur_y += LH_I
        elif typ == 'S':
            add_textbox(slide, text, x, cur_y, w, LH_S,
                        font="Calibri", size=12, color=C_MUTED)
            cur_y += LH_S
        else:  # 'T'
            add_textbox(slide, text, x, cur_y, w, LH_T,
                        font="Calibri", size=16, color=C_DARK, wrap=True)
            cur_y += LH_T
    return cur_y


# ── Slide-type factories ───────────────────────────────────────────────────────

def content_slide(prs, label, title, k8s=False):
    acc = C_K8S_HDR if k8s else C_HEADER
    slide = blank_slide(prs)
    set_bg(slide, C_CREAM)
    # Left accent bar (full height)
    add_rect(slide, Inches(0), Inches(0), Inches(0.10), Inches(7.50), acc)
    # Header band
    add_rect(slide, Inches(0.10), Inches(0), Inches(13.23), Inches(1.50), acc)
    # Thin body accent bar — adds visual depth below header
    add_rect(slide, Inches(0.60), Inches(1.55), Inches(0.05), Inches(5.30),
             C_SUMMARY_BAR)
    # Section label inside header band
    add_textbox(slide, label,
                Inches(0.25), LABEL_Y, Inches(12.80), Inches(0.32),
                font="Calibri", size=11, bold=True, color=C_WHITE)
    # Slide title inside header band
    add_textbox(slide, title,
                Inches(0.22), TITLE_Y, Inches(12.80), Inches(0.80),
                font="Cambria", size=26, bold=True, color=C_WHITE)
    return slide


def section_divider(prs, label, title, subtitle="", subtopics=None, page=0, k8s=False):
    """Modern dark section divider with brand accent bar and corner block."""
    C_LIGHT = RGBColor(0xCB, 0xD5, 0xE1)  # slate 300 — topic list text
    acc = C_K8S if k8s else C_HEADER
    slide = blank_slide(prs)
    set_bg(slide, C_DARK)
    # Thick left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.55), Inches(7.50), acc)
    # Top-right corner accent block
    add_rect(slide, Inches(10.60), Inches(0), Inches(2.73), Inches(1.20), acc)
    # Section label
    add_textbox(slide, label,
                Inches(0.80), Inches(2.00), Inches(8.50), Inches(0.50),
                font="Calibri", size=14, bold=True, color=acc)
    # Section title
    add_textbox(slide, title,
                Inches(0.75), Inches(2.55), Inches(8.80), Inches(1.80),
                font="Cambria", size=42, bold=True, color=C_WHITE)
    # Optional subtitle
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.80), Inches(4.52), Inches(8.50), Inches(0.50),
                    font="Calibri", size=16, color=C_LIGHT)
    # Numbered topic list on right
    if subtopics:
        for i, topic in enumerate(subtopics):
            ny = Inches(2.00 + i * 0.68)
            if ny > Inches(6.60):
                break
            # Horizontal separator between items (except before the first)
            if i > 0:
                add_rect(slide, Inches(9.60), ny - Inches(0.10),
                         Inches(3.15), Inches(0.015),
                         RGBColor(0x4A, 0x5C, 0x6A))
            add_textbox(slide, str(i + 1),
                        Inches(9.60), ny, Inches(0.50), Inches(0.46),
                        font="Cambria", size=18, bold=True, color=acc)
            add_textbox(slide, topic,
                        Inches(10.15), ny, Inches(2.80), Inches(0.46),
                        font="Calibri", size=18, color=C_LIGHT)
    add_footer(slide, page)
    return page


def lab_header_slide(prs, lab_num, title, scenario, page, k8s=False):
    acc = C_K8S if k8s else C_HEADER
    slide = blank_slide(prs)
    set_bg(slide, C_DARK)
    # Left coloured panel
    add_rect(slide, Inches(0), Inches(0), Inches(3.80), Inches(7.50), acc)
    # Thin top accent line on dark right area
    add_rect(slide, Inches(3.80), Inches(0), Inches(9.53), Inches(0.12), acc)
    # Lab badge in left panel
    add_textbox(slide, lab_num,
                Inches(0.15), Inches(2.60), Inches(3.50), Inches(0.55),
                font="Calibri", size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "HANDS-ON LAB",
                Inches(0.15), Inches(3.18), Inches(3.50), Inches(0.35),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    # Lab title on dark right area
    add_textbox(slide, title,
                Inches(4.10), Inches(1.60), Inches(9.00), Inches(1.80),
                font="Cambria", size=36, bold=True, color=C_WHITE, wrap=True)
    # KillerCoda URL
    add_textbox(slide, f"KillerCoda  →  {KC_BASE}/{scenario}",
                Inches(4.10), Inches(3.60), Inches(9.00), Inches(0.40),
                font="Calibri", size=14, italic=True, color=C_BLUE)
    # Instruction hint
    add_textbox(slide, "Follow the steps on KillerCoda — commands run in the browser terminal.",
                Inches(4.10), Inches(4.10), Inches(9.00), Inches(0.40),
                font="Calibri", size=13, color=C_MUTED2)
    add_footer(slide, page)
    return page


def two_column_slide(prs, label, title,
                     left_heading, left_lines,
                     right_heading, right_lines, page, k8s=False):
    slide = content_slide(prs, label, title, k8s=k8s)
    col_w = Inches(5.80)
    lx, rx = Inches(0.22), Inches(6.90)
    by = Inches(1.72)
    add_rect(slide, Inches(6.66), by, Inches(0.03), Inches(4.80), C_MUTED)
    add_textbox(slide, left_heading, lx, by, col_w, Inches(0.45),
                font="Cambria", size=18, bold=True, color=C_DARK)
    add_body_lines(slide, left_lines, x=lx, y=by + Inches(0.48), w=col_w)
    add_textbox(slide, right_heading, rx, by, col_w, Inches(0.45),
                font="Cambria", size=18, bold=True, color=C_DARK)
    add_body_lines(slide, right_lines, x=rx, y=by + Inches(0.48), w=col_w)
    add_footer(slide, page)
    return page


# ══════════════════════════════════════════════════════════════════════════════
#  Architecture diagram slides
# ══════════════════════════════════════════════════════════════════════════════

def slide_vm_vs_container_diagram(prs, page):
    """Visual VM vs Container stack diagram."""
    C_APPBOX = RGBColor(0xDA, 0xEA, 0xFE)
    C_OSBOX  = RGBColor(0x8A, 0x86, 0x7D)
    C_HW     = RGBColor(0x3D, 0x3C, 0x38)

    slide = content_slide(prs, "DOCKER — OVERVIEW", "VM vs Container Architecture")

    # Column header bands
    add_rect(slide, Inches(0.60), Inches(1.90), Inches(5.90), Inches(0.43), RGBColor(0xEC, 0xEB, 0xE6))
    add_textbox(slide, "Virtual Machine Model",
                Inches(0.60), Inches(1.93), Inches(5.90), Inches(0.37),
                font="Cambria", size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(6.80), Inches(1.90), Inches(5.90), Inches(0.43), RGBColor(0xEC, 0xEB, 0xE6))
    add_textbox(slide, "Container Model",
                Inches(6.80), Inches(1.93), Inches(5.90), Inches(0.37),
                font="Cambria", size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # Centre divider
    add_rect(slide, Inches(6.60), Inches(1.90), Inches(0.04), Inches(5.00), C_MUTED)
    add_textbox(slide, "vs", Inches(6.25), Inches(4.10), Inches(0.75), Inches(0.42),
                font="Cambria", size=18, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # ── VM side (3 columns) ──────────────────────────────────────────────────
    vm_xs = [0.65, 2.42, 4.19]
    vm_w  = 1.72
    for xi, app in zip(vm_xs, ["App 1", "App 2", "App 3"]):
        add_rect(slide, Inches(xi), Inches(2.50), Inches(vm_w), Inches(0.54), C_APPBOX)
        add_textbox(slide, app,
                    Inches(xi + 0.05), Inches(2.55), Inches(vm_w - 0.10), Inches(0.40),
                    font="Calibri", size=12, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(xi), Inches(3.10), Inches(vm_w), Inches(0.54), C_OSBOX)
        add_textbox(slide, "Guest OS",
                    Inches(xi + 0.05), Inches(3.15), Inches(vm_w - 0.10), Inches(0.40),
                    font="Calibri", size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.65), Inches(3.72), Inches(5.77), Inches(0.54), C_ORANGE)
    add_textbox(slide, "Hypervisor",
                Inches(0.65), Inches(3.77), Inches(5.77), Inches(0.40),
                font="Calibri", size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.65), Inches(4.34), Inches(5.77), Inches(0.52), C_HW)
    add_textbox(slide, "Hardware",
                Inches(0.65), Inches(4.39), Inches(5.77), Inches(0.38),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Each VM carries a full Guest OS — heavier, slower to start",
                Inches(0.65), Inches(5.00), Inches(5.77), Inches(0.36),
                font="Calibri", size=11, italic=True, color=C_MUTED)

    # ── Container side (3 columns) ───────────────────────────────────────────
    ct_xs = [6.88, 8.65, 10.42]
    ct_w  = 1.72
    for xi, app in zip(ct_xs, ["App 1", "App 2", "App 3"]):
        add_rect(slide, Inches(xi), Inches(2.50), Inches(ct_w), Inches(0.54), C_APPBOX)
        add_textbox(slide, app,
                    Inches(xi + 0.05), Inches(2.55), Inches(ct_w - 0.10), Inches(0.40),
                    font="Calibri", size=12, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(6.88), Inches(3.10), Inches(5.77), Inches(0.54), C_ORANGE)
    add_textbox(slide, "Docker Engine",
                Inches(6.88), Inches(3.15), Inches(5.77), Inches(0.40),
                font="Calibri", size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(6.88), Inches(3.72), Inches(5.77), Inches(0.54), C_OSBOX)
    add_textbox(slide, "Host OS",
                Inches(6.88), Inches(3.77), Inches(5.77), Inches(0.40),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(6.88), Inches(4.34), Inches(5.77), Inches(0.52), C_HW)
    add_textbox(slide, "Hardware",
                Inches(6.88), Inches(4.39), Inches(5.77), Inches(0.38),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Containers share the Host OS kernel — lightweight, fast to start",
                Inches(6.88), Inches(5.00), Inches(5.77), Inches(0.36),
                font="Calibri", size=11, italic=True, color=C_MUTED)

    add_footer(slide, page)
    return page


def slide_docker_engine_diagram(prs, page):
    """Visual Docker Engine: Client → Daemon → Registry."""
    C_BOXBG = RGBColor(0xDA, 0xEA, 0xFE)

    slide = content_slide(prs, "DOCKER — ARCHITECTURE", "Docker Engine: Client → Daemon → Registry")

    # Docker Client box
    add_rect(slide, Inches(0.60), Inches(2.10), Inches(3.20), Inches(2.90), C_BOXBG)
    add_textbox(slide, "Docker Client",
                Inches(0.65), Inches(2.14), Inches(3.10), Inches(0.40),
                font="Cambria", size=14, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    for i, cmd in enumerate(["docker run", "docker build", "docker pull", "docker push", "docker ps"]):
        add_textbox(slide, cmd,
                    Inches(0.90), Inches(2.68 + i * 0.38), Inches(2.70), Inches(0.34),
                    font="Courier New", size=10, color=C_NAVY)

    # Arrow
    add_textbox(slide, "→", Inches(3.88), Inches(3.32), Inches(0.64), Inches(0.44),
                font="Calibri", size=28, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "REST API",
                Inches(3.80), Inches(3.80), Inches(0.80), Inches(0.28),
                font="Calibri", size=9, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Docker Daemon box
    add_rect(slide, Inches(4.60), Inches(1.88), Inches(4.10), Inches(3.70), C_DARK)
    add_textbox(slide, "Docker Daemon  (dockerd)",
                Inches(4.65), Inches(1.92), Inches(4.00), Inches(0.40),
                font="Cambria", size=13, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    obj_data = [
        ("Images",     RGBColor(0x2E, 0x74, 0xB5)),
        ("Containers", RGBColor(0x10, 0x78, 0x54)),
        ("Volumes",    RGBColor(0x7B, 0x54, 0xA0)),
        ("Networks",   RGBColor(0xC0, 0x55, 0x2A)),
    ]
    for idx, (lbl, clr) in enumerate(obj_data):
        bx = 4.75 + (idx % 2) * 1.98
        by = 2.46 + (idx // 2) * 1.18
        add_rect(slide, Inches(bx), Inches(by), Inches(1.83), Inches(0.95), clr)
        add_textbox(slide, lbl,
                    Inches(bx + 0.05), Inches(by + 0.28), Inches(1.73), Inches(0.38),
                    font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Arrow
    add_textbox(slide, "→", Inches(8.78), Inches(3.32), Inches(0.64), Inches(0.44),
                font="Calibri", size=28, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # Registry box
    add_rect(slide, Inches(9.50), Inches(2.10), Inches(3.20), Inches(2.90), C_BOXBG)
    add_textbox(slide, "Registry",
                Inches(9.55), Inches(2.14), Inches(3.10), Inches(0.40),
                font="Cambria", size=14, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)
    for i, lbl in enumerate(["Docker Hub (public)", "AWS ECR", "Azure ACR", "Private registry"]):
        add_textbox(slide, lbl,
                    Inches(9.72), Inches(2.68 + i * 0.43), Inches(2.80), Inches(0.38),
                    font="Calibri", size=11, color=C_NAVY)

    add_textbox(slide,
                "Flow:  docker run  →  Client sends REST call to Daemon  →  Daemon pulls image from Registry  →  Creates & starts Container",
                Inches(0.60), Inches(5.72), Inches(12.10), Inches(0.36),
                font="Calibri", size=11, italic=True, color=C_MUTED)

    add_footer(slide, page)
    return page


def slide_k8s_arch_diagram(prs, page):
    """Visual Kubernetes Cluster Architecture: Control Plane + Worker Nodes."""
    C_CP  = RGBColor(0x1F, 0x49, 0x7D)
    C_CPi = RGBColor(0x2E, 0x74, 0xB5)
    C_WN  = RGBColor(0x1E, 0x4D, 0x3A)
    C_WNi = RGBColor(0x27, 0x6B, 0x52)
    C_POD = RGBColor(0x10, 0x78, 0x54)

    slide = content_slide(prs, "KUBERNETES — ARCHITECTURE",
                          "Kubernetes Cluster: Control Plane & Worker Nodes")

    # Cluster outline
    add_rect(slide, Inches(0.55), Inches(1.88), Inches(12.22), Inches(4.55),
             RGBColor(0xF0, 0xF4, 0xFF))
    add_textbox(slide, "Kubernetes Cluster",
                Inches(0.60), Inches(1.90), Inches(12.00), Inches(0.26),
                font="Calibri", size=9, bold=True, color=C_MUTED)

    # Control Plane
    add_rect(slide, Inches(0.65), Inches(2.22), Inches(3.90), Inches(3.90), C_CP)
    add_textbox(slide, "Control Plane",
                Inches(0.70), Inches(2.25), Inches(3.80), Inches(0.38),
                font="Cambria", size=13, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    cp_items = [("API Server", "Entry point — REST API"),
                ("etcd", "Cluster state store"),
                ("Scheduler", "Assigns Pods to Nodes"),
                ("Controller Mgr", "Reconciles desired state")]
    for i, (lbl, desc) in enumerate(cp_items):
        iy = 2.72 + i * 0.84
        add_rect(slide, Inches(0.80), Inches(iy), Inches(3.60), Inches(0.68), C_CPi)
        add_textbox(slide, lbl,
                    Inches(0.85), Inches(iy + 0.03), Inches(3.50), Inches(0.28),
                    font="Calibri", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc,
                    Inches(0.85), Inches(iy + 0.36), Inches(3.50), Inches(0.24),
                    font="Calibri", size=9, color=C_MUTED2, align=PP_ALIGN.CENTER)

    # Arrow
    add_textbox(slide, "→", Inches(4.65), Inches(4.10), Inches(0.62), Inches(0.44),
                font="Calibri", size=28, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # Worker nodes (2)
    for ni, nx in enumerate([5.42, 9.00]):
        add_rect(slide, Inches(nx), Inches(2.22), Inches(3.25), Inches(3.90), C_WN)
        add_textbox(slide, f"Worker Node {ni + 1}",
                    Inches(nx + 0.05), Inches(2.25), Inches(3.15), Inches(0.38),
                    font="Cambria", size=12, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

        for i, (lbl, desc) in enumerate([("kubelet", "Node agent — runs Pods"),
                                          ("kube-proxy", "Network & routing rules"),
                                          ("Pods", "App containers run here")]):
            iy = 2.72 + i * 1.06
            fill = C_POD if lbl == "Pods" else C_WNi
            add_rect(slide, Inches(nx + 0.10), Inches(iy), Inches(3.05), Inches(0.88), fill)
            add_textbox(slide, lbl,
                        Inches(nx + 0.15), Inches(iy + 0.04), Inches(2.95), Inches(0.28),
                        font="Calibri", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            add_textbox(slide, desc,
                        Inches(nx + 0.15), Inches(iy + 0.46), Inches(2.95), Inches(0.26),
                        font="Calibri", size=9, color=C_MUTED2, align=PP_ALIGN.CENTER)

    add_textbox(slide,
                "kubectl  →  sends commands to API Server  →  which schedules & manages Pods on Worker Nodes",
                Inches(0.60), Inches(6.60), Inches(12.10), Inches(0.30),
                font="Calibri", size=10, italic=True, color=C_FOOTER)

    add_footer(slide, page)
    return page


def slide_k8s_objects_diagram(prs, page):
    """Visual Service → Deployment → Pods relationship diagram."""
    C_SVC  = C_ORANGE
    C_DEPL = RGBColor(0x2E, 0x74, 0xB5)
    C_POD  = RGBColor(0x10, 0x78, 0x54)
    C_CTR  = RGBColor(0x17, 0x9E, 0x6D)
    C_NS   = RGBColor(0xEC, 0xEB, 0xE6)

    slide = content_slide(prs, "KUBERNETES — OBJECTS", "Service → Deployment → Pods")

    # Service
    add_rect(slide, Inches(0.65), Inches(2.00), Inches(8.10), Inches(0.88), C_SVC)
    add_textbox(slide, "Service",
                Inches(0.70), Inches(2.03), Inches(8.00), Inches(0.36),
                font="Cambria", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Stable network endpoint (ClusterIP / NodePort / LoadBalancer)  •  DNS  •  Load balancing",
                Inches(0.70), Inches(2.40), Inches(8.00), Inches(0.38),
                font="Calibri", size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "↓", Inches(4.30), Inches(2.94), Inches(0.60), Inches(0.38),
                font="Calibri", size=22, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # Deployment
    add_rect(slide, Inches(0.65), Inches(3.38), Inches(8.10), Inches(0.88), C_DEPL)
    add_textbox(slide, "Deployment",
                Inches(0.70), Inches(3.41), Inches(8.00), Inches(0.36),
                font="Cambria", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Manages desired replicas  •  Rolling updates  •  Self-healing via ReplicaSet",
                Inches(0.70), Inches(3.78), Inches(8.00), Inches(0.38),
                font="Calibri", size=10, color=C_MUTED2, align=PP_ALIGN.CENTER)

    for xi in [1.55, 4.30, 7.05]:
        add_textbox(slide, "↓", Inches(xi), Inches(4.32), Inches(0.60), Inches(0.38),
                    font="Calibri", size=22, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # 3 Pods
    for i, nx in enumerate([0.65, 3.22, 5.79]):
        add_rect(slide, Inches(nx), Inches(4.76), Inches(2.46), Inches(1.52), C_POD)
        add_textbox(slide, f"Pod {i + 1}",
                    Inches(nx + 0.05), Inches(4.79), Inches(2.36), Inches(0.30),
                    font="Cambria", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(nx + 0.10), Inches(5.18), Inches(2.26), Inches(0.72), C_CTR)
        add_textbox(slide, "Container",
                    Inches(nx + 0.10), Inches(5.38), Inches(2.26), Inches(0.30),
                    font="Calibri", size=10, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Namespace panel on right
    add_rect(slide, Inches(9.00), Inches(2.00), Inches(3.90), Inches(4.40), C_NS)
    add_textbox(slide, "Namespace",
                Inches(9.05), Inches(2.04), Inches(3.80), Inches(0.36),
                font="Cambria", size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide,
                "Isolates resources by environment\n(dev / staging / prod)\n\n"
                "kubectl create namespace dev\nkubectl get pods -n dev\n\n"
                "Each Namespace has its own:\n  Services, Deployments, Pods",
                Inches(9.10), Inches(2.46), Inches(3.70), Inches(2.80),
                font="Calibri", size=11, color=C_DARK)

    add_footer(slide, page)
    return page


def slide_docker_image_layers_diagram(prs, page):
    """Docker image layer stack — FROM base to CMD, showing build cache."""
    slide = content_slide(prs, "DOCKERFILE — LAYERS", "How Docker Image Layers Stack")

    layers = [
        (RGBColor(0x3D, 0x3C, 0x38), C_WHITE,  'CMD ["python", "app.py"]',               "Runtime command"),
        (RGBColor(0x27, 0x6B, 0x52), C_WHITE,  "COPY . .",                                "Source code (changes often)"),
        (RGBColor(0x10, 0x78, 0x54), C_WHITE,  "RUN pip install -r requirements.txt",     "Dependencies installed"),
        (RGBColor(0x2E, 0x74, 0xB5), C_WHITE,  "COPY requirements.txt .",                 "Dep list (rarely changes)"),
        (RGBColor(0x1F, 0x49, 0x7D), C_WHITE,  "WORKDIR /app",                            "Working directory"),
        (C_ORANGE,                   C_WHITE,  "FROM python:3.11-slim",                   "Base image — cached layer"),
    ]

    box_h, box_gap = 0.52, 0.05
    box_w, lx = 7.50, 0.70
    start_y = 5.80

    for i, (fill, text_color, instruction, desc) in enumerate(layers):
        y = start_y - i * (box_h + box_gap)
        add_rect(slide, Inches(lx), Inches(y), Inches(box_w), Inches(box_h), fill)
        add_textbox(slide, instruction,
                    Inches(lx + 0.15), Inches(y + 0.09), Inches(4.20), Inches(0.34),
                    font="Courier New", size=9, color=text_color)
        add_textbox(slide, desc,
                    Inches(lx + 4.50), Inches(y + 0.10), Inches(2.90), Inches(0.30),
                    font="Calibri", size=10, italic=True, color=text_color)

    add_textbox(slide, "Build  ↑  (bottom to top)",
                Inches(8.50), Inches(3.60), Inches(1.80), Inches(0.34),
                font="Calibri", size=11, bold=True, color=C_ORANGE)

    # Cache rules panel
    cx = 8.50
    add_rect(slide, Inches(cx), Inches(1.90), Inches(4.30), Inches(3.34), RGBColor(0xEC, 0xEB, 0xE6))
    add_textbox(slide, "Build Cache Rules",
                Inches(cx + 0.15), Inches(2.00), Inches(4.00), Inches(0.36),
                font="Cambria", size=14, bold=True, color=C_DARK)
    cache_notes = [
        "Each instruction = one image layer",
        "Unchanged layer → reused from cache",
        "Changed layer → all above it rebuild",
        "Put stable steps first (FROM, deps)",
        "Put changing files last (COPY . .)",
    ]
    for i, note in enumerate(cache_notes):
        add_textbox(slide, note,
                    Inches(cx + 0.18), Inches(2.46 + i * 0.42), Inches(4.00), Inches(0.36),
                    font="Calibri", size=11, color=C_DARK)

    add_footer(slide, page)
    return page


def slide_docker_storage_diagram(prs, page):
    """Named Volume vs Bind Mount with host-container relationship."""
    C_HOST   = RGBColor(0x3D, 0x3C, 0x38)
    C_DOCKER = RGBColor(0x1F, 0x49, 0x7D)
    C_CTR    = RGBColor(0x27, 0x6B, 0x52)

    slide = content_slide(prs, "DOCKER STORAGE", "Named Volume vs Bind Mount — Visual")

    # ── Named Volume (left) ────────────────────────────────────────────────────
    add_textbox(slide, "Named Volume",
                Inches(0.65), Inches(1.90), Inches(5.90), Inches(0.38),
                font="Cambria", size=16, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.65), Inches(2.38), Inches(5.90), Inches(1.00), C_DOCKER)
    add_textbox(slide, "Docker Engine",
                Inches(0.70), Inches(2.43), Inches(5.80), Inches(0.36),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "/var/lib/docker/volumes/my-vol",
                Inches(0.70), Inches(2.80), Inches(5.80), Inches(0.36),
                font="Courier New", size=9, color=C_MUTED2, align=PP_ALIGN.CENTER)
    add_textbox(slide, "↕  mount    -v my-vol:/data",
                Inches(1.00), Inches(3.45), Inches(5.00), Inches(0.36),
                font="Calibri", size=12, color=C_ORANGE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.65), Inches(3.90), Inches(5.90), Inches(0.88), C_CTR)
    add_textbox(slide, "Container  —  reads/writes /data",
                Inches(0.70), Inches(4.20), Inches(5.80), Inches(0.32),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Data persists after docker rm",
                Inches(0.70), Inches(4.92), Inches(5.80), Inches(0.32),
                font="Calibri", size=11, italic=True, color=C_MUTED)
    add_textbox(slide, "Best for: databases, persistent app data",
                Inches(0.70), Inches(5.26), Inches(5.80), Inches(0.32),
                font="Calibri", size=11, color=C_DARK)

    # Centre divider
    add_rect(slide, Inches(6.60), Inches(1.90), Inches(0.04), Inches(3.80), C_MUTED)
    add_textbox(slide, "vs",
                Inches(6.24), Inches(3.55), Inches(0.76), Inches(0.42),
                font="Cambria", size=18, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # ── Bind Mount (right) ─────────────────────────────────────────────────────
    add_textbox(slide, "Bind Mount",
                Inches(6.85), Inches(1.90), Inches(5.90), Inches(0.38),
                font="Cambria", size=16, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(6.85), Inches(2.38), Inches(5.90), Inches(1.00), C_HOST)
    add_textbox(slide, "Host OS",
                Inches(6.90), Inches(2.43), Inches(5.80), Inches(0.36),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "$(pwd)/src  — you control this path",
                Inches(6.90), Inches(2.80), Inches(5.80), Inches(0.36),
                font="Courier New", size=9, color=C_MUTED2, align=PP_ALIGN.CENTER)
    add_textbox(slide, "↕  instant sync    -v $(pwd)/src:/app/src",
                Inches(6.90), Inches(3.45), Inches(5.70), Inches(0.36),
                font="Calibri", size=12, color=C_ORANGE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(6.85), Inches(3.90), Inches(5.90), Inches(0.88), C_CTR)
    add_textbox(slide, "Container  —  reads/writes /app/src",
                Inches(6.90), Inches(4.20), Inches(5.80), Inches(0.32),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Changes visible immediately — both directions",
                Inches(6.90), Inches(4.92), Inches(5.80), Inches(0.32),
                font="Calibri", size=11, italic=True, color=C_MUTED)
    add_textbox(slide, "Best for: live code reload, development",
                Inches(6.90), Inches(5.26), Inches(5.80), Inches(0.32),
                font="Calibri", size=11, color=C_DARK)

    add_footer(slide, page)
    return page


def slide_docker_networking_diagram(prs, page):
    """Default bridge (no DNS) vs Custom bridge (DNS by container name)."""
    C_NET = RGBColor(0x2E, 0x74, 0xB5)
    C_CTR = RGBColor(0x27, 0x6B, 0x52)
    C_BAD = RGBColor(0xCC, 0x44, 0x44)
    C_OK  = RGBColor(0x10, 0x78, 0x54)

    slide = content_slide(prs, "DOCKER NETWORKING", "Default Bridge vs Custom Bridge Network")

    # ── Default Bridge (left) ──────────────────────────────────────────────────
    add_rect(slide, Inches(0.60), Inches(1.90), Inches(5.90), Inches(4.30),
             RGBColor(0xF0, 0xF4, 0xFF))
    add_textbox(slide, "Default Bridge  (docker0)",
                Inches(0.65), Inches(1.93), Inches(5.80), Inches(0.36),
                font="Cambria", size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    for i, (name, ip) in enumerate([("container-a", "172.17.0.2"), ("container-b", "172.17.0.3")]):
        cx = 0.80 + i * 2.95
        add_rect(slide, Inches(cx), Inches(2.46), Inches(2.60), Inches(1.10), C_CTR)
        add_textbox(slide, name,
                    Inches(cx + 0.05), Inches(2.52), Inches(2.50), Inches(0.32),
                    font="Calibri", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, ip,
                    Inches(cx + 0.05), Inches(2.86), Inches(2.50), Inches(0.28),
                    font="Courier New", size=9, color=C_MUTED2, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.80), Inches(3.66), Inches(5.60), Inches(0.40), C_NET)
    add_textbox(slide, "docker0 bridge  (172.17.0.0/16)",
                Inches(0.80), Inches(3.70), Inches(5.60), Inches(0.28),
                font="Calibri", size=11, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "✘ No DNS — containers must use IP addresses",
                Inches(0.65), Inches(4.18), Inches(5.80), Inches(0.32),
                font="Calibri", size=11, italic=True, color=C_BAD)
    add_textbox(slide, "ping 172.17.0.3  (hardcoded — breaks on restart!)",
                Inches(0.65), Inches(4.52), Inches(5.80), Inches(0.32),
                font="Courier New", size=9, color=C_NAVY)
    add_textbox(slide, "IP changes every time container restarts",
                Inches(0.65), Inches(4.88), Inches(5.80), Inches(0.32),
                font="Calibri", size=11, color=C_MUTED)

    # Centre divider
    add_rect(slide, Inches(6.60), Inches(1.90), Inches(0.04), Inches(4.30), C_MUTED)
    add_textbox(slide, "vs",
                Inches(6.24), Inches(3.80), Inches(0.76), Inches(0.42),
                font="Cambria", size=18, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # ── Custom Bridge (right) ──────────────────────────────────────────────────
    add_rect(slide, Inches(6.82), Inches(1.90), Inches(5.90), Inches(4.30),
             RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide, "Custom Bridge  (my-net)",
                Inches(6.87), Inches(1.93), Inches(5.80), Inches(0.36),
                font="Cambria", size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    for i, (name, ip) in enumerate([("app1", "172.18.0.2"), ("app2", "172.18.0.3")]):
        cx = 7.00 + i * 2.95
        add_rect(slide, Inches(cx), Inches(2.46), Inches(2.60), Inches(1.10), C_CTR)
        add_textbox(slide, name,
                    Inches(cx + 0.05), Inches(2.52), Inches(2.50), Inches(0.32),
                    font="Calibri", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, ip,
                    Inches(cx + 0.05), Inches(2.86), Inches(2.50), Inches(0.28),
                    font="Courier New", size=9, color=C_MUTED2, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(7.00), Inches(3.66), Inches(5.60), Inches(0.40), C_OK)
    add_textbox(slide, "my-net  —  Docker DNS auto-configured",
                Inches(7.00), Inches(3.70), Inches(5.60), Inches(0.28),
                font="Calibri", size=11, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "✓ DNS by container name — no IPs needed",
                Inches(6.87), Inches(4.18), Inches(5.80), Inches(0.32),
                font="Calibri", size=11, italic=True, color=C_OK)
    add_textbox(slide, "docker exec app1 ping app2  (works by name!)",
                Inches(6.87), Inches(4.52), Inches(5.80), Inches(0.32),
                font="Courier New", size=9, color=C_NAVY)
    add_textbox(slide, "Isolated from containers on other networks",
                Inches(6.87), Inches(4.88), Inches(5.80), Inches(0.32),
                font="Calibri", size=11, color=C_MUTED)

    add_footer(slide, page)
    return page


def slide_docker_compose_arch_diagram(prs, page):
    """Docker Compose multi-service architecture: web + redis + postgres."""
    C_WEB   = RGBColor(0x2E, 0x74, 0xB5)
    C_REDIS = RGBColor(0xCC, 0x44, 0x44)
    C_DB    = RGBColor(0x7B, 0x54, 0xA0)
    C_NET   = RGBColor(0xF0, 0xF4, 0xFF)
    C_VOL   = RGBColor(0xEC, 0xEB, 0xE6)

    slide = content_slide(prs, "DOCKER COMPOSE", "Compose Multi-Service Architecture")

    # Compose network boundary
    add_rect(slide, Inches(0.55), Inches(1.90), Inches(9.00), Inches(3.80), C_NET)
    add_textbox(slide, "Compose Network  (app-net — created automatically)",
                Inches(0.60), Inches(1.94), Inches(8.80), Inches(0.26),
                font="Calibri", size=9, bold=True, color=C_MUTED)

    # web service
    add_rect(slide, Inches(0.75), Inches(2.36), Inches(2.60), Inches(2.60), C_WEB)
    add_textbox(slide, "web",
                Inches(0.80), Inches(2.42), Inches(2.50), Inches(0.38),
                font="Cambria", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "build: .\nports: 3000:3000\ndepends_on:\n  - db\n  - redis",
                Inches(0.85), Inches(2.86), Inches(2.38), Inches(1.00),
                font="Courier New", size=8, color=C_MUTED2)

    # Arrows web → redis and web → db
    add_textbox(slide, "→", Inches(3.42), Inches(2.80), Inches(0.55), Inches(0.42),
                font="Calibri", size=26, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "→", Inches(3.42), Inches(3.98), Inches(0.55), Inches(0.42),
                font="Calibri", size=26, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # redis service
    add_rect(slide, Inches(4.05), Inches(2.36), Inches(2.60), Inches(1.40), C_REDIS)
    add_textbox(slide, "redis",
                Inches(4.10), Inches(2.42), Inches(2.50), Inches(0.38),
                font="Cambria", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "image: redis:7-alpine\nport: 6379",
                Inches(4.15), Inches(2.86), Inches(2.38), Inches(0.70),
                font="Courier New", size=8, color=C_MUTED2)

    # db service
    add_rect(slide, Inches(4.05), Inches(3.90), Inches(2.60), Inches(1.68), C_DB)
    add_textbox(slide, "db  (postgres)",
                Inches(4.10), Inches(3.96), Inches(2.50), Inches(0.38),
                font="Cambria", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "image: postgres:16\nvolumes: pgdata:/var/lib\nhealthcheck: pg_isready",
                Inches(4.15), Inches(4.40), Inches(2.38), Inches(0.90),
                font="Courier New", size=8, color=C_MUTED2)

    # Volume (outside network)
    add_rect(slide, Inches(6.90), Inches(3.90), Inches(2.50), Inches(1.68), C_VOL)
    add_textbox(slide, "pgdata volume",
                Inches(6.95), Inches(3.96), Inches(2.40), Inches(0.38),
                font="Cambria", size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Managed by Docker\nPersists after down",
                Inches(6.95), Inches(4.40), Inches(2.40), Inches(0.70),
                font="Calibri", size=11, color=C_MUTED)
    add_textbox(slide, "↔", Inches(6.70), Inches(4.44), Inches(0.40), Inches(0.42),
                font="Calibri", size=22, bold=True, color=C_NAVY, align=PP_ALIGN.CENTER)

    # External user
    add_rect(slide, Inches(9.75), Inches(2.36), Inches(3.10), Inches(1.00),
             RGBColor(0x21, 0x20, 0x1C))
    add_textbox(slide, "External User",
                Inches(9.80), Inches(2.42), Inches(3.00), Inches(0.32),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "localhost:3000",
                Inches(9.80), Inches(2.78), Inches(3.00), Inches(0.28),
                font="Courier New", size=9, color=C_MUTED2, align=PP_ALIGN.CENTER)
    add_textbox(slide, "←",
                Inches(9.40), Inches(2.70), Inches(0.55), Inches(0.42),
                font="Calibri", size=26, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "All services started with one command:  docker-compose up -d",
                Inches(0.60), Inches(6.00), Inches(12.10), Inches(0.32),
                font="Calibri", size=12, italic=True, color=C_MUTED)

    add_footer(slide, page)
    return page


def slide_k8s_pod_anatomy_diagram(prs, page):
    """Pod anatomy — containers sharing network IP and volumes."""
    C_POD_BG  = RGBColor(0xE8, 0xF5, 0xE9)
    C_CTR     = RGBColor(0x10, 0x78, 0x54)
    C_SIDE    = RGBColor(0x2E, 0x74, 0xB5)
    C_NET_LBL = RGBColor(0x1F, 0x49, 0x7D)
    C_VOL_LBL = RGBColor(0x7B, 0x54, 0xA0)

    slide = content_slide(prs, "PODS — ANATOMY", "Inside a Pod: Containers, Network & Storage")

    # Pod outer boundary
    add_rect(slide, Inches(0.60), Inches(1.90), Inches(8.80), Inches(4.30), C_POD_BG)
    add_textbox(slide, "Pod   (my-app)",
                Inches(0.65), Inches(1.93), Inches(5.00), Inches(0.34),
                font="Cambria", size=14, bold=True, color=C_DARK)
    add_textbox(slide, "Node: worker-1    IP: 10.244.1.5",
                Inches(5.00), Inches(1.93), Inches(4.30), Inches(0.34),
                font="Calibri", size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)

    # Main container
    add_rect(slide, Inches(0.90), Inches(2.44), Inches(3.40), Inches(2.78), C_CTR)
    add_textbox(slide, "Main Container",
                Inches(0.95), Inches(2.50), Inches(3.30), Inches(0.36),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "image: my-app:v2\nport: 8080\n\nRuns the main\napplication",
                Inches(0.95), Inches(2.92), Inches(3.28), Inches(1.10),
                font="Calibri", size=11, color=C_MUTED2)

    # Sidecar container
    add_rect(slide, Inches(4.50), Inches(2.44), Inches(3.40), Inches(2.78), C_SIDE)
    add_textbox(slide, "Sidecar Container",
                Inches(4.55), Inches(2.50), Inches(3.30), Inches(0.36),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "image: log-collector\n\nCollects logs,\nforwards to\ncentral store",
                Inches(4.55), Inches(2.92), Inches(3.28), Inches(1.10),
                font="Calibri", size=11, color=C_MUTED2)

    # Shared network band
    add_rect(slide, Inches(0.90), Inches(5.30), Inches(7.00), Inches(0.54), C_NET_LBL)
    add_textbox(slide, "Shared Network  |  Same IP: 10.244.1.5  |  Same port space  |  Communicate via localhost",
                Inches(0.95), Inches(5.36), Inches(6.90), Inches(0.36),
                font="Calibri", size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Shared volume stripe
    add_rect(slide, Inches(8.00), Inches(2.44), Inches(1.28), Inches(3.40), C_VOL_LBL)
    add_textbox(slide, "Shared\nVolume\n/logs",
                Inches(8.05), Inches(3.30), Inches(1.18), Inches(0.80),
                font="Calibri", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Right info panel
    add_rect(slide, Inches(9.55), Inches(1.90), Inches(3.40), Inches(4.30), RGBColor(0xEC, 0xEB, 0xE6))
    add_textbox(slide, "Shared Resources",
                Inches(9.60), Inches(1.95), Inches(3.30), Inches(0.36),
                font="Cambria", size=14, bold=True, color=C_DARK)
    shared = [
        ("Network",   "Same IP & ports. Talk to each other via localhost."),
        ("Volumes",   "Read/write same files (e.g. log files)."),
        ("Lifecycle", "Start and stop together — all or nothing."),
    ]
    for i, (lbl, desc) in enumerate(shared):
        add_textbox(slide, lbl + ":",
                    Inches(9.65), Inches(2.44 + i * 1.04), Inches(3.20), Inches(0.30),
                    font="Calibri", size=12, bold=True, color=C_DARK)
        add_textbox(slide, desc,
                    Inches(9.65), Inches(2.76 + i * 1.04), Inches(3.20), Inches(0.44),
                    font="Calibri", size=11, color=C_MUTED, wrap=True)

    add_footer(slide, page)
    return page


def slide_k8s_replicaset_diagram(prs, page):
    """ReplicaSet self-healing control loop visualisation."""
    C_RUN  = RGBColor(0x10, 0x78, 0x54)
    C_DEAD = RGBColor(0xCC, 0x44, 0x44)
    C_NEW  = RGBColor(0x27, 0x6B, 0x52)
    C_CTRL = RGBColor(0x1F, 0x49, 0x7D)
    pw, ph = 1.70, 0.80

    slide = content_slide(prs, "DEPLOYMENTS — SELF-HEALING", "ReplicaSet Control Loop")

    # Before — 3 running pods
    add_textbox(slide, "BEFORE  (desired: 3 — actual: 3)",
                Inches(0.60), Inches(2.00), Inches(5.60), Inches(0.34),
                font="Cambria", size=13, bold=True, color=C_DARK)
    for i in range(3):
        add_rect(slide, Inches(0.60 + i * 1.86), Inches(2.44), Inches(pw), Inches(ph), C_RUN)
        add_textbox(slide, f"Pod {i+1}\nRunning",
                    Inches(0.65 + i * 1.86), Inches(2.58), Inches(pw - 0.10), Inches(0.48),
                    font="Calibri", size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Arrow
    add_textbox(slide, "→  Pod 2 crashes  →",
                Inches(5.78), Inches(2.68), Inches(2.40), Inches(0.36),
                font="Calibri", size=12, italic=True, color=C_DEAD, align=PP_ALIGN.CENTER)

    # Detected — pod 2 dead
    add_textbox(slide, "DETECTED  (desired: 3 — actual: 2  → gap!)",
                Inches(8.32), Inches(2.00), Inches(4.90), Inches(0.34),
                font="Cambria", size=13, bold=True, color=C_DARK)
    for i, status in enumerate(["Running", "DEAD", "Running"]):
        clr = C_DEAD if status == "DEAD" else C_RUN
        add_rect(slide, Inches(8.32 + i * 1.56), Inches(2.44), Inches(1.46), Inches(ph), clr)
        add_textbox(slide, f"Pod {i+1}\n{status}",
                    Inches(8.37 + i * 1.56), Inches(2.58), Inches(1.36), Inches(0.48),
                    font="Calibri", size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Controller box
    add_rect(slide, Inches(4.70), Inches(3.60), Inches(3.90), Inches(1.00), C_CTRL)
    add_textbox(slide, "ReplicaSet Controller",
                Inches(4.75), Inches(3.65), Inches(3.80), Inches(0.36),
                font="Cambria", size=13, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Watches pods  •  Detects gap  •  Creates replacement",
                Inches(4.75), Inches(4.02), Inches(3.80), Inches(0.44),
                font="Calibri", size=10, color=C_MUTED2, align=PP_ALIGN.CENTER)

    add_textbox(slide, "↑  detects crash",
                Inches(9.60), Inches(3.35), Inches(1.80), Inches(0.26),
                font="Calibri", size=10, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, "↓  creates replacement",
                Inches(4.80), Inches(4.66), Inches(3.70), Inches(0.28),
                font="Calibri", size=10, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # Healed state
    add_textbox(slide, "HEALED  (desired: 3 — actual: 3  — automatically!)",
                Inches(0.60), Inches(5.08), Inches(12.10), Inches(0.34),
                font="Cambria", size=13, bold=True, color=RGBColor(0x10, 0x78, 0x54))
    for i, (lbl, clr) in enumerate([("Pod 1   Running", C_RUN), ("Pod 3   Running", C_RUN), ("Pod 4   Running  (new)", C_NEW)]):
        add_rect(slide, Inches(0.60 + i * 4.20), Inches(5.52), Inches(3.90), Inches(ph), clr)
        add_textbox(slide, lbl,
                    Inches(0.65 + i * 4.20), Inches(5.80), Inches(3.80), Inches(0.28),
                    font="Calibri", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_footer(slide, page)
    return page


def slide_k8s_rolling_update_diagram(prs, page):
    """Rolling update step-by-step: old RS fades out, new RS ramps up."""
    C_OLD  = RGBColor(0x8A, 0x86, 0x7D)
    C_NEW  = RGBColor(0x10, 0x78, 0x54)
    C_DEPL = RGBColor(0x1F, 0x49, 0x7D)
    pw, ph = 0.84, 0.54

    slide = content_slide(prs, "ROLLOUTS — MECHANISM", "Rolling Update — Step by Step")

    steps = [
        ("Start\n(all v1)",          3, 0),
        ("Step 1\n+1 new, -1 old",   2, 1),
        ("Step 2\n+1 new, -1 old",   1, 2),
        ("Done\n(all v2)",           0, 3),
    ]

    for si, (label, old_n, new_n) in enumerate(steps):
        col_x = 0.55 + si * 3.18

        add_textbox(slide, label,
                    Inches(col_x), Inches(1.90), Inches(2.90), Inches(0.52),
                    font="Calibri", size=11, bold=True, color=C_DARK)

        add_textbox(slide, f"Old RS (v1) ×{old_n}",
                    Inches(col_x), Inches(2.52), Inches(2.80), Inches(0.26),
                    font="Calibri", size=9, color=C_OLD)
        for i in range(old_n):
            add_rect(slide, Inches(col_x + i * 0.92), Inches(2.80), Inches(pw), Inches(ph), C_OLD)
            add_textbox(slide, "v1",
                        Inches(col_x + i * 0.92), Inches(2.95), Inches(pw), Inches(0.26),
                        font="Calibri", size=9, color=C_WHITE, align=PP_ALIGN.CENTER)

        add_textbox(slide, f"New RS (v2) ×{new_n}",
                    Inches(col_x), Inches(3.54), Inches(2.80), Inches(0.26),
                    font="Calibri", size=9, color=C_NEW)
        for i in range(new_n):
            add_rect(slide, Inches(col_x + i * 0.92), Inches(3.82), Inches(pw), Inches(ph), C_NEW)
            add_textbox(slide, "v2",
                        Inches(col_x + i * 0.92), Inches(3.97), Inches(pw), Inches(0.26),
                        font="Calibri", size=9, color=C_WHITE, align=PP_ALIGN.CENTER)

        if si < 3:
            add_textbox(slide, "→",
                        Inches(col_x + 2.98), Inches(3.10), Inches(0.40), Inches(0.40),
                        font="Calibri", size=22, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # Deployment controller band
    add_rect(slide, Inches(0.55), Inches(4.72), Inches(12.10), Inches(0.62), C_DEPL)
    add_textbox(slide, "Deployment Controller  —  kubectl set image nginx=nginx:v2"
                       "  |  Old RS scaled to 0 but kept for rollback history",
                Inches(0.60), Inches(4.86), Inches(12.00), Inches(0.34),
                font="Calibri", size=10, color=C_ORANGE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Zero downtime: old and new pods run simultaneously during transition  "
                       "|  maxSurge/maxUnavailable control the pace",
                Inches(0.55), Inches(5.46), Inches(12.10), Inches(0.32),
                font="Calibri", size=11, italic=True, color=C_MUTED)

    add_footer(slide, page)
    return page


def slide_k8s_service_types_diagram(prs, page):
    """K8s Service types: ClusterIP, NodePort, LoadBalancer layered view."""
    C_LB   = RGBColor(0xD9, 0x77, 0x57)
    C_NP   = RGBColor(0x2E, 0x74, 0xB5)
    C_CI   = RGBColor(0x10, 0x78, 0x54)
    C_POD  = RGBColor(0x1E, 0x4D, 0x3A)

    slide = content_slide(prs, "SERVICES — TYPES", "Service Types: ClusterIP / NodePort / LoadBalancer")

    # Internet band
    add_rect(slide, Inches(0.55), Inches(2.00), Inches(12.10), Inches(0.60),
             RGBColor(0x3D, 0x3C, 0x38))
    add_textbox(slide, "Internet  /  External Users",
                Inches(0.60), Inches(2.10), Inches(12.00), Inches(0.34),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # LoadBalancer
    add_rect(slide, Inches(0.55), Inches(2.68), Inches(12.10), Inches(0.66), C_LB)
    add_textbox(slide, "LoadBalancer  (cloud LB — AWS ELB / GCP / Azure)  "
                       "|  External IP : 80  →  NodePort  →  ClusterIP  →  Pods",
                Inches(0.60), Inches(2.80), Inches(12.00), Inches(0.36),
                font="Calibri", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # NodePort
    add_rect(slide, Inches(0.55), Inches(3.42), Inches(12.10), Inches(0.66), C_NP)
    add_textbox(slide, "NodePort  |  Node IP : 30080  "
                       "|  Range 30000-32767  |  Reachable from outside cluster without cloud LB",
                Inches(0.60), Inches(3.54), Inches(12.00), Inches(0.36),
                font="Calibri", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # ClusterIP
    add_rect(slide, Inches(0.55), Inches(4.16), Inches(12.10), Inches(0.66), C_CI)
    add_textbox(slide, "ClusterIP  |  10.96.0.1 : 80  "
                       "|  Internal only  |  DNS: my-svc.default.svc.cluster.local  |  Load balances across Pods",
                Inches(0.60), Inches(4.28), Inches(12.00), Inches(0.36),
                font="Calibri", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Pods
    add_textbox(slide, "Pods  ↓",
                Inches(0.55), Inches(4.90), Inches(1.20), Inches(0.30),
                font="Calibri", size=10, color=C_MUTED)
    for i in range(3):
        add_rect(slide, Inches(1.90 + i * 3.10), Inches(4.90), Inches(2.60), Inches(0.68), C_POD)
        add_textbox(slide, f"Pod {i+1}  :8080",
                    Inches(1.95 + i * 3.10), Inches(5.10), Inches(2.50), Inches(0.34),
                    font="Calibri", size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Use-case legend
    legend = [
        ("ClusterIP",    "C_CI",  "Internal service-to-service calls (default type)"),
        ("NodePort",     "C_NP",  "External access without cloud — testing / on-prem"),
        ("LoadBalancer", "C_LB",  "Production on managed K8s (AKS, EKS, GKE)"),
    ]
    for i, (typ, _, use) in enumerate(legend):
        add_textbox(slide, f"{typ}:",
                    Inches(0.60), Inches(5.72 + i * 0.26), Inches(1.60), Inches(0.24),
                    font="Calibri", size=10, bold=True, color=C_DARK)
        add_textbox(slide, use,
                    Inches(2.25), Inches(5.72 + i * 0.26), Inches(10.00), Inches(0.24),
                    font="Calibri", size=10, color=C_MUTED)

    add_footer(slide, page)
    return page


def slide_k8s_pv_pvc_diagram(prs, page):
    """PersistentVolume → PVC → Pod binding chain."""
    C_PV  = RGBColor(0x7B, 0x54, 0xA0)
    C_PVC = RGBColor(0x2E, 0x74, 0xB5)
    C_POD = RGBColor(0x10, 0x78, 0x54)

    slide = content_slide(prs, "KUBERNETES STORAGE", "PV → PVC → Pod: Storage Binding Chain")

    items = [
        (C_PV,  "PersistentVolume\n(PV)",       "Cluster Admin creates\n1Gi hostPath storage\nReclaimPolicy: Retain",   "Cluster Resource"),
        (C_PVC, "PersistentVolume\nClaim (PVC)", "Developer requests\n500Mi ReadWriteOnce\nK8s finds matching PV",      "Pod’s Storage Request"),
        (C_POD, "Pod",                            "Mounts PVC at /data\nReads & writes files\nData survives pod restart", "App Container"),
    ]
    bw, bh, ys = 2.80, 1.90, 2.44

    for i, (clr, title, body, role) in enumerate(items):
        bx = 0.60 + i * 4.18
        add_rect(slide, Inches(bx), Inches(ys), Inches(bw), Inches(bh), clr)
        add_textbox(slide, title,
                    Inches(bx + 0.10), Inches(ys + 0.10), Inches(bw - 0.20), Inches(0.52),
                    font="Cambria", size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, body,
                    Inches(bx + 0.12), Inches(ys + 0.68), Inches(bw - 0.24), Inches(0.96),
                    font="Calibri", size=11, color=C_MUTED2, wrap=True)
        add_textbox(slide, f"Role: {role}",
                    Inches(bx + 0.10), Inches(ys + 1.66), Inches(bw - 0.20), Inches(0.26),
                    font="Calibri", size=9, italic=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        if i < 2:
            add_textbox(slide, "→\nbinds",
                        Inches(bx + bw + 0.22), Inches(ys + 0.60), Inches(0.76), Inches(0.64),
                        font="Calibri", size=16, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

    # Storage backend panel
    add_rect(slide, Inches(9.60), Inches(ys), Inches(3.20), Inches(bh), RGBColor(0xEC, 0xEB, 0xE6))
    add_textbox(slide, "Storage Backends",
                Inches(9.65), Inches(ys + 0.10), Inches(3.10), Inches(0.36),
                font="Cambria", size=13, bold=True, color=C_DARK)
    for i, b in enumerate(["hostPath  (single node / dev)", "NFS  (shared across nodes)",
                            "AWS EBS / Azure Disk", "GCP Persistent Disk", "Ceph / Rook  (on-prem)"]):
        add_textbox(slide, b,
                    Inches(9.68), Inches(ys + 0.54 + i * 0.28), Inches(3.04), Inches(0.26),
                    font="Calibri", size=10, color=C_DARK)

    add_textbox(slide,
                "Data survives:  pod restart  •  pod deletion  •  node reboot  "
                "(depending on storage backend and ReclaimPolicy)",
                Inches(0.60), Inches(4.56), Inches(12.10), Inches(0.34),
                font="Calibri", size=11, italic=True, color=C_MUTED)
    add_textbox(slide,
                "Static: Admin creates PV manually.   "
                "Dynamic: StorageClass auto-creates PV when PVC is submitted.",
                Inches(0.60), Inches(4.94), Inches(12.10), Inches(0.34),
                font="Calibri", size=10, color=C_DARK)

    add_footer(slide, page)
    return page


def slide_k8s_job_flow_diagram(prs, page):
    """CronJob → Job → Pods execution flow with comparison table."""
    C_CRON = RGBColor(0xD9, 0x77, 0x57)
    C_JOB  = RGBColor(0x2E, 0x74, 0xB5)
    C_POD  = RGBColor(0x10, 0x78, 0x54)
    C_DONE = RGBColor(0x8A, 0x86, 0x7D)

    slide = content_slide(prs, "JOBS & CRONJOBS", "CronJob → Job → Pods: Execution Flow")

    # CronJob box
    add_rect(slide, Inches(0.60), Inches(2.24), Inches(2.80), Inches(1.82), C_CRON)
    add_textbox(slide, "CronJob",
                Inches(0.65), Inches(2.30), Inches(2.70), Inches(0.40),
                font="Cambria", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "schedule:\n  \"*/5 * * * *\"\nRuns every 5 min",
                Inches(0.65), Inches(2.76), Inches(2.70), Inches(0.88),
                font="Courier New", size=9, color=C_MUTED2)

    add_textbox(slide, "→\nTriggers",
                Inches(3.48), Inches(2.68), Inches(1.30), Inches(0.72),
                font="Calibri", size=13, color=C_DARK, align=PP_ALIGN.CENTER)

    # Job box
    add_rect(slide, Inches(4.85), Inches(2.24), Inches(2.80), Inches(1.82), C_JOB)
    add_textbox(slide, "Job",
                Inches(4.90), Inches(2.30), Inches(2.70), Inches(0.40),
                font="Cambria", size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "completions: 3\nparallelism: 2\nrestartPolicy:\n  Never",
                Inches(4.90), Inches(2.76), Inches(2.70), Inches(0.88),
                font="Courier New", size=9, color=C_MUTED2)

    add_textbox(slide, "→\nCreates",
                Inches(7.73), Inches(2.68), Inches(1.00), Inches(0.72),
                font="Calibri", size=13, color=C_DARK, align=PP_ALIGN.CENTER)

    # Pods (2 parallel + 1 sequential)
    for i, (py, status, clr) in enumerate([(2.10, "Running", C_POD),
                                            (2.90, "Running", C_POD),
                                            (3.70, "Completed ✓", C_DONE)]):
        add_rect(slide, Inches(8.85), Inches(py), Inches(2.10), Inches(0.66), clr)
        add_textbox(slide, f"Pod {i+1}  {status}",
                    Inches(8.90), Inches(py + 0.18), Inches(2.00), Inches(0.30),
                    font="Calibri", size=10, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "parallelism: 2 → Pods 1 & 2 run simultaneously",
                Inches(8.85), Inches(4.45), Inches(4.00), Inches(0.28),
                font="Calibri", size=9, italic=True, color=C_MUTED)

    # Comparison table
    add_rect(slide, Inches(0.60), Inches(4.48), Inches(8.10), Inches(1.78),
             RGBColor(0xEC, 0xEB, 0xE6))
    add_textbox(slide, "Workload Comparison",
                Inches(0.65), Inches(4.54), Inches(7.90), Inches(0.36),
                font="Cambria", size=13, bold=True, color=C_DARK)
    rows = [
        ("Deployment", "Keep N pods always running (self-healing)", "Web server, API"),
        ("Job",        "Run pods until N completions succeed",       "DB migration, batch job"),
        ("CronJob",    "Create a Job on a cron schedule",           "Backup, cleanup, report"),
    ]
    for i, (name, desc, eg) in enumerate(rows):
        add_textbox(slide, name + ":",
                    Inches(0.70), Inches(4.98 + i * 0.38), Inches(1.40), Inches(0.34),
                    font="Calibri", size=11, bold=True, color=C_DARK)
        add_textbox(slide, desc,
                    Inches(2.10), Inches(4.98 + i * 0.38), Inches(3.80), Inches(0.34),
                    font="Calibri", size=11, color=C_DARK)
        add_textbox(slide, f"e.g. {eg}",
                    Inches(5.90), Inches(4.98 + i * 0.38), Inches(2.70), Inches(0.34),
                    font="Calibri", size=10, italic=True, color=C_MUTED)

    add_footer(slide, page)
    return page


# ══════════════════════════════════════════════════════════════════════════════
#  Slide definitions — Day 1
# ══════════════════════════════════════════════════════════════════════════════

def build_day1(prs):
    p = 0

    # ── 1. Cover ──────────────────────────────────────────────────────────────
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_DARK)
    add_circle(slide, 11.8, 0.6, 2.3, C_ORANGE, alpha=15)
    add_circle(slide, 12.5, 4.7, 1.7, C_ORANGE, alpha=10)
    add_textbox(slide, f"WSQ COURSE  .  {TGS}",
                Inches(0.90), Inches(1.50), Inches(10.0), Inches(0.35),
                font="Calibri", size=14, bold=True, color=C_ORANGE)
    add_textbox(slide, "Application Integration with",
                Inches(0.85), Inches(2.00), Inches(11.2), Inches(0.85),
                font="Cambria", size=50, bold=True, color=C_WHITE)
    add_textbox(slide, "Docker and Kubernetes",
                Inches(0.85), Inches(2.85), Inches(11.2), Inches(0.85),
                font="Cambria", size=50, bold=True, color=C_WHITE)
    add_textbox(slide, "2-Day Course  |  Day 1: Docker  |  Day 2: Kubernetes",
                Inches(0.90), Inches(3.90), Inches(11.0), Inches(0.40),
                font="Calibri", size=18, color=C_MUTED2)
    add_textbox(slide, "Tertiary Infotech Pte. Ltd.",
                Inches(0.90), Inches(5.70), Inches(11.0), Inches(0.35),
                font="Calibri", size=15, bold=True, color=C_WHITE)
    add_textbox(slide, "UEN: 201200696W",
                Inches(0.90), Inches(6.05), Inches(11.0), Inches(0.30),
                font="Calibri", size=13, color=C_MUTED)
    add_textbox(slide, "www.tertiarycourses.com.sg   |   enquiry@tertiaryinfotech.com   |   +65 6100 0613",
                Inches(0.90), Inches(6.58), Inches(11.5), Inches(0.30),
                font="Calibri", size=12, color=C_MUTED)
    add_textbox(slide, "Version 3.0",
                Inches(10.93), Inches(1.50), Inches(1.80), Inches(0.30),
                font="Calibri", size=12, color=C_MUTED)

    # ── 2. Housekeeping ───────────────────────────────────────────────────────
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_CREAM)
    add_section_label(slide, "HOUSEKEEPING")
    add_main_title(slide, "Before We Start", color=C_DARK)
    cols = [
        ("Digital Attendance",
         "Scan the QR code at the front to mark attendance.\n"
         "AM (9am) + PM (2pm) sign-in every session.\n\n"
         "Minimum 75% attendance required for SSG funding.\n"
         "Missing attendance = forfeited subsidy — you pay full course fee.\n\n"
         "URL:  attendance.ssg.gov.sg\n"
         "Course Code:  TGS-2021010366"),
        ("Ground Rules",
         "Phones on silent mode during sessions.\n"
         "Questions welcome at any time — no such thing as a bad question.\n"
         "Breaks: 15 min AM · 1 hr lunch · 15 min PM.\n"
         "Camera on during activities, labs, and assessments.\n"
         "KillerCoda terminal only — close unrelated browser tabs."),
        ("LMS Access",
         "Courseware & Assessment:\nhttps://lms.tertiaryinfotech.com\n\n"
         "Hands-on Labs (KillerCoda):\nhttps://killercoda.com/tertiary-labs\n\n"
         "GitHub Lab Files:\nhttps://github.com/tertiarycourses\n\n"
         "Login credentials on the whiteboard."),
    ]
    for i, (heading, body) in enumerate(cols):
        x = Inches(0.55 + i * 4.25)
        add_rect(slide, x, Inches(1.85), Inches(4.00), Inches(5.10), C_DARK)
        add_textbox(slide, heading, x + Inches(0.20), Inches(2.02), Inches(3.70), Inches(0.45),
                    font="Cambria", size=17, bold=True, color=C_ORANGE)
        add_textbox(slide, body, x + Inches(0.20), Inches(2.58), Inches(3.68), Inches(4.00),
                    font="Calibri", size=12, color=C_MUTED2, wrap=True)
    add_footer(slide, p)

    # ── 3. About the Trainer ──────────────────────────────────────────────────
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_CREAM)
    add_section_label(slide, "INTRODUCTION")
    add_main_title(slide, "About the Trainer", color=C_DARK)
    add_rect(slide, Inches(0.60), Inches(1.90), Inches(3.60), Inches(4.80), C_DARK)
    add_textbox(slide, "Mohan Pothula",
                Inches(0.70), Inches(2.85), Inches(3.50), Inches(0.55),
                font="Cambria", size=22, bold=True, color=C_WHITE)
    add_textbox(slide, "Enterprise AI & DevOps Trainer",
                Inches(0.70), Inches(3.42), Inches(3.50), Inches(0.35),
                font="Calibri", size=12, color=C_ORANGE)
    add_textbox(slide, "DBS Bank  •  SingTel  •  Mediacorp  •  SPH Media",
                Inches(0.70), Inches(3.82), Inches(3.50), Inches(0.35),
                font="Calibri", size=11, color=C_MUTED2)
    add_textbox(slide, "linkedin.com/in/mohanpothula",
                Inches(0.70), Inches(4.22), Inches(3.50), Inches(0.30),
                font="Calibri", size=11, italic=True, color=C_BLUE)
    add_textbox(slide, "mohanpothula@gmail.com",
                Inches(0.70), Inches(4.56), Inches(3.50), Inches(0.30),
                font="Calibri", size=11, color=C_MUTED2)
    items = [
        ("Experience",
         "20+ years of hands-on enterprise experience across DBS Bank, SingTel, Mediacorp, and SPH Media. "
         "Brings real production knowledge into the classroom — students learn from someone who has actually "
         "deployed 800+ ML models, managed 10,000-node Hadoop clusters, architected RAG systems on AWS Bedrock, "
         "and built Agentic AI pilots with OpenClaw and Hermes Agents."),
        ("Specialisation",
         "Bridging the gap between foundational concepts and real-world enterprise application. "
         "Training programmes tailored for engineers, architects, and business-technology professionals."),
        ("Expertise",      "Docker  •  Kubernetes  •  DevOps / CI-CD  •  Cloud (AWS)  •  ML Platforms  •  Agentic AI"),
        ("Certifications", "Certified Kubernetes Administrator (CKA)  •  Docker Certified Associate (DCA)  •  AWS Solutions Architect"),
    ]
    y = 2.00
    for label, text in items:
        add_textbox(slide, f"{label}:",
                    Inches(4.60), Inches(y), Inches(2.20), Inches(0.30),
                    font="Calibri", size=12, bold=True, color=C_DARK)
        lines = len(text) // 90 + 1
        add_textbox(slide, text,
                    Inches(4.60), Inches(y + 0.28), Inches(8.20), Inches(lines * 0.32),
                    font="Calibri", size=12, color=C_DARK, wrap=True)
        y += 0.30 + lines * 0.32 + 0.14
    add_footer(slide, p)

    # ── 4. Course Outline ─────────────────────────────────────────────────────
    p += 1
    slide = content_slide(prs, "ROADMAP", "Day 1 — Course Outline")
    left_topics = [
        "1.  What is Docker? — VMs vs Containers",
        "2.  Docker Architecture & Core Concepts",
        "3.  Running & Managing Containers",
        "4.  Building Images with Dockerfile",
        "5.  Docker Storage — Volumes & Bind Mounts",
    ]
    right_topics = [
        "6.  Docker Networking & Port Mapping",
        "7.  Configuration — Environment Variables",
        "8.  Docker Hub (Reference)",
        "9.  Docker Compose — Multi-Service Apps",
        "10. Knowledge Checks",
    ]
    for i, text in enumerate(left_topics):
        add_textbox(slide, text, Inches(0.70), Inches(1.95 + i * 0.72), Inches(6.0), Inches(0.55),
                    font="Calibri", size=15, color=C_DARK)
    for i, text in enumerate(right_topics):
        add_textbox(slide, text, Inches(6.90), Inches(1.95 + i * 0.72), Inches(6.0), Inches(0.55),
                    font="Calibri", size=15, color=C_DARK)
    # Bottom summary bar
    add_rect(slide, Inches(0.22), Inches(6.44), Inches(12.88), Inches(0.52), C_SUMMARY_BAR)
    add_textbox(slide,
                "Day 1: 12 Labs + 2 Knowledge Checks  |  Day 2: 7 Labs + 2 Knowledge Checks"
                "  |  Hands-on via KillerCoda browser environment",
                Inches(0.40), Inches(6.51), Inches(12.52), Inches(0.38),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_footer(slide, p)

    # ══ SECTION 1 — WHAT IS DOCKER? ═══════════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 1", "What is Docker?",
                    subtitle="Containers, images, and why they matter",
                    subtopics=["VMs vs Containers", "How Docker works", "Images & Layers",
                               "Docker Hub & Registry"],
                    page=p)

    # 6. Problem slide
    p += 1
    slide = content_slide(prs, "DOCKER — OVERVIEW", "The Problem: 'It Works on My Machine'")
    add_body_lines(slide, [
        ('T', "Docker is a software platform that allows you to build, test, and deploy "
              "applications quickly. It packages software into standardised units called "
              "containers that have everything the software needs to run, including "
              "libraries, system tools, code, and runtime.   (Source: aws.amazon.com/docker)"),
        ('BL', ''),
        ('B', "Common environment mismatch problems:"),
        ('T', "  - Works on developer laptop, breaks on the test server"),
        ('T', "  - Different Python / Node / Java versions installed"),
        ('T', "  - Missing environment variables or config files"),
        ('T', "  - Dependency conflicts between projects"),
        ('BL', ''),
        ('I', "Docker solves this by packaging the app AND its environment into one portable, "
              "standardised container that runs identically everywhere."),
    ])
    add_footer(slide, p)

    # 7. VMs vs Containers
    p += 1
    two_column_slide(prs,
        label="DOCKER — OVERVIEW", title="Virtual Machines vs Containers",
        left_heading="Virtual Machines",
        left_lines=[
            ('T', "Full OS kernel per VM"),
            ('T', "Boot time: seconds to minutes"),
            ('T', "Size: gigabytes"),
            ('T', "Hypervisor overhead"),
            ('T', "Strong hardware-level isolation"),
            ('BL', ''),
            ('S', "Examples: VMware, VirtualBox, AWS EC2"),
        ],
        right_heading="Containers",
        right_lines=[
            ('T', "Share the host OS kernel"),
            ('T', "Start time: milliseconds"),
            ('T', "Size: megabytes"),
            ('T', "Minimal overhead"),
            ('T', "Process-level isolation"),
            ('BL', ''),
            ('S', "Examples: Docker, Podman, containerd"),
        ],
        page=p)

    # VM vs Container visual diagram
    p += 1
    slide_vm_vs_container_diagram(prs, p)

    # 8. Docker Architecture
    p += 1
    slide = content_slide(prs, "DOCKER — OVERVIEW", "Docker Architecture")
    add_body_lines(slide, [
        ('B', "Docker Client"),
        ('T', "The CLI you use. Sends commands (docker run, docker build) to the daemon."),
        ('BL', ''),
        ('B', "Docker Daemon  (dockerd)"),
        ('T', "Runs on the host. Manages images, containers, networks, and volumes."),
        ('BL', ''),
        ('B', "Docker Registry"),
        ('T', "Central image store. Docker Hub is the default public registry."),
        ('BL', ''),
        ('I', "Flow:  docker pull -> Registry -> Daemon  |  docker run -> Daemon creates container"),
    ])
    add_footer(slide, p)

    # Docker Engine visual diagram
    p += 1
    slide_docker_engine_diagram(prs, p)

    # 9. Core Concepts
    p += 1
    slide = content_slide(prs, "DOCKER — OVERVIEW", "Core Concepts")
    add_body_lines(slide, [
        ('B', "Image"),
        ('T', "A standardised package that includes all of the files, binaries, libraries, and "
              "configurations to run a container.  Immutable — changes require a new image.   (docs.docker.com)"),
        ('BL', ''),
        ('B', "Container"),
        ('T', "An isolated process for each of your app's components — self-contained, isolated, "
              "independent, and portable.   (docs.docker.com)"),
        ('BL', ''),
        ('B', "Volume"),
        ('T', "Persistent data stores for containers, created and managed by Docker. "
              "Data survives container removal.   (docs.docker.com/engine/storage/volumes)"),
        ('BL', ''),
        ('B', "Network"),
        ('T', "Virtual network connecting containers. Custom bridge networks provide automatic "
              "DNS — containers reach each other by name, not IP."),
    ])
    add_footer(slide, p)

    # ══ SECTION 2 — DOCKER FUNDAMENTALS LABS 1 & 2 ════════════════════════════
    p += 1
    section_divider(prs, "SECTION 2", "Running Your First Containers",
                    subtitle="Labs 1 & 2 — Docker Fundamentals",
                    subtopics=["Lab 1: Run your first container", "Lab 2: Explore a container",
                               "Lab 3: Build a custom image", "Lab 4: Flask web app"],
                    page=p)

    # 11. What is a Container?
    p += 1
    slide = content_slide(prs, "DOCKER — FUNDAMENTALS", "What is a Docker Container?")
    add_body_lines(slide, [
        ('T', "Containers are isolated processes for each of your app's components. "
              "Each component runs in its own isolated environment, completely separate "
              "from everything else on the machine.   (docs.docker.com)"),
        ('BL', ''),
        ('B', "Four key properties (docs.docker.com):"),
        ('T', "  Self-contained  — has everything required to function without host dependencies"),
        ('T', "  Isolated        — runs separately; minimal impact on host and other containers"),
        ('T', "  Independent     — managed individually; removing one doesn't affect others"),
        ('T', "  Portable        — runs consistently across dev machines, data centres, and cloud"),
        ('BL', ''),
        ('B', "Container lifecycle:"),
        ('C', "docker run   # create + start\ndocker stop  # graceful shutdown\ndocker start # restart a stopped container\ndocker rm    # delete the container (not the image)"),
        ('BL', ''),
        ('I', "Image = read-only blueprint.  Container = running instance with a thin writable layer on top."),
    ])
    add_footer(slide, p)

    # 12. Docker Images & the Registry
    p += 1
    slide = content_slide(prs, "DOCKER — FUNDAMENTALS", "Docker Images & the Registry")
    add_body_lines(slide, [
        ('T', "A container image is a standardised package that includes all of the files, "
              "binaries, libraries, and configurations to run a container.   (docs.docker.com)"),
        ('BL', ''),
        ('B', "Images are composed of layers — each layer represents a set of filesystem changes:"),
        ('T', "  Base OS layer  ->  runtime layer  ->  dependencies  ->  app code"),
        ('T', "  Unchanged layers are reused from cache — fast rebuilds, reduced bandwidth."),
        ('BL', ''),
        ('B', "Image naming:  name:tag"),
        ('C', "nginx:latest        # official Nginx, latest tag\npython:3.11-slim   # slim Python 3.11 image\nmyapp:v1.0         # your own custom image"),
        ('BL', ''),
        ('B', "Registry — centralised location for storing and sharing container images:"),
        ('T', "  Docker Hub (hub.docker.com) = default public registry, 100,000+ images"),
        ('T', "  docker pull  downloads an image | docker push  uploads your image"),
    ])
    add_footer(slide, p)

    # 13. Lab 1 header
    p += 1
    lab_header_slide(prs, "LAB 1", "Run Your First Container",
                     "day1-01-docker-fundamentals", p)

    # 12. Lab 1 — Commands
    p += 1
    slide = content_slide(prs, "LAB 1 — DOCKER FUNDAMENTALS", "Run Your First Container")
    add_lab_body_lines(slide, [
        ('B', "Run an interactive Ubuntu container:"),
        ('C', "docker run -it ubuntu:latest"),
        ('T', "Inside the container:"),
        ('C', "cat /etc/hosts\nexit"),
        ('B', "Check container status on the host:"),
        ('C', "docker ps             # only running containers\ndocker ps -a          # all containers including stopped\ndocker ps -a | grep ubuntu"),
        ('BL', ''),
        ('I', "Key concept: -it = interactive terminal. exit stops but does NOT delete the container."),
    ])
    kc_note(slide, "day1-01-docker-fundamentals")
    add_footer(slide, p)

    # 13. Lab 2 header
    p += 1
    lab_header_slide(prs, "LAB 2", "Nginx + docker cp",
                     "day1-01-docker-fundamentals", p)

    # 14. Lab 2 — Commands
    p += 1
    slide = content_slide(prs, "LAB 2 — DOCKER FUNDAMENTALS", "Run Nginx and Copy a File")
    add_lab_body_lines(slide, [
        ('B', "Run Nginx in the background (detached):"),
        ('C', "docker run -d --name my-nginx nginx:latest\ndocker ps"),
        ('B', "Copy a file from container to host:"),
        ('C', "docker cp my-nginx:/usr/share/nginx/html/index.html ./index.html\ncat index.html"),
        ('B', "Inspect logs then clean up:"),
        ('C', "docker logs my-nginx\ndocker stop my-nginx\ndocker rm my-nginx"),
        ('BL', ''),
        ('I', "Key concept: docker cp works host -> container and container -> host. Container need not be running."),
    ])
    kc_note(slide, "day1-01-docker-fundamentals")
    add_footer(slide, p)

    # ══ SECTION 3 — DOCKERFILE ════════════════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 3", "Building Images with Dockerfile",
                    subtitle="Labs 3 & 4 — Custom images and Flask app",
                    subtopics=["Dockerfile instructions", "Image layers & cache",
                               "CMD vs ENTRYPOINT", "Lab 3: Custom image",
                               "Lab 4: Flask web app"],
                    page=p)

    # 16. What is a Dockerfile?
    p += 1
    slide = content_slide(prs, "DOCKERFILE — OVERVIEW", "What is a Dockerfile?")
    add_body_lines(slide, [
        ('T', "A Dockerfile is a text document containing instructions that Docker uses to "
              "automatically construct images — all the commands a user could call on the "
              "command line to assemble an image.   (docs.docker.com/reference/dockerfile)"),
        ('BL', ''),
        ('B', "Key instructions:"),
        ('T', "  FROM      initiates a build stage and sets the base image (must come first)"),
        ('T', "  WORKDIR   sets working directory for RUN, CMD, COPY, ADD instructions that follow"),
        ('T', "  COPY      transfers files from the build context into the image filesystem"),
        ('T', "  RUN       executes commands during build, committing results as a new layer"),
        ('T', "  EXPOSE    documents which port the containerised app listens on (does NOT publish)"),
        ('T', "  ENV       sets environment variables that persist into the final container"),
        ('T', "  CMD       specifies the default command executed when a container starts"),
        ('BL', ''),
        ('I', "Each instruction creates a cached layer — only changed layers rebuild."),
    ])
    add_footer(slide, p)

    # 17. Image Layers
    p += 1
    slide = content_slide(prs, "DOCKERFILE — LAYERS", "Image Layers & Build Cache")
    add_body_lines(slide, [
        ('T', "Each layer in an image contains a set of filesystem changes — additions, deletions, "
              "or modifications. Layers are reused between images, reducing build time, storage, "
              "and bandwidth.   (docs.docker.com/get-started/docker-concepts/building-images)"),
        ('BL', ''),
        ('B', "Cache rule: if a layer is unchanged -> reused from cache."),
        ('B', "If any layer changes -> all subsequent layers are rebuilt from scratch."),
        ('BL', ''),
        ('B', "Best practice — order instructions for fast rebuilds:"),
        ('C', "FROM python:3.11-slim\nWORKDIR /app\nCOPY requirements.txt .   # changes rarely\nRUN pip install -r requirements.txt\nCOPY . .                  # source code changes often"),
        ('BL', ''),
        ('I', "Put stable instructions first. Frequently-changing files (source code) go last."),
    ])
    add_footer(slide, p)

    # Image Layers visual diagram
    p += 1
    slide_docker_image_layers_diagram(prs, p)

    # 18. CMD vs ENTRYPOINT
    p += 1
    two_column_slide(prs,
        label="DOCKERFILE — INSTRUCTIONS", title="CMD vs ENTRYPOINT",
        left_heading="CMD",
        left_lines=[
            ('T', "Default command — can be overridden at runtime."),
            ('BL', ''),
            ('C', 'CMD ["python", "app.py"]'),
            ('BL', ''),
            ('T', "Override at runtime:"),
            ('C', "docker run myimage bash"),
            ('BL', ''),
            ('S', "Use when: you want a replaceable default."),
        ],
        right_heading="ENTRYPOINT",
        right_lines=[
            ('T', "Fixed executable — CMD or runtime args become its arguments."),
            ('BL', ''),
            ('C', 'ENTRYPOINT ["python"]\nCMD ["app.py"]'),
            ('BL', ''),
            ('T', "Result: python app.py"),
            ('BL', ''),
            ('S', "Use when: the container IS a command (e.g. a CLI tool)."),
        ],
        page=p)

    # 19. Lab 3 header
    p += 1
    lab_header_slide(prs, "LAB 3", "Build a Python Image",
                     "day1-01-docker-fundamentals", p)

    # 20. Lab 3 — Dockerfile content
    p += 1
    slide = content_slide(prs, "LAB 3 — DOCKERFILE", "Build a Python Image — Dockerfile")
    add_lab_body_lines(slide, [
        ('B', "File: main.py"),
        ('C', 'print("Hello from inside Docker!")'),
        ('BL', ''),
        ('B', "File: Dockerfile"),
        ('C', "FROM python:3.11-slim\nWORKDIR /app\nCOPY main.py .\nCMD [\"python\", \"main.py\"]"),
        ('BL', ''),
        ('I', "Each line is a cached image layer. Run docker build to compile into an image."),
    ])
    kc_note(slide, "day1-01-docker-fundamentals")
    add_footer(slide, p)

    # 21. Lab 3 — Commands
    p += 1
    slide = content_slide(prs, "LAB 3 — DOCKERFILE", "Build a Python Image — Commands")
    add_lab_body_lines(slide, [
        ('B', "Create the project folder and files (on KillerCoda):"),
        ('C', "mkdir lab3 && cd lab3\n# Create main.py and Dockerfile (see previous slide)"),
        ('B', "Build the image:"),
        ('C', "docker build -t lab3 ."),
        ('B', "Run the container:"),
        ('C', "docker run lab3"),
        ('B', "Verify:"),
        ('C', "docker images | grep lab3\ndocker ps -a | grep lab3\ncd .."),
        ('BL', ''),
        ('I', "Expected output:  Hello from inside Docker!"),
    ])
    kc_note(slide, "day1-01-docker-fundamentals")
    add_footer(slide, p)

    # 22. Lab 4 header
    p += 1
    lab_header_slide(prs, "LAB 4", "Flask App in Docker",
                     "day1-01-docker-fundamentals", p)

    # 23. Lab 4 — Dockerfile
    p += 1
    slide = content_slide(prs, "LAB 4 — DOCKERFILE", "Flask App in Docker — Files")
    add_lab_body_lines(slide, [
        ('B', "File: app.py"),
        ('C', "from flask import Flask\napp = Flask(__name__)\n@app.route(\"/\")\ndef home():\n    return \"<h1>Hello from Flask in Docker!</h1>\"\nif __name__ == \"__main__\":\n    app.run(host=\"0.0.0.0\", port=5000)"),
        ('BL', ''),
        ('B', "File: requirements.txt"),
        ('C', "flask"),
        ('BL', ''),
        ('B', "File: Dockerfile"),
        ('C', "FROM python:3.11-slim\nWORKDIR /app\nCOPY requirements.txt .\nRUN pip install -r requirements.txt\nCOPY . .\nEXPOSE 5000\nCMD [\"python\", \"app.py\"]"),
    ])
    kc_note(slide, "day1-01-docker-fundamentals")
    add_footer(slide, p)

    # 24. Lab 4 — Commands
    p += 1
    slide = content_slide(prs, "LAB 4 — DOCKERFILE", "Flask App in Docker — Commands")
    add_lab_body_lines(slide, [
        ('C', "mkdir lab4 && cd lab4\n# Create app.py, requirements.txt, Dockerfile (see previous slide)"),
        ('B', "Build and run:"),
        ('C', "docker build -t lab4 .\ndocker run -d -p 5001:5000 --name flask-app lab4"),
        ('B', "Test the app:"),
        ('C', "curl http://localhost:5001"),
        ('B', "Clean up:"),
        ('C', "docker stop flask-app && docker rm flask-app && cd .."),
        ('BL', ''),
        ('I', "Key concept: EXPOSE 5000 documents the port. -p 5001:5000 maps host -> container."),
    ])
    kc_note(slide, "day1-01-docker-fundamentals")
    add_footer(slide, p)

    # ══ SECTION 4 — DOCKER STORAGE ════════════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 4", "Docker Storage",
                    subtitle="Lab 5 — Named volumes and bind mounts",
                    subtopics=["Why volumes?", "Named volumes", "Bind mounts",
                               "Lab 5a: Named volume", "Lab 5b: Bind mount"],
                    page=p)

    # 26. Why Volumes?
    p += 1
    slide = content_slide(prs, "DOCKER STORAGE", "Why Do We Need Volumes?")
    add_body_lines(slide, [
        ('T', "By default, all data written inside a container lives in its writable layer. "
              "When the container is removed, that data is gone."),
        ('BL', ''),
        ('B', "Docker provides two persistence mechanisms   (docs.docker.com/engine/storage):"),
        ('T', "  Named Volumes  — persistent data stores completely managed by Docker; "
              "isolated from the host filesystem; easier to back up and migrate"),
        ('T', "  Bind Mounts    — you specify the exact host path; instant two-way file sync; "
              "dependent on the host's directory structure and OS"),
        ('BL', ''),
        ('B', "When to use each:"),
        ('T', "  Volumes:      databases, persistent app data, production workloads"),
        ('T', "  Bind mounts:  live code sync and hot-reload during development"),
    ])
    add_footer(slide, p)

    # 27. Lab 5a header
    p += 1
    lab_header_slide(prs, "LAB 5a", "Named Volumes",
                     "day1-02-docker-storage", p)

    # 28. Lab 5a — Commands
    p += 1
    slide = content_slide(prs, "LAB 5a — DOCKER STORAGE", "Named Volumes — Commands")
    add_lab_body_lines(slide, [
        ('B', "Create and inspect a volume:"),
        ('C', "docker volume create my-data\ndocker volume ls\ndocker volume inspect my-data"),
        ('B', "Write data with container app1:"),
        ('C', "docker run -d --name app1 -v my-data:/data alpine sleep 3600\ndocker exec app1 sh -c \"echo 'hello from app1' > /data/test.txt\"\ndocker exec app1 cat /data/test.txt"),
        ('B', "Remove app1 then read data with app2 (volume persists!):"),
        ('C', "docker stop app1 && docker rm app1\ndocker run -d --name app2 -v my-data:/data alpine sleep 3600\ndocker exec app2 cat /data/test.txt"),
        ('B', "Clean up:"),
        ('C', "docker stop app2 && docker rm app2 && docker volume rm my-data"),
        ('I', "Key concept: Named volumes outlive containers. Use for databases and persistent data."),
    ])
    kc_note(slide, "day1-02-docker-storage")
    add_footer(slide, p)

    # Volume comparison — concept bridge between Lab 5a and Lab 5b
    p += 1
    two_column_slide(prs,
        label="DOCKER STORAGE — SUMMARY", title="Named Volumes vs Bind Mounts",
        left_heading="Named Volume",
        left_lines=[
            ('T', "Managed by Docker"),
            ('T', "Stored in /var/lib/docker/volumes"),
            ('T', "Portable across hosts (with plugins)"),
            ('T', "Best for: production data, databases"),
            ('BL', ''),
            ('C', "-v my-vol:/app/data"),
        ],
        right_heading="Bind Mount",
        right_lines=[
            ('T', "You specify the exact host path"),
            ('T', "Host filesystem is the source"),
            ('T', "Instant two-way file sync"),
            ('T', "Best for: dev, live code reload"),
            ('BL', ''),
            ('C', "-v $(pwd)/src:/app/src"),
        ],
        page=p)

    # Docker Storage visual diagram — before Lab 5b
    p += 1
    slide_docker_storage_diagram(prs, p)

    # Lab 5b header
    p += 1
    lab_header_slide(prs, "LAB 5b", "Bind Mounts",
                     "day1-02-docker-storage", p)

    # Lab 5b — Commands
    p += 1
    slide = content_slide(prs, "LAB 5b — DOCKER STORAGE", "Bind Mounts — Commands")
    add_lab_body_lines(slide, [
        ('B', "Create a host file and read it inside a container:"),
        ('C', "mkdir -p myfiles\necho \"hello from host\" > myfiles/note.txt\ndocker run --rm -v $(pwd)/myfiles:/data alpine cat /data/note.txt"),
        ('B', "Write from inside the container back to the host:"),
        ('C', "docker run --rm -v $(pwd)/myfiles:/data alpine \\\n  sh -c \"echo 'hello from container' > /data/reply.txt\"\ncat myfiles/reply.txt\nrm -rf myfiles"),
        ('BL', ''),
        ('I', "Changes in either direction are instant — no copy step needed."),
    ])
    kc_note(slide, "day1-02-docker-storage")
    add_footer(slide, p)

    # ══ SECTION 5 — DOCKER NETWORKING ═════════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 5", "Docker Networking",
                    subtitle="Labs 6 & 7 — Custom networks and port mapping",
                    subtopics=["Network drivers", "Custom bridge networks",
                               "Container DNS", "Port mapping",
                               "Lab 6: Custom network", "Lab 7: Port mapping"],
                    page=p)

    # 33. Networking Concepts
    p += 1
    slide = content_slide(prs, "DOCKER NETWORKING", "Networking Concepts")
    add_body_lines(slide, [
        ('B', "Default networks:"),
        ('T', "  bridge  — default for standalone containers (no inter-container DNS)"),
        ('T', "  host    — container shares the host network stack (Linux only)"),
        ('T', "  none    — no networking"),
        ('BL', ''),
        ('B', "Custom bridge network (recommended):"),
        ('T', "  Containers on the same custom network reach each other by container name"),
        ('T', "  Docker provides automatic DNS — no hardcoded IP addresses needed"),
        ('T', "  Isolated from containers on other networks"),
        ('BL', ''),
        ('C', "docker network ls\ndocker network create my-net\ndocker network inspect my-net"),
    ])
    add_footer(slide, p)

    # Docker Networking visual diagram
    p += 1
    slide_docker_networking_diagram(prs, p)

    # 34. Lab 6 header
    p += 1
    lab_header_slide(prs, "LAB 6", "Custom Networks",
                     "day1-03-docker-networking", p)

    # 35. Lab 6 — Commands
    p += 1
    slide = content_slide(prs, "LAB 6 — NETWORKING", "Custom Networks — Commands")
    add_lab_body_lines(slide, [
        ('B', "Create a network and two containers:"),
        ('C', "docker network create my-net\ndocker run -d --name app1 --network my-net busybox sleep 3600\ndocker run -d --name app2 --network my-net busybox sleep 3600"),
        ('B', "Ping by container name (no IP needed):"),
        ('C', "docker exec app1 ping -c 3 app2"),
        ('B', "Disconnect / reconnect:"),
        ('C', "docker network disconnect my-net app2\ndocker exec app1 ping -c 2 app2     # fails -- expected\ndocker network connect my-net app2\ndocker exec app1 ping -c 2 app2     # works again"),
        ('B', "Clean up:"),
        ('C', "docker stop app1 app2 && docker rm app1 app2 && docker network rm my-net"),
        ('I', "Key concept: Custom bridge networks provide automatic DNS by container name."),
    ])
    kc_note(slide, "day1-03-docker-networking")
    add_footer(slide, p)

    # Port mapping concept slide
    p += 1
    slide = content_slide(prs, "DOCKER NETWORKING", "Port Mapping — Connecting Containers to the Host")
    add_body_lines(slide, [
        ('T', "Containers run inside an isolated network namespace. By default, no traffic from "
              "the outside world can reach them. Port mapping punches a hole."),
        ('BL', ''),
        ('B', "Syntax:  -p <host_port>:<container_port>"),
        ('C', "docker run -p 8080:80 nginx     # host:8080 -> container:80\ndocker run -p 5001:5000 myapp  # host:5001 -> container:5000"),
        ('BL', ''),
        ('B', "EXPOSE vs -p:"),
        ('T', "  EXPOSE 5000 in Dockerfile  — documents the port, does NOT publish it"),
        ('T', "  -p 5001:5000 at runtime    — actually makes the port reachable"),
        ('BL', ''),
        ('B', "Shortcut flags:"),
        ('C', "docker run -P myapp       # -P maps ALL EXPOSE'd ports to random host ports\ndocker port myapp        # show current port mappings"),
        ('BL', ''),
        ('I', "You can run multiple containers from the same image on different host ports simultaneously."),
    ])
    add_footer(slide, p)

    # Lab 7 header
    p += 1
    lab_header_slide(prs, "LAB 7", "Port Mapping",
                     "day1-03-docker-networking", p)

    # 37. Lab 7 — Commands
    p += 1
    slide = content_slide(prs, "LAB 7 — NETWORKING", "Port Mapping — Commands")
    add_lab_body_lines(slide, [
        ('T', "Uses the Flask image built in Lab 4. Syntax: -p <host>:<container>"),
        ('B', "Run two instances on different host ports:"),
        ('C', "docker run -d -p 5001:5000 --name web1 lab4\ncurl http://localhost:5001\ndocker run -d -p 5002:5000 --name web2 lab4\ncurl http://localhost:5002"),
        ('B', "View port mappings:"),
        ('C', "docker port web1\ndocker port web2"),
        ('B', "Let Docker assign a random host port (-P):"),
        ('C', "docker run -d -P --name web3 lab4\ndocker port web3"),
        ('B', "Clean up:"),
        ('C', "docker stop web1 web2 web3 && docker rm web1 web2 web3"),
        ('I', "-P maps all EXPOSE'd ports to random host ports."),
    ])
    kc_note(slide, "day1-03-docker-networking")
    add_footer(slide, p)

    # ══ SECTION 6 — CONFIGURATION ═════════════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 6", "Configuration with Environment Variables",
                    subtitle="Lab 8 — Runtime config without rebuilding",
                    subtopics=["ENV in Dockerfile", "Runtime -e flag",
                               ".env files", "Lab 8: Environment config"],
                    page=p)

    # 39. Env vars theory
    p += 1
    slide = content_slide(prs, "CONFIGURATION", "Why Environment Variables?")
    add_body_lines(slide, [
        ('T', "Environment variables let you configure a container at runtime — no rebuild needed."),
        ('BL', ''),
        ('B', "Three ways to pass env vars:"),
        ('T', "  1.  ENV in Dockerfile — bakes in defaults"),
        ('C', "ENV MY_NAME=World"),
        ('T', "  2.  -e flag at runtime — override per container"),
        ('C', "docker run -e MY_NAME=Alfred myimage"),
        ('T', "  3.  --env-file — pass many vars from a file"),
        ('C', "docker run --env-file .env myimage"),
        ('BL', ''),
        ('I', "Never hardcode secrets (passwords, API keys) in a Dockerfile."),
    ])
    add_footer(slide, p)

    # 40. Lab 8 header
    p += 1
    lab_header_slide(prs, "LAB 8", "Environment Variables",
                     "day1-04-docker-config", p)

    # 41. Lab 8 — Dockerfile
    p += 1
    slide = content_slide(prs, "LAB 8 — CONFIGURATION", "Environment Variables — Files")
    add_lab_body_lines(slide, [
        ('B', "File: main.py"),
        ('C', "import os\nname = os.environ.get(\"MY_NAME\", \"World\")\nenv  = os.environ.get(\"MY_ENV\",  \"development\")\nprint(f\"Hello, {name}! Environment: {env}\")"),
        ('BL', ''),
        ('B', "File: Dockerfile"),
        ('C', "FROM python:3.11-slim\nWORKDIR /app\nENV MY_NAME=World\nENV MY_ENV=development\nCOPY main.py .\nCMD [\"python\", \"main.py\"]"),
    ])
    kc_note(slide, "day1-04-docker-config")
    add_footer(slide, p)

    # 42. Lab 8 — Commands
    p += 1
    slide = content_slide(prs, "LAB 8 — CONFIGURATION", "Environment Variables — Commands")
    add_lab_body_lines(slide, [
        ('C', "mkdir lab8 && cd lab8\n# Create main.py and Dockerfile (see previous slide)"),
        ('B', "Build the image:"),
        ('C', "docker build -t lab8 ."),
        ('B', "Run with different env var combinations:"),
        ('C', "docker run lab8                                      # uses Dockerfile defaults\ndocker run -e MY_NAME=Alfred lab8                    # override name\ndocker run -e MY_NAME=Alfred -e MY_ENV=production lab8"),
        ('B', "Use an env file:"),
        ('C', "echo \"MY_NAME=Student\\nMY_ENV=staging\" > .env\ndocker run --env-file .env lab8"),
        ('B', "View all environment variables inside the container:"),
        ('C', "docker run lab8 env"),
        ('B', "Clean up:"),
        ('C', "cd .. && docker system prune -f"),
    ])
    kc_note(slide, "day1-04-docker-config")
    add_footer(slide, p)

    # 43. Docker Hub (theory)
    p += 1
    slide = content_slide(prs, "DOCKER HUB", "Sharing Images — Docker Hub")
    add_body_lines(slide, [
        ('T', "A registry is a centralised location for storing and sharing container images. "
              "Docker Hub is the default public registry — over 100,000 images including "
              "Docker Official Images and Verified Publishers.   (docs.docker.com)"),
        ('BL', ''),
        ('B', "Registry vs Repository:"),
        ('T', "  Registry  = the overall system (e.g. Docker Hub, Amazon ECR, Azure ACR)"),
        ('T', "  Repository = a collection of related images within a registry (like a folder)"),
        ('BL', ''),
        ('B', "Push workflow (Lab 9 — Reference Only):"),
        ('C', "docker login\ndocker tag lab4 <username>/lab4:v1\ndocker push <username>/lab4:v1"),
        ('B', "Pull from anywhere:"),
        ('C', "docker pull <username>/lab4:v1\ndocker run -p 5001:5000 <username>/lab4:v1"),
        ('BL', ''),
        ('I', "Lab 9 is Reference Only — requires a Docker Hub account and is not on KillerCoda."),
    ])
    add_footer(slide, p)

    # 44. Knowledge Check 1
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_AMBER)
    add_section_label(slide, "KNOWLEDGE CHECK  1", color=C_DARK)
    add_main_title(slide, "Quick Review — Docker Basics", color=C_DARK)
    add_body_lines(slide, [
        ('B', "Q1.  What is the difference between a Docker image and a container?"),
        ('BL', ''),
        ('B', "Q2.  Which command shows ALL containers including stopped ones?"),
        ('C', "A)  docker ps        B)  docker ps -a        C)  docker images        D)  docker ls -a"),
        ('BL', ''),
        ('B', "Q3.  Which Dockerfile instruction sets the default command when a container starts?"),
        ('BL', ''),
        ('B', "Q4.  True / False:  A named volume is deleted when you run docker rm."),
        ('BL', ''),
        ('B', "Q5.  What flag maps host port 8080 to container port 3000?"),
        ('C', "A)  -p 3000:8080        B)  -p 8080:3000        C)  -P 8080:3000        D)  --port 8080/3000"),
    ], y=Inches(1.80))
    add_footer(slide, p)

    # ══ SECTION 7 — DOCKER COMPOSE ════════════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 7", "Docker Compose",
                    subtitle="Labs 10, 11 & 12 — Multi-service apps",
                    subtopics=["Compose file structure", "Services & networks",
                               "Lab 10: First Compose app", "Lab 11: Multi-service",
                               "Lab 12: Full-stack app"],
                    page=p)

    # 46. What is Compose?
    p += 1
    slide = content_slide(prs, "DOCKER COMPOSE", "What is Docker Compose?")
    add_body_lines(slide, [
        ('T', "Docker Compose is a tool for defining and running multi-container applications. "
              "With Compose, you can define all of your containers and their configurations "
              "in a single YAML file.   (docs.docker.com/compose)"),
        ('BL', ''),
        ('B', "Without Compose — run each step manually:"),
        ('C', "docker network create app-net\ndocker run -d --name redis --network app-net redis:7\ndocker run -d -p 5001:5000 --network app-net web"),
        ('BL', ''),
        ('B', "With Compose — single-command deployment (anyone who clones the repo can run):"),
        ('C', "docker-compose up -d"),
        ('BL', ''),
        ('I', "Compose operates declaratively — specify desired state; Compose applies changes "
              "intelligently. Handles networking, volumes, dependency order, and env vars."),
    ])
    add_footer(slide, p)

    # Docker Compose architecture diagram
    p += 1
    slide_docker_compose_arch_diagram(prs, p)

    # 47. Compose YAML structure
    p += 1
    slide = content_slide(prs, "DOCKER COMPOSE", "docker-compose.yml Structure")
    add_body_lines(slide, [
        ('C', "services:\n  web:\n    build: .              # build from Dockerfile\n    ports:\n      - \"5001:5000\"     # host:container\n    depends_on:\n      - redis\n    environment:\n      - MY_ENV=production\n\n  redis:\n    image: redis:7-alpine # pull from Docker Hub\n\nvolumes:\n  mydata:               # named volume declaration"),
    ])
    add_footer(slide, p)

    # 48. Compose Commands
    p += 1
    slide = content_slide(prs, "DOCKER COMPOSE", "Key Compose Commands")
    add_body_lines(slide, [
        ('C', "docker-compose up -d       # start all services in background\ndocker-compose ps          # list service status\ndocker-compose logs        # view all service logs\ndocker-compose logs web    # logs for one service\ndocker-compose exec web sh # open shell in service\ndocker-compose stop        # stop without removing\ndocker-compose down        # stop and remove containers + networks\ndocker-compose down -v     # also remove volumes"),
        ('BL', ''),
        ('I', "docker-compose down removes the containers and networks created by up. Volumes remain unless -v."),
    ])
    add_footer(slide, p)

    # 49. Lab 10 header
    p += 1
    lab_header_slide(prs, "LAB 10", "Single Service Compose",
                     "day1-05-docker-compose", p)

    # 50. Lab 10 — Files
    p += 1
    slide = content_slide(prs, "LAB 10 — DOCKER COMPOSE", "Single Service Compose — Files")
    add_lab_body_lines(slide, [
        ('B', "File: app.py"),
        ('C', "from flask import Flask\napp = Flask(__name__)\n@app.route(\"/\")\ndef home():\n    return \"<h1>Hello from Docker Compose!</h1>\"\nif __name__ == \"__main__\":\n    app.run(host=\"0.0.0.0\", port=5000)"),
        ('BL', ''),
        ('B', "File: docker-compose.yml"),
        ('C', "services:\n  web:\n    build: .\n    ports:\n      - \"5001:5000\""),
    ])
    kc_note(slide, "day1-05-docker-compose")
    add_footer(slide, p)

    # 51. Lab 10 — Commands
    p += 1
    slide = content_slide(prs, "LAB 10 — DOCKER COMPOSE", "Single Service Compose — Commands")
    add_lab_body_lines(slide, [
        ('C', "mkdir lab10 && cd lab10\n# Create app.py, requirements.txt, Dockerfile, docker-compose.yml"),
        ('B', "Start the service:"),
        ('C', "docker-compose up -d\ndocker-compose ps\ncurl http://localhost:5001"),
        ('B', "Inspect logs and run commands inside:"),
        ('C', "docker-compose logs\ndocker-compose exec web python -c \"print('Running inside container')\""),
        ('B', "Stop and remove:"),
        ('C', "docker-compose down && cd .."),
    ])
    kc_note(slide, "day1-05-docker-compose")
    add_footer(slide, p)

    # Multi-service communication concept slide
    p += 1
    slide = content_slide(prs, "DOCKER COMPOSE", "Multi-Service Communication")
    add_body_lines(slide, [
        ('T', "When you define multiple services in a docker-compose.yml, Compose automatically "
              "creates a shared network and assigns each service a DNS hostname equal to its service name."),
        ('BL', ''),
        ('B', "How services find each other:"),
        ('C', "# In docker-compose.yml:\nservices:\n  web:\n    ...\n  redis:\n    image: redis:7-alpine\n\n# In app.py — use the SERVICE NAME as hostname:\nr = redis.Redis(host=\"redis\", port=6379)"),
        ('BL', ''),
        ('B', "depends_on — controls startup order:"),
        ('T', "  depends_on: [redis]   starts web only after redis container is created"),
        ('T', "  Important: 'created' is NOT the same as 'ready'. Use healthchecks for true readiness."),
        ('BL', ''),
        ('I', "Service names become DNS hostnames — no IP addresses, no environment variables needed."),
    ])
    add_footer(slide, p)

    # Lab 11 header
    p += 1
    lab_header_slide(prs, "LAB 11", "Flask + Redis",
                     "day1-05-docker-compose", p)

    # 53. Lab 11 — Files
    p += 1
    slide = content_slide(prs, "LAB 11 — DOCKER COMPOSE", "Flask + Redis — Files")
    add_lab_body_lines(slide, [
        ('B', "File: app.py"),
        ('C', "from flask import Flask\nimport redis\napp = Flask(__name__)\nr = redis.Redis(host=\"redis\", port=6379)\n@app.route(\"/\")\ndef home():\n    count = r.incr(\"visits\")\n    return f\"<h1>Visit count: {count}</h1>\""),
        ('BL', ''),
        ('B', "File: docker-compose.yml"),
        ('C', "services:\n  web:\n    build: .\n    ports:\n      - \"5001:5000\"\n    depends_on:\n      - redis\n  redis:\n    image: redis:7-alpine"),
    ])
    kc_note(slide, "day1-05-docker-compose")
    add_footer(slide, p)

    # 54. Lab 11 — Commands
    p += 1
    slide = content_slide(prs, "LAB 11 — DOCKER COMPOSE", "Flask + Redis — Commands")
    add_lab_body_lines(slide, [
        ('C', "mkdir lab11 && cd lab11\n# Create app.py, requirements.txt, Dockerfile, docker-compose.yml"),
        ('B', "Start services:"),
        ('C', "docker-compose up -d\ndocker-compose ps"),
        ('B', "Test the visit counter (increments each request):"),
        ('C', "curl http://localhost:5001      # Visit count: 1\ncurl http://localhost:5001      # Visit count: 2\ncurl http://localhost:5001      # Visit count: 3"),
        ('B', "Stop and remove:"),
        ('C', "docker-compose down && cd .."),
        ('BL', ''),
        ('I', "Redis hostname in app.py = service name 'redis'. Compose handles the DNS automatically."),
    ])
    kc_note(slide, "day1-05-docker-compose")
    add_footer(slide, p)

    # Health checks & startup order concept slide
    p += 1
    slide = content_slide(prs, "DOCKER COMPOSE", "Health Checks & Startup Dependencies")
    add_body_lines(slide, [
        ('T', "Databases (PostgreSQL, MySQL) take several seconds to be truly ready to accept connections "
              "after their container starts. Without proper ordering, your app will crash on startup."),
        ('BL', ''),
        ('B', "Solution: healthcheck + condition: service_healthy"),
        ('C', "db:\n  image: postgres:16-alpine\n  healthcheck:\n    test: [\"CMD-SHELL\", \"pg_isready -U user -d mydb\"]\n    interval: 5s\n    retries: 5\n\nweb:\n  depends_on:\n    db:\n      condition: service_healthy   # waits until pg_isready passes"),
        ('BL', ''),
        ('B', "Three depends_on conditions:"),
        ('T', "  service_started   — container is running (default, not truly ready)"),
        ('T', "  service_healthy   — healthcheck is passing (database is accepting connections)"),
        ('T', "  service_completed_successfully — for one-shot init containers"),
        ('BL', ''),
        ('I', "Always use service_healthy for databases and message queues in full-stack Compose apps."),
    ])
    add_footer(slide, p)

    # Lab 12 header
    p += 1
    lab_header_slide(prs, "LAB 12", "Full-Stack — Web + PostgreSQL + Redis",
                     "day1-05-docker-compose", p)

    # 56. Lab 12 — docker-compose.yml
    p += 1
    slide = content_slide(prs, "LAB 12 — DOCKER COMPOSE", "Full-Stack — docker-compose.yml")
    add_lab_body_lines(slide, [
        ('C', "services:\n  web:\n    build: .\n    ports:\n      - \"3000:3000\"\n    environment:\n      DATABASE_URL: postgresql://user:pass@db:5432/mydb\n      REDIS_URL: redis://redis:6379\n    depends_on:\n      db:\n        condition: service_healthy\n      redis:\n        condition: service_started\n  db:\n    image: postgres:16-alpine\n    environment:\n      POSTGRES_USER: user\n      POSTGRES_PASSWORD: pass\n      POSTGRES_DB: mydb\n    healthcheck:\n      test: [\"CMD-SHELL\", \"pg_isready -U user -d mydb\"]\n      interval: 5s\n      retries: 5\n    volumes:\n      - pgdata:/var/lib/postgresql/data\n  redis:\n    image: redis:7-alpine\nvolumes:\n  pgdata:"),
    ])
    kc_note(slide, "day1-05-docker-compose")
    add_footer(slide, p)

    # 57. Lab 12 — Commands & endpoints
    p += 1
    slide = content_slide(prs, "LAB 12 — DOCKER COMPOSE", "Full-Stack — Commands & Endpoints")
    add_lab_body_lines(slide, [
        ('B', "The Node.js app exposes three endpoints:"),
        ('C', "GET /        -> visit counter (Redis incr)\nGET /db      -> current PostgreSQL timestamp\nGET /health  -> {\"status\": \"ok\"}"),
        ('B', "Start and test:"),
        ('C', "mkdir lab12 && cd lab12\n# Create app.js, package.json, Dockerfile, docker-compose.yml\ndocker-compose up -d\ndocker-compose ps\ncurl http://localhost:3000\ncurl http://localhost:3000/db\ncurl http://localhost:3000/health"),
        ('B', "Stop and remove (including volumes):"),
        ('C', "docker-compose down -v && cd .."),
        ('I', "depends_on with service_healthy ensures web waits for PostgreSQL to be ready."),
    ])
    kc_note(slide, "day1-05-docker-compose")
    add_footer(slide, p)

    # 58. Knowledge Check 2
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_AMBER)
    add_section_label(slide, "KNOWLEDGE CHECK  2", color=C_DARK)
    add_main_title(slide, "Quick Review — Storage, Networking & Compose", color=C_DARK)
    add_body_lines(slide, [
        ('B', "Q1.  Which volume type is fully managed by Docker?"),
        ('C', "A)  Bind mount        B)  Named volume        C)  tmpfs        D)  Host volume"),
        ('BL', ''),
        ('B', "Q2.  How do containers on a custom bridge network reach each other?"),
        ('BL', ''),
        ('B', "Q3.  What does  -p 8080:80  mean?"),
        ('BL', ''),
        ('B', "Q4.  Which key ensures a service starts only after another is healthy?"),
        ('C', "A)  needs        B)  after        C)  depends_on with condition: service_healthy"),
        ('BL', ''),
        ('B', "Q5.  Which command removes containers AND networks created by docker-compose up?"),
        ('C', "A)  docker-compose stop        B)  docker-compose rm        C)  docker-compose down"),
    ], y=Inches(1.80))
    add_footer(slide, p)

    # 59. Day 1 Summary
    p += 1
    slide = content_slide(prs, "WRAP UP", "Day 1 — Summary")
    add_body_lines(slide, [
        ('B', "What we covered today:"),
        ('T', "  - Docker architecture: client, daemon, registry"),
        ('T', "  - Running and managing containers  (run, ps, stop, rm, logs, cp)"),
        ('T', "  - Building custom images with Dockerfile  (build, layers, caching)"),
        ('T', "  - Persistent storage: named volumes and bind mounts"),
        ('T', "  - Networking: custom bridge networks, container DNS, port mapping"),
        ('T', "  - Runtime configuration: environment variables and env files"),
        ('T', "  - Docker Compose: defining and running multi-service apps"),
        ('BL', ''),
        ('B', "Tomorrow — Kubernetes (Day 2):"),
        ('T', "  Pods, Namespaces, Deployments, Rollouts, Services, Storage, Jobs"),
    ])
    add_footer(slide, p)

    # Closing (Assessment is at end of Day 2 only)
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_DARK)
    add_circle(slide, 11.6, 0.5, 2.3, C_ORANGE, alpha=12)
    add_circle(slide, 12.3, 4.6, 1.7, C_ORANGE, alpha=8)
    add_textbox(slide, "END OF DAY 1",
                Inches(0.90), Inches(1.80), Inches(11.0), Inches(0.40),
                font="Calibri", size=14, bold=True, color=C_ORANGE)
    add_textbox(slide, "Questions?",
                Inches(0.85), Inches(2.30), Inches(11.0), Inches(1.10),
                font="Cambria", size=52, bold=True, color=C_WHITE)
    add_textbox(slide, "See you tomorrow for Day 2 — Kubernetes.",
                Inches(0.90), Inches(3.60), Inches(11.0), Inches(0.50),
                font="Calibri", size=18, color=C_MUTED2)
    add_textbox(slide, "enquiry@tertiaryinfotech.com   |   +65 6100 0613",
                Inches(0.90), Inches(4.50), Inches(11.0), Inches(0.35),
                font="Calibri", size=14, color=C_MUTED)
    add_footer(slide, p)

    return p


# ══════════════════════════════════════════════════════════════════════════════
#  Slide definitions — Day 2 (Kubernetes)
# ══════════════════════════════════════════════════════════════════════════════

def build_day2(prs, start_page=0):
    p = start_page

    # ── DAY 2 TRANSITION SLIDE ────────────────────────────────────────────────
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_DARK)
    add_circle(slide, 11.8, 0.6, 2.3, C_ORANGE, alpha=15)
    add_circle(slide, 12.5, 4.7, 1.7, C_ORANGE, alpha=10)
    add_textbox(slide, "DAY 2",
                Inches(0.90), Inches(1.80), Inches(11.0), Inches(0.50),
                font="Calibri", size=18, bold=True, color=C_ORANGE)
    add_textbox(slide, "Kubernetes",
                Inches(0.85), Inches(2.40), Inches(11.0), Inches(1.40),
                font="Cambria", size=52, bold=True, color=C_WHITE)
    add_textbox(slide, "Container orchestration — deploy, scale, and heal your apps.",
                Inches(0.90), Inches(4.00), Inches(11.0), Inches(0.50),
                font="Calibri", size=18, color=C_MUTED2)
    add_footer(slide, p)

    # ── Day 2 Course Outline ──────────────────────────────────────────────────
    p += 1
    slide = content_slide(prs, "ROADMAP", "Day 2 — Course Outline")
    left_topics = [
        "1.  What is Kubernetes? — Architecture & objects",
        "2.  kubectl — The K8s CLI",
        "3.  Pods — Lifecycle, YAML, anatomy  (Lab 13)",
        "4.  Namespaces — Environment isolation  (Lab 14)",
        "5.  Deployments — Self-healing, scaling  (Lab 15)",
    ]
    right_topics = [
        "6.  Rolling Updates & Rollbacks  (Lab 16)",
        "7.  Services — ClusterIP, NodePort  (Lab 17)",
        "8.  Storage — emptyDir, PV/PVC  (Lab 18)",
        "9.  Jobs & CronJobs  (Lab 19)",
        "10. Knowledge Checks  •  Assessment",
    ]
    for i, text in enumerate(left_topics):
        add_textbox(slide, text, Inches(0.70), Inches(1.95 + i * 0.72), Inches(6.0), Inches(0.55),
                    font="Calibri", size=15, color=C_DARK)
    for i, text in enumerate(right_topics):
        add_textbox(slide, text, Inches(6.90), Inches(1.95 + i * 0.72), Inches(6.0), Inches(0.55),
                    font="Calibri", size=15, color=C_DARK)
    # Bottom summary bar
    add_rect(slide, Inches(0.22), Inches(6.44), Inches(12.88), Inches(0.52), C_SUMMARY_BAR)
    add_textbox(slide,
                "7 Labs (Lab 13–19) + 2 Knowledge Checks  |  Assessment on final day"
                "  |  Hands-on via KillerCoda browser environment",
                Inches(0.40), Inches(6.51), Inches(12.52), Inches(0.38),
                font="Calibri", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_footer(slide, p)

    # ══ SECTION 1 — WHAT IS KUBERNETES? ═══════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 1", "What is Kubernetes?",
                    subtitle="Container orchestration at scale",
                    subtopics=["Architecture overview", "Control plane", "Worker nodes",
                               "kubectl & Core objects", "K8s vs Docker Compose"],
                    page=p)

    # 4. The Problem
    p += 1
    slide = content_slide(prs, "KUBERNETES — OVERVIEW", "The Container Management Problem")
    add_body_lines(slide, [
        ('T', "Docker is excellent for running a single container on one host. But in production:"),
        ('BL', ''),
        ('B', "Challenges at scale:"),
        ('T', "  - Hundreds of containers across dozens of servers"),
        ('T', "  - A container crashes — who restarts it?"),
        ('T', "  - Traffic spikes — who scales the app up?"),
        ('T', "  - New version released — how to update with zero downtime?"),
        ('T', "  - Which host has enough memory to run the next container?"),
        ('BL', ''),
        ('I', "Kubernetes (K8s) automates all of this — deployment, scaling, healing, and networking."),
    ])
    add_footer(slide, p)

    # 5. What is K8s?
    p += 1
    slide = content_slide(prs, "KUBERNETES — OVERVIEW", "What is Kubernetes?")
    add_body_lines(slide, [
        ('T', "Kubernetes is a portable, extensible, open source platform for managing containerised "
              "workloads and services that facilitate both declarative configuration and automation. "
              "Google open-sourced it in 2014, combining 15+ years of production experience. "
              "(kubernetes.io/docs/concepts/overview)"),
        ('BL', ''),
        ('B', "Key capabilities   (kubernetes.io):"),
        ('T', "  - Service discovery & load balancing  — expose via DNS or IP; distribute traffic"),
        ('T', "  - Storage orchestration  — automatically mount local or cloud storage"),
        ('T', "  - Automated rollouts & rollbacks  — control state changes at a managed rate"),
        ('T', "  - Self-healing  — restarts failed containers; replaces non-responsive ones"),
        ('T', "  - Horizontal scaling  — scale via command, UI, or automatic CPU-based triggers"),
        ('T', "  - Secret & config management  — store passwords, tokens, keys securely"),
        ('BL', ''),
        ('S', "Used by: Google (GKE), Microsoft Azure (AKS), Amazon (EKS), and virtually every large-scale cloud deployment."),
    ])
    add_footer(slide, p)

    # 6. Architecture Overview
    p += 1
    slide = content_slide(prs, "KUBERNETES — ARCHITECTURE", "K8s Cluster Architecture")
    add_body_lines(slide, [
        ('B', "A Kubernetes cluster has two types of machines:"),
        ('BL', ''),
        ('B', "Control Plane  (Master Node)"),
        ('T', "  Manages the cluster — scheduling, desired state, API."),
        ('T', "  Components: API Server, etcd, Scheduler, Controller Manager"),
        ('BL', ''),
        ('B', "Worker Nodes"),
        ('T', "  Run the actual application containers inside Pods."),
        ('T', "  Components: kubelet, kube-proxy, container runtime (containerd)"),
        ('BL', ''),
        ('I', "kubectl talks to the API Server -> API Server stores state in etcd -> Scheduler assigns Pods to Nodes -> kubelet starts containers."),
    ])
    add_footer(slide, p)

    # K8s Architecture visual diagram
    p += 1
    slide_k8s_arch_diagram(prs, p)

    # 7. Control Plane
    p += 1
    slide = content_slide(prs, "KUBERNETES — ARCHITECTURE", "Control Plane Components")
    add_body_lines(slide, [
        ('B', "API Server  (kube-apiserver)"),
        ('T', "  The front door to K8s. Every kubectl command is an API call here."),
        ('BL', ''),
        ('B', "etcd"),
        ('T', "  Distributed key-value store. The single source of truth for all cluster state."),
        ('BL', ''),
        ('B', "Scheduler  (kube-scheduler)"),
        ('T', "  Watches for unscheduled Pods and assigns them to a Node based on resources."),
        ('BL', ''),
        ('B', "Controller Manager  (kube-controller-manager)"),
        ('T', "  Runs controllers that reconcile actual state with desired state "
              "(e.g. ReplicaSet controller ensures N replicas are always running)."),
    ])
    add_footer(slide, p)

    # 8. Worker Node
    p += 1
    slide = content_slide(prs, "KUBERNETES — ARCHITECTURE", "Worker Node Components")
    add_body_lines(slide, [
        ('B', "kubelet"),
        ('T', "  Agent running on every node. Receives Pod specs from the API Server "
              "and ensures containers are running and healthy."),
        ('BL', ''),
        ('B', "kube-proxy"),
        ('T', "  Maintains network rules on the node. Implements the Service abstraction "
              "by routing traffic to the right Pod IPs."),
        ('BL', ''),
        ('B', "Container Runtime"),
        ('T', "  Software that actually runs containers. Default is containerd (replaces Docker in K8s 1.24+)."),
        ('BL', ''),
        ('I', "Every worker node runs all three: kubelet + kube-proxy + container runtime."),
    ])
    add_footer(slide, p)

    # 9. kubectl
    p += 1
    slide = content_slide(prs, "KUBERNETES — TOOLS", "kubectl — The K8s CLI")
    add_body_lines(slide, [
        ('T', "kubectl is the command-line tool for interacting with any Kubernetes cluster."),
        ('BL', ''),
        ('B', "Essential commands:"),
        ('C', "kubectl get pods              # list pods\nkubectl get pods -o wide      # with node and IP info\nkubectl describe pod <name>   # detailed info + events\nkubectl apply -f file.yaml    # create/update from YAML\nkubectl delete -f file.yaml   # delete resources\nkubectl logs <pod>            # view container logs\nkubectl exec -it <pod> -- sh  # open shell in pod"),
        ('BL', ''),
        ('I', "kubectl --help  or  kubectl <command> --help  shows all options."),
    ])
    add_footer(slide, p)

    # 10. Core Objects
    p += 1
    slide = content_slide(prs, "KUBERNETES — OBJECTS", "Core Kubernetes Objects")
    add_body_lines(slide, [
        ('B', "Pod"),
        ('T', "  Smallest deployable unit. Wraps one or more containers sharing network and storage."),
        ('B', "Deployment"),
        ('T', "  Manages a set of identical Pods. Handles scaling, updates, and self-healing."),
        ('B', "Service"),
        ('T', "  Stable network endpoint for a group of Pods. Provides load balancing and DNS."),
        ('B', "Namespace"),
        ('T', "  Virtual cluster for isolating resources (e.g. dev, staging, prod environments)."),
        ('B', "PersistentVolume / PVC"),
        ('T', "  Decouple storage lifecycle from Pod lifecycle — data survives Pod restarts."),
    ])
    add_footer(slide, p)

    # K8s Objects visual diagram (Service / Deployment / Pods)
    p += 1
    slide_k8s_objects_diagram(prs, p)

    # 11. K8s vs Docker Compose
    p += 1
    two_column_slide(prs,
        label="KUBERNETES — OVERVIEW", title="Kubernetes vs Docker Compose",
        left_heading="Docker Compose",
        left_lines=[
            ('T', "Single host only"),
            ('T', "Manual restart on failure"),
            ('T', "No auto-scaling"),
            ('T', "Simple YAML format"),
            ('T', "Great for local dev"),
            ('BL', ''),
            ('S', "Best for: development, testing"),
        ],
        right_heading="Kubernetes",
        right_lines=[
            ('T', "Multi-host cluster"),
            ('T', "Self-healing (auto-restart)"),
            ('T', "Horizontal auto-scaling"),
            ('T', "Rich YAML object model"),
            ('T', "Production-grade"),
            ('BL', ''),
            ('S', "Best for: production, CI/CD"),
        ],
        page=p)

    # 12. Labels & Selectors
    p += 1
    slide = content_slide(prs, "KUBERNETES — CONCEPTS", "Labels & Selectors")
    add_body_lines(slide, [
        ('T', "Labels are key-value pairs attached to K8s objects. Selectors filter objects by label."),
        ('BL', ''),
        ('B', "Labels in YAML:"),
        ('C', "metadata:\n  labels:\n    app: my-app\n    env: production\n    tier: frontend"),
        ('B', "Select by label (kubectl):"),
        ('C', "kubectl get pods -l app=my-app\nkubectl get pods -l env=production,tier=frontend"),
        ('BL', ''),
        ('I', "Services and Deployments use selectors to find their Pods — the glue that connects K8s objects."),
    ])
    add_footer(slide, p)

    # ══ SECTION 2 — PODS & NAMESPACES ═════════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 2", "Pods & Namespaces",
                    subtitle="Labs 13 & 14 — The foundational K8s units",
                    subtopics=["What is a Pod?", "Pod YAML structure", "Pod lifecycle",
                               "Lab 13: Pods (imperative & declarative)",
                               "Lab 14: Namespaces"],
                    page=p)

    # 14. What is a Pod?
    p += 1
    slide = content_slide(prs, "PODS", "What is a Pod?")
    add_body_lines(slide, [
        ('T', "A Pod is the smallest deployable unit in Kubernetes — a wrapper around one or more containers."),
        ('BL', ''),
        ('B', "All containers in a Pod share:"),
        ('T', "  - The same network namespace (same IP, same ports)"),
        ('T', "  - The same storage volumes"),
        ('T', "  - The same lifecycle (start and stop together)"),
        ('BL', ''),
        ('B', "Pods are ephemeral — do NOT count on them living forever:"),
        ('T', "  - Delete a Pod manually: it is gone."),
        ('T', "  - Node fails: Pods on that node are gone."),
        ('T', "  - Use a Deployment to ensure Pods are recreated automatically."),
    ])
    add_footer(slide, p)

    # Pod anatomy visual diagram
    p += 1
    slide_k8s_pod_anatomy_diagram(prs, p)

    # 15. Pod YAML
    p += 1
    slide = content_slide(prs, "PODS — YAML", "Pod YAML Structure")
    add_body_lines(slide, [
        ('C', "apiVersion: v1\nkind: Pod\nmetadata:\n  name: my-app\n  labels:\n    app: my-app\nspec:\n  containers:\n    - name: nginx\n      image: nginx\n      ports:\n        - containerPort: 80"),
        ('BL', ''),
        ('B', "Key fields:"),
        ('T', "  apiVersion  — API group (v1 for core objects)"),
        ('T', "  kind        — object type (Pod, Deployment, Service ...)"),
        ('T', "  metadata    — name + labels"),
        ('T', "  spec        — desired state (containers, volumes, ...)"),
    ])
    add_footer(slide, p)

    # 16. Pod Lifecycle
    p += 1
    slide = content_slide(prs, "PODS — LIFECYCLE", "Pod Phase & States")
    add_body_lines(slide, [
        ('B', "Pod phases (kubectl get pods STATUS column):"),
        ('BL', ''),
        ('T', "  Pending     — scheduled but containers not yet started"),
        ('T', "  Running     — at least one container is running"),
        ('T', "  Succeeded   — all containers exited with code 0"),
        ('T', "  Failed      — at least one container exited non-zero"),
        ('T', "  Unknown     — cannot reach the node"),
        ('BL', ''),
        ('B', "Container states (kubectl describe pod):"),
        ('T', "  Waiting  |  Running  |  Terminated"),
        ('BL', ''),
        ('I', "CrashLoopBackOff means the container keeps crashing and K8s is waiting before restarting."),
    ])
    add_footer(slide, p)

    # 17. Lab 13 header
    p += 1
    lab_header_slide(prs, "LAB 13", "Pods — Imperative & Declarative",
                     "day2-01-k8s-pods-namespaces", p)

    # 18. Lab 13a — Imperative
    p += 1
    slide = content_slide(prs, "LAB 13 — PODS", "Pods — Imperative Commands")
    add_lab_body_lines(slide, [
        ('B', "Run and inspect a pod:"),
        ('C', "kubectl run my-nginx --image=nginx\nkubectl get pods\nkubectl get pods -o wide\nkubectl describe pod my-nginx"),
        ('B', "Logs and exec:"),
        ('C', "kubectl logs my-nginx\nkubectl exec my-nginx -- cat /etc/hostname\nkubectl exec -it my-nginx -- /bin/sh\n# Inside: ls /usr/share/nginx/html && exit"),
        ('B', "Delete and verify:"),
        ('C', "kubectl delete pod my-nginx\nkubectl get pods"),
        ('BL', ''),
        ('I', "Key concept: Pods are mortal. Delete one and it is gone. Deployments keep Pods alive."),
    ])
    kc_note(slide, "day2-01-k8s-pods-namespaces")
    add_footer(slide, p)

    # 19. Lab 13b — Declarative
    p += 1
    slide = content_slide(prs, "LAB 13 — PODS", "Pods — Declarative YAML")
    add_lab_body_lines(slide, [
        ('B', "File: pod.yaml"),
        ('C', "apiVersion: v1\nkind: Pod\nmetadata:\n  name: my-app\n  labels:\n    app: my-app\nspec:\n  containers:\n    - name: nginx\n      image: nginx\n      ports:\n        - containerPort: 80"),
        ('B', "Apply, inspect, delete:"),
        ('C', "kubectl apply -f pod.yaml\nkubectl get pods\nkubectl describe pod my-app\nkubectl delete -f pod.yaml"),
        ('BL', ''),
        ('I', "kubectl apply -f is idempotent — safe to run multiple times. Always prefer YAML for production."),
    ])
    kc_note(slide, "day2-01-k8s-pods-namespaces")
    add_footer(slide, p)

    # 20. Namespaces theory
    p += 1
    slide = content_slide(prs, "NAMESPACES", "What is a Namespace?")
    add_body_lines(slide, [
        ('T', "A Namespace is a virtual cluster inside your K8s cluster. "
              "Resources in different namespaces are isolated by default."),
        ('BL', ''),
        ('B', "Default namespaces:"),
        ('T', "  default      — where objects go if you don't specify a namespace"),
        ('T', "  kube-system  — internal K8s components (DNS, proxy, scheduler)"),
        ('T', "  kube-public  — publicly readable cluster info"),
        ('BL', ''),
        ('B', "Why use namespaces?"),
        ('T', "  - Separate environments: dev / staging / prod"),
        ('T', "  - Team isolation: team-a / team-b"),
        ('T', "  - Apply quotas and policies per namespace"),
    ])
    add_footer(slide, p)

    # 21. Lab 14 header
    p += 1
    lab_header_slide(prs, "LAB 14", "Namespaces",
                     "day2-01-k8s-pods-namespaces", p)

    # 22. Lab 14 — Commands
    p += 1
    slide = content_slide(prs, "LAB 14 — NAMESPACES", "Namespaces — Commands")
    add_lab_body_lines(slide, [
        ('B', "Create namespace and run a pod inside it:"),
        ('C', "kubectl get namespaces\nkubectl create namespace dev\nkubectl run my-nginx --image=nginx -n dev\nkubectl get pods -n dev\nkubectl get pods --all-namespaces"),
        ('B', "Delete namespace (removes everything inside!):"),
        ('C', "kubectl delete namespace dev"),
        ('B', "Declarative namespace + pod:"),
        ('C', "# namespace.yaml  ->  kind: Namespace / name: staging\n# pod-staging.yaml ->  namespace: staging in metadata\nkubectl apply -f namespace.yaml\nkubectl apply -f pod-staging.yaml\nkubectl get pods -n staging\nkubectl delete -f namespace.yaml"),
        ('I', "Deleting a namespace removes everything inside it — use with care in production."),
    ])
    kc_note(slide, "day2-01-k8s-pods-namespaces")
    add_footer(slide, p)

    # ══ SECTION 3 — DEPLOYMENTS ════════════════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 3", "Deployments & Scaling",
                    subtitle="Lab 15 — Self-healing, replicas, YAML",
                    subtopics=["Why Deployments?", "Naked Pod vs Deployment",
                               "ReplicaSets", "Lab 15: Create & scale"],
                    page=p)

    # 24. Deployment theory
    p += 1
    slide = content_slide(prs, "DEPLOYMENTS", "What is a Deployment?")
    add_body_lines(slide, [
        ('T', "A Deployment is the standard way to run stateless applications in Kubernetes. "
              "It manages a ReplicaSet which in turn manages Pods."),
        ('BL', ''),
        ('B', "What a Deployment gives you:"),
        ('T', "  - Declare desired count: \"I want 3 replicas of this app\""),
        ('T', "  - Self-healing: if a Pod dies, a new one is created automatically"),
        ('T', "  - Scaling: change replicas up or down with one command"),
        ('T', "  - Rolling updates: update the image with zero downtime"),
        ('T', "  - Rollback: revert to a previous version instantly"),
        ('BL', ''),
        ('I', "Rule: never run a naked Pod in production. Always use a Deployment."),
    ])
    add_footer(slide, p)

    # 25. Deployment vs Pod
    p += 1
    two_column_slide(prs,
        label="DEPLOYMENTS", title="Naked Pod vs Deployment",
        left_heading="Naked Pod",
        left_lines=[
            ('T', "Runs one container instance"),
            ('T', "Dies and stays dead"),
            ('T', "Manual recreation"),
            ('T', "No rolling update"),
            ('T', "No scaling"),
            ('BL', ''),
            ('S', "Only for: debugging, one-shot tasks"),
        ],
        right_heading="Deployment",
        right_lines=[
            ('T', "Manages N replicas"),
            ('T', "Self-heals on failure"),
            ('T', "Automatic recreation"),
            ('T', "Rolling update built-in"),
            ('T', "Scale with one command"),
            ('BL', ''),
            ('S', "Always use for production apps"),
        ],
        page=p)

    # 26. ReplicaSet
    p += 1
    slide = content_slide(prs, "DEPLOYMENTS — INTERNALS", "ReplicaSet & Self-Healing")
    add_body_lines(slide, [
        ('T', "A Deployment creates and manages a ReplicaSet. The ReplicaSet watches your Pods and ensures "
              "the desired count is always running."),
        ('BL', ''),
        ('B', "Self-healing loop:"),
        ('T', "  1.  ReplicaSet watches Pods with matching labels"),
        ('T', "  2.  Pod count drops below desired  ->  create new Pod"),
        ('T', "  3.  Pod count exceeds desired  ->  delete extra Pod"),
        ('BL', ''),
        ('B', "View the ReplicaSet:"),
        ('C', "kubectl get replicasets\nkubectl describe replicaset <name>"),
        ('BL', ''),
        ('I', "You rarely interact with ReplicaSets directly — manage them through Deployments."),
    ])
    add_footer(slide, p)

    # ReplicaSet self-healing visual diagram
    p += 1
    slide_k8s_replicaset_diagram(prs, p)

    # 27. Lab 15 header
    p += 1
    lab_header_slide(prs, "LAB 15", "Deployments — Imperative & Declarative",
                     "day2-02-k8s-deployments", p)

    # 28. Lab 15a — Imperative
    p += 1
    slide = content_slide(prs, "LAB 15 — DEPLOYMENTS", "Deployments — Imperative Commands")
    add_lab_body_lines(slide, [
        ('B', "Create and inspect:"),
        ('C', "kubectl create deployment my-nginx --image=nginx\nkubectl get deployments\nkubectl get pods\nkubectl describe deployment my-nginx"),
        ('B', "Scale and watch:"),
        ('C', "kubectl scale deployment my-nginx --replicas=3\nkubectl get pods -w   # Ctrl+C when all Running"),
        ('B', "Self-healing: delete a pod and watch it recreate:"),
        ('C', "kubectl delete pod $(kubectl get pods -l app=my-nginx -o name | head -1 | cut -d/ -f2)\nkubectl get pods"),
        ('B', "Scale down and clean up:"),
        ('C', "kubectl scale deployment my-nginx --replicas=1\nkubectl delete deployment my-nginx"),
        ('I', "Key concept: Deployments use a ReplicaSet to enforce desired count. Pod deletion triggers recreation."),
    ])
    kc_note(slide, "day2-02-k8s-deployments")
    add_footer(slide, p)

    # 29. Lab 15b — Declarative
    p += 1
    slide = content_slide(prs, "LAB 15 — DEPLOYMENTS", "Deployments — Declarative YAML")
    add_lab_body_lines(slide, [
        ('B', "File: deployment.yaml"),
        ('C', "apiVersion: apps/v1\nkind: Deployment\nmetadata:\n  name: my-app\nspec:\n  replicas: 3\n  selector:\n    matchLabels:\n      app: my-app\n  template:\n    metadata:\n      labels:\n        app: my-app\n    spec:\n      containers:\n        - name: nginx\n          image: nginx\n          ports:\n            - containerPort: 80"),
        ('B', "Apply and scale via YAML:"),
        ('C', "kubectl apply -f deployment.yaml\nkubectl get deployments && kubectl get pods\n# Edit replicas: 5 in YAML then:\nkubectl apply -f deployment.yaml\nkubectl delete -f deployment.yaml"),
    ])
    kc_note(slide, "day2-02-k8s-deployments")
    add_footer(slide, p)

    # 30. Knowledge Check 1
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_AMBER)
    add_section_label(slide, "KNOWLEDGE CHECK  1", color=C_DARK)
    add_main_title(slide, "Quick Review — K8s Basics, Pods & Deployments", color=C_DARK)
    add_body_lines(slide, [
        ('B', "Q1.  Which control plane component is responsible for assigning Pods to nodes?"),
        ('C', "A) API Server     B) etcd     C) Scheduler     D) Controller Manager"),
        ('BL', ''),
        ('B', "Q2.  What is the difference between a Pod and a Deployment?"),
        ('BL', ''),
        ('B', "Q3.  True / False:  kubectl apply -f is safe to run multiple times on the same file."),
        ('BL', ''),
        ('B', "Q4.  Which command shows pods in ALL namespaces?"),
        ('C', "A) kubectl get pods -n all     B) kubectl get pods --all-namespaces     C) kubectl get all pods"),
        ('BL', ''),
        ('B', "Q5.  What happens when you delete a namespace?"),
    ], y=Inches(1.80))
    add_footer(slide, p)

    # ══ SECTION 4 — ROLLING UPDATES & ROLLBACKS ═══════════════════════════════
    p += 1
    section_divider(prs, "SECTION 4", "Rolling Updates & Rollbacks",
                    subtitle="Lab 16 — Zero-downtime updates and instant revert",
                    subtopics=["Rolling vs Recreate", "How rolling update works",
                               "Lab 16a: Update image", "Lab 16b: Rollback"],
                    page=p)

    # 32. Update Strategies
    p += 1
    two_column_slide(prs,
        label="ROLLOUTS — STRATEGIES", title="Update Strategies",
        left_heading="Rolling Update  (default)",
        left_lines=[
            ('T', "Replace Pods gradually"),
            ('T', "Old and new run simultaneously"),
            ('T', "Zero downtime"),
            ('T', "Slower — controlled pace"),
            ('BL', ''),
            ('C', "strategy:\n  type: RollingUpdate\n  maxSurge: 1\n  maxUnavailable: 0"),
        ],
        right_heading="Recreate",
        right_lines=[
            ('T', "Stop all old Pods first"),
            ('T', "Then start new Pods"),
            ('T', "Brief downtime"),
            ('T', "Faster — all-or-nothing"),
            ('BL', ''),
            ('C', "strategy:\n  type: Recreate"),
        ],
        page=p)

    # 33. How Rolling Update works
    p += 1
    slide = content_slide(prs, "ROLLOUTS — MECHANISM", "How Rolling Updates Work")
    add_body_lines(slide, [
        ('T', "When you update the image in a Deployment, K8s creates a new ReplicaSet and gradually "
              "shifts Pods from the old RS to the new one."),
        ('BL', ''),
        ('B', "Step-by-step:"),
        ('T', "  1.  New ReplicaSet created with 0 replicas"),
        ('T', "  2.  Scale new RS up by 1, scale old RS down by 1"),
        ('T', "  3.  Repeat until new RS = desired count, old RS = 0"),
        ('BL', ''),
        ('B', "Track live progress:"),
        ('C', "kubectl rollout status deployment my-nginx"),
        ('BL', ''),
        ('I', "K8s keeps old ReplicaSets as rollback history. Default: 10 revisions kept."),
    ])
    add_footer(slide, p)

    # Rolling update step-by-step visual diagram
    p += 1
    slide_k8s_rolling_update_diagram(prs, p)

    # 34. Lab 16 header
    p += 1
    lab_header_slide(prs, "LAB 16", "Rolling Updates & Rollbacks",
                     "day2-03-k8s-rollouts", p)

    # 35. Lab 16a — Rolling Update
    p += 1
    slide = content_slide(prs, "LAB 16 — ROLLOUTS", "Rolling Update — Commands")
    add_lab_body_lines(slide, [
        ('B', "Deploy nginx 1.24, then update to 1.25:"),
        ('C', "kubectl create deployment my-nginx --image=nginx:1.24 --replicas=3\nkubectl get deployments\nkubectl describe deployment my-nginx | grep Image"),
        ('B', "Update image and track rollout:"),
        ('C', "kubectl set image deployment my-nginx nginx=nginx:1.25\nkubectl rollout status deployment my-nginx\nkubectl describe deployment my-nginx | grep Image"),
        ('B', "Update again to 1.27 and check history:"),
        ('C', "kubectl set image deployment my-nginx nginx=nginx:1.27\nkubectl rollout status deployment my-nginx\nkubectl rollout history deployment my-nginx"),
    ])
    kc_note(slide, "day2-03-k8s-rollouts")
    add_footer(slide, p)

    # 36. Lab 16b — Rollback
    p += 1
    slide = content_slide(prs, "LAB 16 — ROLLOUTS", "Rollback — Commands")
    add_lab_body_lines(slide, [
        ('B', "Rollback to previous revision (nginx 1.25):"),
        ('C', "kubectl rollout undo deployment my-nginx\nkubectl describe deployment my-nginx | grep Image"),
        ('B', "Rollback to a specific revision (nginx 1.24):"),
        ('C', "kubectl rollout history deployment my-nginx\nkubectl rollout undo deployment my-nginx --to-revision=1\nkubectl describe deployment my-nginx | grep Image"),
        ('B', "Declarative approach — edit deployment.yaml and apply:"),
        ('C', "# Change image: nginx:1.24 -> nginx:1.25 in deployment.yaml\nkubectl apply -f deployment.yaml\nkubectl rollout status deployment my-app\nkubectl rollout undo deployment my-app\nkubectl delete deployment my-nginx\nkubectl delete -f deployment.yaml"),
        ('I', "rollout undo is instant — K8s switches to the previous ReplicaSet, no rebuild."),
    ])
    kc_note(slide, "day2-03-k8s-rollouts")
    add_footer(slide, p)

    # ══ SECTION 5 — SERVICES ══════════════════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 5", "Services & Networking",
                    subtitle="Lab 17 — ClusterIP and NodePort",
                    subtopics=["Why Services?", "ClusterIP", "NodePort",
                               "Lab 17a: ClusterIP", "Lab 17b: NodePort",
                               "Lab 17c: Service YAML"],
                    page=p)

    # 38. Why Services?
    p += 1
    slide = content_slide(prs, "SERVICES", "Why Do We Need Services?")
    add_body_lines(slide, [
        ('T', "Pod IPs are ephemeral. Every time a Pod is recreated (crash, update, scale), it gets a new IP."),
        ('BL', ''),
        ('B', "Problem without Services:"),
        ('T', "  - Your frontend must discover the backend IP dynamically"),
        ('T', "  - 3 backend replicas = 3 different IPs to track"),
        ('T', "  - Rolling update changes IPs mid-request"),
        ('BL', ''),
        ('B', "What a Service provides:"),
        ('T', "  - Stable virtual IP (ClusterIP) that never changes"),
        ('T', "  - DNS name: <service-name>.<namespace>.svc.cluster.local"),
        ('T', "  - Load balancing across all matching Pods"),
        ('BL', ''),
        ('I', "Services select Pods via labels — they automatically track which Pods are healthy."),
    ])
    add_footer(slide, p)

    # 39. Service Types
    p += 1
    slide = content_slide(prs, "SERVICES — TYPES", "Service Types Overview")
    add_body_lines(slide, [
        ('B', "ClusterIP  (default)"),
        ('T', "  Internal-only. Reachable only from within the cluster. Best for service-to-service calls."),
        ('BL', ''),
        ('B', "NodePort"),
        ('T', "  Exposes service on each node's IP at a static port (30000-32767). "
              "Reachable from outside the cluster. Good for testing."),
        ('BL', ''),
        ('B', "LoadBalancer"),
        ('T', "  Creates an external cloud load balancer (AWS ELB, GCP LB). "
              "Production use on managed K8s (AKS, EKS, GKE)."),
        ('BL', ''),
        ('B', "ExternalName"),
        ('T', "  Maps a service to a DNS name outside the cluster (e.g. an external database)."),
    ])
    add_footer(slide, p)

    # Service types visual diagram
    p += 1
    slide_k8s_service_types_diagram(prs, p)

    # 40. Lab 17 header
    p += 1
    lab_header_slide(prs, "LAB 17", "Services — ClusterIP & NodePort",
                     "day2-04-k8s-services", p)

    # 41. Lab 17a — ClusterIP
    p += 1
    slide = content_slide(prs, "LAB 17 — SERVICES", "ClusterIP — Internal Service")
    add_lab_body_lines(slide, [
        ('B', "Create a deployment and expose it with ClusterIP:"),
        ('C', "kubectl create deployment my-nginx --image=nginx --replicas=2\nkubectl expose deployment my-nginx --port=80 --target-port=80\nkubectl get services\nkubectl describe service my-nginx\nkubectl get endpoints my-nginx"),
        ('B', "Test from inside the cluster (temp pod):"),
        ('C', "kubectl run test-pod --image=busybox --rm -it --restart=Never \\\n  -- wget -O- my-nginx:80"),
        ('B', "Clean up:"),
        ('C', "kubectl delete service my-nginx\nkubectl delete deployment my-nginx"),
        ('I', "ClusterIP provides a stable virtual IP and DNS name. Service load-balances across all matching Pods."),
    ])
    kc_note(slide, "day2-04-k8s-services")
    add_footer(slide, p)

    # 42. Lab 17b — NodePort
    p += 1
    slide = content_slide(prs, "LAB 17 — SERVICES", "NodePort — External Access")
    add_lab_body_lines(slide, [
        ('B', "Expose deployment as NodePort:"),
        ('C', "kubectl create deployment my-nginx --image=nginx --replicas=2\nkubectl expose deployment my-nginx \\\n  --type=NodePort --port=80 --target-port=80\nkubectl get service my-nginx   # note the 30xxx port"),
        ('B', "Access via node IP and port:"),
        ('C', "kubectl get nodes -o wide\ncurl http://$(kubectl get nodes -o jsonpath='{.items[0].status.addresses[0].address}'):$(kubectl get svc my-nginx -o jsonpath='{.spec.ports[0].nodePort}')"),
        ('B', "Clean up:"),
        ('C', "kubectl delete service my-nginx\nkubectl delete deployment my-nginx"),
    ])
    kc_note(slide, "day2-04-k8s-services")
    add_footer(slide, p)

    # 43. Lab 17c — YAML
    p += 1
    slide = content_slide(prs, "LAB 17 — SERVICES", "Service + Deployment — YAML")
    add_lab_body_lines(slide, [
        ('B', "service.yaml"),
        ('C', "apiVersion: v1\nkind: Service\nmetadata:\n  name: my-app-svc\nspec:\n  type: NodePort\n  selector:\n    app: my-app\n  ports:\n    - port: 80\n      targetPort: 80\n      nodePort: 30080"),
        ('B', "Apply and test:"),
        ('C', "kubectl apply -f deployment.yaml\nkubectl apply -f service.yaml\nkubectl get services\nkubectl get endpoints my-app-svc\ncurl http://localhost:30080\nkubectl delete -f service.yaml\nkubectl delete -f deployment.yaml"),
    ])
    kc_note(slide, "day2-04-k8s-services")
    add_footer(slide, p)

    # 44. Knowledge Check 2
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_AMBER)
    add_section_label(slide, "KNOWLEDGE CHECK  2", color=C_DARK)
    add_main_title(slide, "Quick Review — Rollouts & Services", color=C_DARK)
    add_body_lines(slide, [
        ('B', "Q1.  What is the default Deployment update strategy?"),
        ('C', "A) Recreate     B) RollingUpdate     C) BlueGreen     D) Canary"),
        ('BL', ''),
        ('B', "Q2.  Which command rolls back a Deployment to a specific revision?"),
        ('BL', ''),
        ('B', "Q3.  True / False:  A ClusterIP service is accessible from outside the cluster."),
        ('BL', ''),
        ('B', "Q4.  NodePort values must be in what range?"),
        ('C', "A) 1-1024     B) 8000-9000     C) 30000-32767     D) Any free port"),
        ('BL', ''),
        ('B', "Q5.  How does a Service know which Pods to send traffic to?"),
    ], y=Inches(1.80))
    add_footer(slide, p)

    # ══ SECTION 6 — STORAGE & JOBS ════════════════════════════════════════════
    p += 1
    section_divider(prs, "SECTION 6", "Storage & Jobs",
                    subtitle="Labs 18 & 19 — PV/PVC, Jobs, CronJobs",
                    subtopics=["emptyDir volumes", "PersistentVolume & PVC",
                               "Lab 18: Storage", "Jobs & CronJobs",
                               "Lab 19: Jobs"],
                    page=p)

    # 46. K8s Storage Overview
    p += 1
    slide = content_slide(prs, "KUBERNETES STORAGE", "Storage in Kubernetes")
    add_body_lines(slide, [
        ('T', "Container storage is ephemeral — files written inside a container are lost when it restarts. "
              "K8s offers several volume types to persist data."),
        ('BL', ''),
        ('B', "Volume types used in this course:"),
        ('BL', ''),
        ('B', "emptyDir"),
        ('T', "  Temporary directory shared between containers in the same Pod. "
              "Deleted when the Pod is removed."),
        ('BL', ''),
        ('B', "PersistentVolume  (PV)  +  PersistentVolumeClaim  (PVC)"),
        ('T', "  Decouple storage from Pod lifecycle. Data persists across Pod restarts and deletions."),
    ])
    add_footer(slide, p)

    # 47. emptyDir
    p += 1
    slide = content_slide(prs, "KUBERNETES STORAGE", "emptyDir — Shared Temporary Storage")
    add_body_lines(slide, [
        ('T', "emptyDir is created when a Pod is assigned to a node and deleted when the Pod is removed. "
              "All containers in the Pod can read/write the same directory."),
        ('BL', ''),
        ('B', "Common use cases:"),
        ('T', "  - Share files between a main container and a sidecar"),
        ('T', "  - Cache data that can be recomputed"),
        ('T', "  - Scratch space for batch processing"),
        ('BL', ''),
        ('B', "Key point: emptyDir survives container restarts (crash/liveness probe), "
              "but is deleted when the whole Pod is removed."),
    ])
    add_footer(slide, p)

    # 48. PV & PVC
    p += 1
    slide = content_slide(prs, "KUBERNETES STORAGE", "PersistentVolume & PVC")
    add_body_lines(slide, [
        ('B', "PersistentVolume  (PV)  — cluster resource"),
        ('T', "  A piece of storage provisioned by an admin (hostPath, NFS, cloud disk). "
              "Exists independently of any Pod."),
        ('BL', ''),
        ('B', "PersistentVolumeClaim  (PVC)  — Pod's storage request"),
        ('T', "  A request for storage: \"I need 500Mi with ReadWriteOnce access.\""),
        ('T', "  K8s binds the PVC to a matching PV automatically."),
        ('BL', ''),
        ('B', "Lifecycle:"),
        ('T', "  PV created -> PVC created -> PVC binds to PV -> Pod mounts PVC -> data persists across pod restarts"),
    ])
    add_footer(slide, p)

    # PV → PVC → Pod binding chain visual diagram
    p += 1
    slide_k8s_pv_pvc_diagram(prs, p)

    # 49. Lab 18 header
    p += 1
    lab_header_slide(prs, "LAB 18", "Storage — emptyDir & PV/PVC",
                     "day2-05-k8s-storage-jobs", p)

    # 50. Lab 18a — emptyDir
    p += 1
    slide = content_slide(prs, "LAB 18 — STORAGE", "emptyDir — Shared Pod Volume")
    add_lab_body_lines(slide, [
        ('B', "File: emptydir-pod.yaml — two containers sharing /data:"),
        ('C', "spec:\n  containers:\n    - name: writer\n      image: busybox\n      command: [\"sh\",\"-c\",\"echo 'hello from writer' > /data/message.txt && sleep 3600\"]\n      volumeMounts:\n        - name: shared-data\n          mountPath: /data\n    - name: reader\n      image: busybox\n      command: [\"sh\",\"-c\",\"sleep 5 && cat /data/message.txt && sleep 3600\"]\n      volumeMounts:\n        - name: shared-data\n          mountPath: /data\n  volumes:\n    - name: shared-data\n      emptyDir: {}"),
        ('B', "Apply and verify:"),
        ('C', "kubectl apply -f emptydir-pod.yaml\nkubectl get pod shared-pod\nkubectl exec shared-pod -c reader -- cat /data/message.txt\nkubectl logs shared-pod -c reader\nkubectl delete pod shared-pod"),
    ])
    kc_note(slide, "day2-05-k8s-storage-jobs")
    add_footer(slide, p)

    # 51. Lab 18b — PV & PVC
    p += 1
    slide = content_slide(prs, "LAB 18 — STORAGE", "PV & PVC — Commands")
    add_lab_body_lines(slide, [
        ('B', "Create PV (1Gi hostPath) and PVC (requests 500Mi):"),
        ('C', "kubectl apply -f pv.yaml\nkubectl apply -f pvc.yaml\nkubectl get pv\nkubectl get pvc    # both should show Bound"),
        ('B', "Pod uses PVC:"),
        ('C', "kubectl apply -f pod-with-pvc.yaml\nkubectl exec pvc-pod -- cat /data/file.txt"),
        ('B', "Prove data survives pod deletion:"),
        ('C', "kubectl delete pod pvc-pod\nkubectl apply -f pod-with-pvc.yaml\nkubectl exec pvc-pod -- cat /data/file.txt   # still there!"),
        ('B', "Clean up:"),
        ('C', "kubectl delete pod pvc-pod\nkubectl delete pvc my-pvc\nkubectl delete pv my-pv"),
    ])
    kc_note(slide, "day2-05-k8s-storage-jobs")
    add_footer(slide, p)

    # 52. Jobs theory
    p += 1
    slide = content_slide(prs, "JOBS", "What is a Kubernetes Job?")
    add_body_lines(slide, [
        ('T', "A Job creates one or more Pods and ensures they run to successful completion. "
              "Unlike a Deployment, a Job does not restart Pods after they succeed."),
        ('BL', ''),
        ('B', "Key Job parameters:"),
        ('T', "  completions   — total number of successful Pod runs required"),
        ('T', "  parallelism   — how many Pods can run simultaneously"),
        ('T', "  restartPolicy — must be Never or OnFailure for Jobs"),
        ('BL', ''),
        ('B', "Common use cases:"),
        ('T', "  Database migrations, batch data processing, report generation, one-time setup tasks"),
    ])
    add_footer(slide, p)

    # Job / CronJob flow visual diagram
    p += 1
    slide_k8s_job_flow_diagram(prs, p)

    # 53. CronJobs theory
    p += 1
    slide = content_slide(prs, "CRONJOBS", "What is a CronJob?")
    add_body_lines(slide, [
        ('T', "A CronJob creates Jobs on a repeating schedule defined in standard cron format."),
        ('BL', ''),
        ('B', "Cron schedule format:"),
        ('C', "# minute  hour  day-of-month  month  day-of-week\n*/1 * * * *    # every minute\n0 2 * * *      # every day at 2 AM\n0 0 * * 0      # every Sunday midnight\n*/5 * * * 1-5  # every 5 min, Mon-Fri"),
        ('BL', ''),
        ('B', "Common use cases:"),
        ('T', "  Scheduled backups, cleanup jobs, report emails, cache invalidation"),
        ('BL', ''),
        ('I', "CronJob -> creates a Job -> Job creates Pods -> Pods run to completion."),
    ])
    add_footer(slide, p)

    # 54. Lab 19 header
    p += 1
    lab_header_slide(prs, "LAB 19", "Jobs & CronJobs",
                     "day2-05-k8s-storage-jobs", p)

    # 55. Lab 19a — Jobs
    p += 1
    slide = content_slide(prs, "LAB 19 — JOBS", "Jobs — Commands")
    add_lab_body_lines(slide, [
        ('B', "Imperative one-shot job:"),
        ('C', "kubectl create job my-job --image=busybox -- echo \"Hello from Job!\"\nkubectl get jobs\nkubectl get pods\nkubectl logs job/my-job\nkubectl delete job my-job"),
        ('B', "Declarative — 3 completions, 2 in parallel (job.yaml):"),
        ('C', "# spec: completions: 3  parallelism: 2\n# containers: command: [\"sh\",\"-c\",\"echo 'Processing...' && sleep 3 && echo 'Done!'\"]\n# restartPolicy: Never\nkubectl apply -f job.yaml\nkubectl get pods -w\nkubectl get jobs\nkubectl logs job/countdown\nkubectl delete -f job.yaml"),
        ('I', "completions: 3 = 3 successful runs total. parallelism: 2 = 2 Pods run at the same time."),
    ])
    kc_note(slide, "day2-05-k8s-storage-jobs")
    add_footer(slide, p)

    # 56. Lab 19b — CronJobs
    p += 1
    slide = content_slide(prs, "LAB 19 — CRONJOBS", "CronJobs — Commands")
    add_lab_body_lines(slide, [
        ('B', "Imperative — every minute:"),
        ('C', "kubectl create cronjob my-cron \\\n  --image=busybox \\\n  --schedule=\"*/1 * * * *\" \\\n  -- echo \"Hello from CronJob!\"\nkubectl get cronjobs\n# Wait 1-2 minutes for a Job to appear\nkubectl get jobs\nkubectl logs job/$(kubectl get jobs --sort-by=.metadata.creationTimestamp \\\n  -o jsonpath='{.items[-1].metadata.name}')\nkubectl delete cronjob my-cron"),
        ('B', "Declarative — every 2 minutes (cronjob.yaml):"),
        ('C', "# schedule: \"*/2 * * * *\"\n# command: [\"sh\",\"-c\",\"echo 'Backup at $(date)' && sleep 3 && echo 'Done!'\"]\nkubectl apply -f cronjob.yaml\nkubectl get cronjobs\n# Wait 2-4 minutes\nkubectl get jobs\nkubectl delete -f cronjob.yaml"),
    ])
    kc_note(slide, "day2-05-k8s-storage-jobs")
    add_footer(slide, p)

    # 57. Knowledge Check 3
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_AMBER)
    add_section_label(slide, "KNOWLEDGE CHECK  3", color=C_DARK)
    add_main_title(slide, "Quick Review — Storage & Jobs", color=C_DARK)
    add_body_lines(slide, [
        ('B', "Q1.  Which volume type is shared between containers in the same Pod but lost when the Pod is deleted?"),
        ('C', "A) PV     B) PVC     C) emptyDir     D) hostPath"),
        ('BL', ''),
        ('B', "Q2.  What K8s resource stores state that persists across Pod deletions?"),
        ('BL', ''),
        ('B', "Q3.  True / False:  A Job restarts Pods after they complete successfully."),
        ('BL', ''),
        ('B', "Q4.  What cron expression runs a job every day at midnight?"),
        ('C', "A) * * * * *     B) 0 0 * * *     C) */1 0 * * *     D) 0 * * * 0"),
        ('BL', ''),
        ('B', "Q5.  What does  parallelism: 2  mean in a Job spec?"),
    ], y=Inches(1.80))
    add_footer(slide, p)

    # ── Bonus: ConfigMaps & Secrets ───────────────────────────────────────────
    p += 1
    slide = content_slide(prs, "CONFIGURATION", "ConfigMaps & Secrets")
    add_body_lines(slide, [
        ('B', "ConfigMap — non-sensitive configuration"),
        ('T', "  Decouple config from image. Mount as env vars or files inside Pods."),
        ('C', "kubectl create configmap app-config --from-literal=APP_ENV=production\nkubectl get configmap app-config -o yaml"),
        ('BL', ''),
        ('B', "Secret — sensitive data (base64-encoded, not encrypted by default)"),
        ('T', "  Passwords, API keys, TLS certs. Referenced by Pods via env vars or mounted files."),
        ('C', "kubectl create secret generic db-secret \\\n  --from-literal=DB_PASSWORD=secret123\nkubectl get secret db-secret -o yaml"),
        ('BL', ''),
        ('I', "For production: use external secret managers (Vault, AWS Secrets Manager) with K8s."),
    ])
    add_footer(slide, p)

    # ── Bonus: Resource Requests & Limits ────────────────────────────────────
    p += 1
    slide = content_slide(prs, "RESOURCE MANAGEMENT", "Resource Requests & Limits")
    add_body_lines(slide, [
        ('T', "Define how much CPU and memory a container needs and can use."),
        ('BL', ''),
        ('B', "In your container spec:"),
        ('C', "resources:\n  requests:\n    memory: \"64Mi\"    # minimum to schedule\n    cpu: \"250m\"       # 250 millicores = 0.25 CPU\n  limits:\n    memory: \"128Mi\"   # container OOM-killed if exceeded\n    cpu: \"500m\"       # throttled if exceeded"),
        ('BL', ''),
        ('B', "Why it matters:"),
        ('T', "  - requests: used by scheduler to find a Node with enough capacity"),
        ('T', "  - limits: prevents a greedy container from starving others"),
        ('BL', ''),
        ('I', "Always set requests and limits in production. No limits = noisy neighbour problem."),
    ])
    add_footer(slide, p)

    # ── Day 2 Summary ─────────────────────────────────────────────────────────
    p += 1
    slide = content_slide(prs, "WRAP UP", "Day 2 — Summary")
    add_body_lines(slide, [
        ('B', "What we covered today:"),
        ('T', "  - Kubernetes architecture: control plane, worker nodes, kubectl"),
        ('T', "  - Pods: smallest unit, lifecycle, imperative & declarative (Labs 13)"),
        ('T', "  - Namespaces: environment isolation  (Lab 14)"),
        ('T', "  - Deployments: self-healing, scaling, ReplicaSets  (Lab 15)"),
        ('T', "  - Rolling updates & rollbacks: zero-downtime deploys  (Lab 16)"),
        ('T', "  - Services: ClusterIP (internal) and NodePort (external)  (Lab 17)"),
        ('T', "  - Storage: emptyDir, PersistentVolume & PVC  (Lab 18)"),
        ('T', "  - Jobs & CronJobs: batch and scheduled tasks  (Lab 19)"),
        ('T', "  - ConfigMaps, Secrets, Resource Requests & Limits"),
    ])
    add_footer(slide, p)

    # ── Assessment ────────────────────────────────────────────────────────────
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_LTBLUE)
    add_section_label(slide, "ASSESSMENT", color=C_DARK)
    add_main_title(slide, "Assessment & Funding", color=C_DARK)
    add_body_lines(slide, [
        ('B', "Final Assessment (now):"),
        ('T', "  Short written / MCQ assessment — must be attempted for SSG funding"),
        ('T', "  Assessment is on the LMS: https://lms.tertiaryinfotech.com"),
        ('BL', ''),
        ('B', "SSG Attendance:"),
        ('T', "  Minimum 75% attendance across all sessions (AM + PM each day)"),
        ('BL', ''),
        ('B', "TRAQOM Survey (mandatory before you leave):"),
        ('T', "  Complete the survey — link on the LMS. Required for course completion."),
        ('BL', ''),
        ('I', "Certificate of Attendance will be issued via email within 2 weeks."),
    ])
    add_footer(slide, p)

    # ── Closing ───────────────────────────────────────────────────────────────
    p += 1
    slide = blank_slide(prs)
    set_bg(slide, C_DARK)
    add_circle(slide, 11.6, 0.5, 2.3, C_ORANGE, alpha=12)
    add_circle(slide, 12.3, 4.6, 1.7, C_ORANGE, alpha=8)
    add_textbox(slide, "END OF COURSE",
                Inches(0.90), Inches(1.80), Inches(11.0), Inches(0.40),
                font="Calibri", size=14, bold=True, color=C_ORANGE)
    add_textbox(slide, "Thank You!",
                Inches(0.85), Inches(2.30), Inches(11.0), Inches(1.10),
                font="Cambria", size=52, bold=True, color=C_WHITE)
    add_textbox(slide, "Keep building — containers are only the beginning.",
                Inches(0.90), Inches(3.60), Inches(11.0), Inches(0.50),
                font="Calibri", size=18, color=C_MUTED2)
    add_textbox(slide, "Mohan Pothula   |   Tertiary Infotech Pte. Ltd.",
                Inches(0.90), Inches(4.20), Inches(11.0), Inches(0.35),
                font="Calibri", size=14, color=C_MUTED)
    add_textbox(slide, "linkedin.com/in/mohanpothula   |   mohanpothula@gmail.com",
                Inches(0.90), Inches(4.58), Inches(11.0), Inches(0.32),
                font="Calibri", size=13, italic=True, color=C_BLUE)
    add_textbox(slide, "enquiry@tertiaryinfotech.com   |   +65 6100 0613   |   www.tertiarycourses.com.sg",
                Inches(0.90), Inches(4.94), Inches(11.0), Inches(0.35),
                font="Calibri", size=13, color=C_MUTED)
    add_footer(slide, p)

    return p


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    # Single combined PPTX for both days
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    build_day1(prs)
    day1_count = len(prs.slides)
    build_day2(prs, start_page=day1_count)
    total = len(prs.slides)
    day2_count = total - day1_count

    out = "Docker_Kubernetes_v3.pptx"
    try:
        prs.save(out)
    except PermissionError:
        out = "Docker_Kubernetes_v3_new.pptx"
        prs.save(out)
        print(f"Note: file was locked -- saved as {out}")

    print(f"Saved  {out}  ({total} slides)")
    print(f"  Day 1 (Docker):     {day1_count} slides")
    print(f"  Day 2 (Kubernetes): {day2_count} slides")
