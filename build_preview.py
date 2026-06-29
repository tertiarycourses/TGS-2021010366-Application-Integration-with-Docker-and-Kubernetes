"""
build_preview.py  —  Apply Agentic AI v12 design to first 30 slides of the
Docker & Kubernetes v2 deck and save as Docker_K8s_Preview_v3.pptx.

Design tokens (Agentic AI v12 style, Windows-safe fonts):
  Titles  → Georgia     (≈ Caladea)
  Body    → Calibri     (≈ Carlito)
  Code    → Courier New (unchanged)
  Accent  → #D97757 warm orange (replaces WSQ dark blue #1F4E79)
  Body fg → #21201C near-black
  Footer  → #6B6862 warm gray, 9pt, centered
"""

import os, shutil
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Paths ─────────────────────────────────────────────────────────────────────
SRC  = r'c:\Users\mohan\agents\Courseware\Docker_Kubernetes_Courseware_v2.pptx'
DST  = r'c:\Users\mohan\agents\Courseware\Docker_K8s_Preview_v3.pptx'
KEEP = 30

COURSE    = "Application Integration with Docker & Kubernetes"
COPYRIGHT = "© Tertiary Infotech Pte. Ltd.  |  TGS-2021010366"

# ── Design tokens ─────────────────────────────────────────────────────────────
TITLE_FONT = 'Georgia'
BODY_FONT  = 'Calibri'
CODE_FONT  = 'Courier New'

COL_ACCENT   = RGBColor(0xD9, 0x77, 0x57)   # #D97757 warm orange
COL_BODY     = RGBColor(0x21, 0x20, 0x1C)   # #21201C near-black body
COL_FOOTER   = RGBColor(0x6B, 0x68, 0x62)   # #6B6862 warm gray
COL_FOOTER_D = RGBColor(0x8A, 0x86, 0x7D)   # lighter gray for dark slides
COL_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
COL_OLD_BAR  = RGBColor(0x1F, 0x4E, 0x79)   # old WSQ blue lab bar → replace

R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

# ── Section labels for slides 0-29 ────────────────────────────────────────────
# Spaced-caps orange label added top-left (light-background slides only)
SECTION_LABELS = {
     0: "W S Q   C O U R S E",
     1: "I N T R O D U C T I O N",
     2: "I N T R O D U C T I O N",
     3: "I N T R O D U C T I O N",
     4: "O V E R V I E W",
     5: "O V E R V I E W",
     6: "O V E R V I E W",
     7: "O V E R V I E W",
     8: "D A Y   1   —   D O C K E R",
     9: "D A Y   1   —   D O C K E R",
    10: "D A Y   1   —   D O C K E R",
    11: "D A Y   1   —   D O C K E R",
    12: "D A Y   1   —   D O C K E R",
    13: "D A Y   1   —   D O C K E R",
    14: "D A Y   1   —   D O C K E R",
    15: "D A Y   1   —   D O C K E R",
    16: "L A B S",
    17: "D O C K E R   F U N D A M E N T A L S",
    18: "D O C K E R   F U N D A M E N T A L S",
    19: "D O C K E R   F U N D A M E N T A L S",
    20: "D O C K E R   F U N D A M E N T A L S",
    21: "D O C K E R   F U N D A M E N T A L S",
    22: "D O C K E R   F U N D A M E N T A L S",
    23: "D O C K E R   F U N D A M E N T A L S",
    24: "D O C K E R   F U N D A M E N T A L S",
    25: "D O C K E R   F U N D A M E N T A L S",
    26: "D O C K E R   F U N D A M E N T A L S",
    27: "D O C K E R   F U N D A M E N T A L S",
    28: "D O C K E R   F U N D A M E N T A L S",
    29: "D O C K E R   F U N D A M E N T A L S",
}

# ── Load ──────────────────────────────────────────────────────────────────────
assert os.path.exists(SRC), f"v2 not found: {SRC}\nRun build_courseware.py first."
shutil.copy2(SRC, DST)
prs = Presentation(DST)
print(f"Loaded {len(prs.slides)} slides from v2")

# ── 1. Trim to KEEP slides ────────────────────────────────────────────────────
sldIdLst = prs.slides._sldIdLst
while len(prs.slides) > KEEP:
    last = sldIdLst[-1]
    rId  = last.get(f'{{{R_NS}}}id')
    prs.part.drop_rel(rId)
    sldIdLst.remove(last)
print(f"Trimmed to {len(prs.slides)} slides")

# ── 2. Helpers ────────────────────────────────────────────────────────────────

def _no_bullet(para):
    pPr = para._p.get_or_add_pPr()
    for tag in ('a:buChar','a:buAutoNum','a:buFont','a:buClr','a:buNone','a:buBlip'):
        for el in pPr.findall(tag, {'a': A_NS}):
            pPr.remove(el)
    pPr.append(etree.fromstring(f'<a:buNone xmlns:a="{A_NS}"/>'))


def is_dark_slide(slide):
    """Heuristic: white title text → dark background slide."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        if run.font.color.rgb == COL_WHITE:
                            return True
                    except Exception:
                        pass
    # Also check solid background fill luminance
    try:
        rgb = slide.background.fill.fore_color.rgb
        lum = 0.299*rgb[0] + 0.587*rgb[1] + 0.114*rgb[2]
        if lum < 80:
            return True
    except Exception:
        pass
    return False


def recolor_blue_shapes(slide):
    """Replace old WSQ blue (#1F4E79) shape fills with orange accent."""
    for shape in slide.shapes:
        try:
            if shape.fill.fore_color.rgb == COL_OLD_BAR:
                shape.fill.fore_color.rgb = COL_ACCENT
                shape.line.color.rgb      = COL_ACCENT
        except Exception:
            pass


def restyle_fonts(slide):
    """
    For each text run on the slide:
      - Title placeholder (idx=0) or large text (>22pt) → Georgia
      - Everything else except Courier New              → Calibri
      - Body text that is plain black                   → near-black #21201C
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            idx = shape.placeholder_format.idx if shape.placeholder_format else -1
        except Exception:
            idx = -1

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                # Never touch code font
                if run.font.name == CODE_FONT:
                    continue

                # Font family
                pt = run.font.size.pt if run.font.size else 0
                if idx == 0 or pt > 22:
                    run.font.name = TITLE_FONT
                else:
                    run.font.name = BODY_FONT

                # Soften pure black → near-black on non-dark slides
                try:
                    if run.font.color.type is not None:
                        rgb = run.font.color.rgb
                        if rgb in (RGBColor(0,0,0), RGBColor(0x26,0x26,0x26)):
                            run.font.color.rgb = COL_BODY
                except Exception:
                    pass


def add_section_label(slide, label):
    """Spaced-caps orange label, 12pt Calibri Bold, top-left."""
    txb = slide.shapes.add_textbox(
        left=Inches(0.3), top=Inches(0.06),
        width=Inches(7), height=Inches(0.25)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    _no_bullet(p)
    r = p.add_run()
    r.text           = label
    r.font.size      = Pt(12)
    r.font.name      = BODY_FONT
    r.font.bold      = True
    r.font.color.rgb = COL_ACCENT


def add_footer(slide, page_num, dark_bg):
    """Three-part footer: copyright | course | page, 9pt, centered."""
    sw = prs.slide_width
    sh = prs.slide_height
    txb = slide.shapes.add_textbox(
        left=Inches(0.3),
        top=sh - Inches(0.32),
        width=sw - Inches(0.6),
        height=Inches(0.25)
    )
    tf = txb.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    _no_bullet(p)
    r = p.add_run()
    r.text           = f"{COPYRIGHT}    {COURSE}    {page_num}"
    r.font.size      = Pt(9)
    r.font.name      = BODY_FONT
    r.font.color.rgb = COL_FOOTER_D if dark_bg else COL_FOOTER
    p.alignment      = PP_ALIGN.CENTER


# ── 3. Apply design to all kept slides ────────────────────────────────────────
slides = list(prs.slides)
for i, slide in enumerate(slides):
    dark = is_dark_slide(slide)
    recolor_blue_shapes(slide)
    restyle_fonts(slide)
    add_footer(slide, i + 1, dark)
    label = SECTION_LABELS.get(i, "")
    if label and not dark:
        add_section_label(slide, label)
    print(f"  Slide {i+1:2d}  {'[dark]' if dark else '[light]'}  {label or ''}")

# ── 4. Save ───────────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nSaved: {DST}")
