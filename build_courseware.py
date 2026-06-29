import shutil
from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

SRC = r'c:\Users\mohan\agents\Courseware\WSQ - Dr. Alfred Ang  - Application Integration with Docker and Kubernetes - v20.pptx'
DST = r'c:\Users\mohan\agents\Courseware\Docker_Kubernetes_Courseware_v2.pptx'
KC  = "https://killercoda.com/tertiary-labs/course/killercoda"

# ── Font standards (match v20 source — Arial throughout) ─────────────────────
BODY_FONT = 'Arial'
CODE_FONT = 'Courier New'

# ── Colours ───────────────────────────────────────────────────────────────────
COL_CODE    = RGBColor(0x1F, 0x49, 0x7D)   # dark navy-blue for code
COL_WARN    = RGBColor(0xC0, 0x00, 0x00)   # red for warnings / ref-only
COL_LAB_BAR = RGBColor(0x1F, 0x4E, 0x79)  # WSQ dark-blue lab header bar
COL_BAR_TXT = RGBColor(0xFF, 0xFF, 0xFF)   # white text on bar
COL_KC_LINE = RGBColor(0x2E, 0x74, 0xB5)  # blue for KillerCoda URL line

A_NS = 'http://schemas.openxmlformats.org/drawingml/2006/main'

shutil.copy2(SRC, DST)
prs = Presentation(DST)
S   = list(prs.slides)
print(f"Loaded {len(S)} slides")

# ── helpers ───────────────────────────────────────────────────────────────────

def _body(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            return ph
    for ph in slide.placeholders:
        if ph.placeholder_format.idx not in (0, 12, 13, 14, 15):
            return ph
    return None

def _no_bullet(para):
    """Remove any bullet from a paragraph so text is flush-left."""
    pPr = para._p.get_or_add_pPr()
    # Remove existing bullet-related children
    for tag in ('a:buChar','a:buAutoNum','a:buFont','a:buClr',
                'a:buSzPct','a:buSzPts','a:buNone','a:buBlip'):
        for el in pPr.findall(tag, {'a': A_NS}):
            pPr.remove(el)
    # Insert buNone to explicitly suppress bullets
    pPr.append(etree.fromstring(f'<a:buNone xmlns:a="{A_NS}"/>'))

def _indent(para, level=1):
    """Add left indent (in EMU) — used to visually indent code blocks."""
    pPr = para._p.get_or_add_pPr()
    indent = level * 228600   # 0.25 inch per level in EMU
    pPr.set('marL', str(indent))
    pPr.set('indent', '0')

def write(slide, rows):
    ph = _body(slide)
    if not ph:
        return
    tf = ph.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for row in rows:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _no_bullet(p)                        # suppress all bullets

        txt  = row.get('t', '')
        code = row.get('c', False)
        bold = row.get('b', False)
        ital = row.get('i', False)
        warn = row.get('warn', False)

        if not txt:
            continue

        if code:
            _indent(p, level=1)              # indent code blocks

        r = p.add_run()
        r.text = txt
        r.font.size   = Pt(10 if code else 13)
        r.font.name   = CODE_FONT if code else BODY_FONT
        r.font.bold   = bold
        r.font.italic = ital

        if code:
            r.font.color.rgb = COL_CODE
        elif warn:
            r.font.color.rgb = COL_WARN
        elif ital and not bold:
            r.font.color.rgb = COL_KC_LINE   # colour KillerCoda / key-concept lines

def add_lab_bar(slide, label):
    """Add a WSQ-style coloured header bar + LAB label to a lab slide."""
    sw = prs.slide_width
    bar = slide.shapes.add_shape(
        1,                                   # MSO_SHAPE_TYPE.RECTANGLE = 1
        left=0, top=0, width=sw, height=Inches(0.42)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COL_LAB_BAR
    bar.line.color.rgb      = COL_LAB_BAR

    tf = bar.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    _no_bullet(p)
    r  = p.add_run()
    r.text         = label
    r.font.size    = Pt(13)
    r.font.name    = BODY_FONT
    r.font.bold    = True
    r.font.color.rgb = COL_BAR_TXT
    p.alignment    = PP_ALIGN.LEFT

def mark_ref_only(slide):
    """Add a red Reference-Only banner as a text box."""
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(0.72), Inches(9.4), Inches(0.32))
    tf = txBox.text_frame
    p  = tf.paragraphs[0]
    _no_bullet(p)
    r  = p.add_run()
    r.text           = "[ Reference Only — Students practice this independently, not in KillerCoda ]"
    r.font.bold      = True
    r.font.size      = Pt(11)
    r.font.italic    = True
    r.font.name      = BODY_FONT
    r.font.color.rgb = COL_WARN

def T(text, **kw): return {'t': text, **kw}
def C(text):       return {'t': text, 'c': True}
def B(text):       return {'t': text, 'b': True}
def I(text):       return {'t': text, 'i': True}
def BL():          return {'t': ''}
def KC_NOTE(scenario): return I(f"KillerCoda: {scenario}  |  {KC}")

# ── Lab bar labels ────────────────────────────────────────────────────────────
LAB_BARS = {
    16:  "Download Labs",
    47:  "LAB 1  —  Run Your First Container",
    49:  "LAB 2  —  Run Nginx and Copy Files",
    59:  "LAB 3  —  Build a Python Image",
    64:  "LAB 4  —  Flask App in Docker",
    72:  "LAB 5  —  Docker Volumes",
    74:  "LAB 5  —  Docker Bind Mounts",
    79:  "LAB 6  —  Custom Networks",
    82:  "LAB 7  —  Port Mapping",
    85:  "LAB 8  —  Environment Variables",
    110: "LAB 10  —  Docker Compose (Single Service)",
    114: "LAB 11  —  Docker Compose: Flask + Redis",
    117: "LAB 12  —  Docker Compose: Full Stack",
    165: "LAB 13  —  Kubernetes Pods (Imperative)",
    167: "LAB 13  —  Kubernetes Pods (Declarative)",
    177: "LAB 14  —  Namespaces",
    188: "LAB 15  —  Deployments (Imperative)",
    190: "LAB 15  —  Deployments (Declarative YAML)",
    197: "LAB 16  —  Rolling Updates",
    198: "LAB 16  —  Rollbacks",
    225: "LAB 17  —  ClusterIP Service",
    226: "LAB 17  —  NodePort Service",
    227: "LAB 17  —  Service YAML",
    250: "LAB 18  —  emptyDir Volumes",
    251: "LAB 18  —  PersistentVolume & PVC",
    271: "LAB 19  —  Jobs",
    272: "LAB 19  —  CronJobs",
}

# ── Slide 16 (0-idx): Download Labs ──────────────────────────────────────────
print("Updating Download Labs slide...")
write(S[16], [
    B("Lab Environment: KillerCoda (browser-based, no installation needed)"),
    BL(),
    T("Access your labs at:"),
    C(KC),
    BL(),
    T("Click the link  →  Select your scenario  →  Click Start"),
    T("A live terminal opens instantly in your browser. Everything is pre-installed."),
    BL(),
    B("GitHub Repository (reference & source materials):"),
    C("https://github.com/tertiarycourses/TGS-2021010366-Application-Integration-with-Docker-and-Kubernetes"),
    BL(),
    I("Note: Each KillerCoda session starts fresh — re-run setup commands each time."),
    I("On KillerCoda always use  docker-compose  (with hyphen), not  docker compose"),
])

# ── LAB 1 ────────────────────────────────────────────────────────────────────
print("Updating Lab 1...")
write(S[47], [
    KC_NOTE("day1-01-docker-fundamentals  →  Lab 1"),
    BL(),
    B("Run an interactive Ubuntu container:"),
    C("docker run -it ubuntu:latest"),
    BL(),
    B("Inside the container:"),
    C("cat /etc/hosts"),
    C("exit"),
    BL(),
    B("Back on the host — check containers:"),
    C("docker ps              # running containers only"),
    C("docker ps -a           # ALL containers including stopped"),
    BL(),
    B("Verify:"),
    C("docker ps -a | grep ubuntu"),
    BL(),
    I("Key: -it = interactive terminal. Container stops when you exit the shell."),
])

write(S[48], [
    B("Why does ubuntu NOT appear in  docker ps ?"),
    BL(),
    T("docker ps  shows only RUNNING containers."),
    T("Ubuntu's PID 1 is the shell — once you type exit, the shell ends → container stops."),
    BL(),
    B("Containers are NOT virtual machines:"),
    T("  • A VM keeps running after you log out"),
    T("  • A container stops when its main process (PID 1) exits"),
    BL(),
    B("To keep a container running in the background:"),
    C("docker run -d --name my-nginx nginx   # -d = detached / background"),
    C("docker ps                              # nginx shows as running"),
    BL(),
    I("Rule: A container lives as long as PID 1 is alive."),
])

# ── LAB 2 ────────────────────────────────────────────────────────────────────
print("Updating Lab 2...")
write(S[49], [
    KC_NOTE("day1-01-docker-fundamentals  →  Lab 2"),
    BL(),
    B("Run Nginx in the background:"),
    C("docker run -d --name my-nginx nginx:latest"),
    C("docker ps"),
    BL(),
    B("Copy a file FROM container TO host:"),
    C("docker cp my-nginx:/usr/share/nginx/html/index.html ./index.html"),
    C("cat index.html"),
    BL(),
    B("Inspect logs and clean up:"),
    C("docker logs my-nginx"),
    C("docker stop my-nginx"),
    C("docker rm my-nginx"),
    BL(),
    I("Key: docker cp works both directions. Container need not be running."),
])

# ── LAB 3: Build Python Container ────────────────────────────────────────────
print("Updating Lab 3...")
write(S[59], [
    KC_NOTE("day1-01-docker-fundamentals  →  Lab 3"),
    BL(),
    B("Create project:"),
    C("mkdir lab3 && cd lab3"),
    BL(),
    B("main.py:"),
    C('print("Hello from inside Docker!")'),
    BL(),
    B("Dockerfile:"),
    C("FROM python:3.11-slim"),
    C("WORKDIR /app"),
    C("COPY main.py ."),
    C('CMD ["python", "main.py"]'),
])

write(S[60], [
    B("Build and run:"),
    C("docker build -t lab3 ."),
    C("docker run lab3"),
    BL(),
    T('Expected output:  "Hello from inside Docker!"'),
    BL(),
    B("Inspect:"),
    C("docker images | grep lab3"),
    C("docker ps -a | grep lab3"),
    C("cd .."),
    BL(),
    I("Key: Each Dockerfile line = one cached layer. Only changed layers rebuild — fast incremental builds."),
])

# Lab 3-2 (slides 61, 62) — Tetris Game — not in KillerCoda, keep unchanged

# ── LAB 4: Flask App in Docker ───────────────────────────────────────────────
print("Updating Lab 4...")
write(S[64], [
    KC_NOTE("day1-01-docker-fundamentals  →  Lab 4"),
    BL(),
    B("Create project:"),
    C("mkdir lab4 && cd lab4"),
    BL(),
    B("app.py:"),
    C("from flask import Flask"),
    C("app = Flask(__name__)"),
    C("@app.route('/')"),
    C("def home():"),
    C("    return '<h1>Hello from Flask in Docker!</h1>'"),
    C("if __name__ == '__main__':"),
    C("    app.run(host='0.0.0.0', port=5000)"),
    BL(),
    B("requirements.txt:"),
    C("flask"),
])

write(S[65], [
    B("Dockerfile:"),
    C("FROM python:3.11-slim"),
    C("WORKDIR /app"),
    C("COPY requirements.txt ."),
    C("RUN pip install -r requirements.txt"),
    C("COPY . ."),
    C("EXPOSE 5000"),
    C('CMD ["python", "app.py"]'),
])

write(S[66], [
    B("Build and run:"),
    C("docker build -t lab4 ."),
    C("docker run -d -p 5001:5000 --name flask-app lab4"),
    C("curl http://localhost:5001"),
    BL(),
    B("Clean up:"),
    C("docker stop flask-app && docker rm flask-app"),
    C("cd .."),
    BL(),
    T("Port format:  HOST_PORT : CONTAINER_PORT"),
    T("  EXPOSE 5000   →  documents the port (does not publish it)"),
    T("  -p 5001:5000  →  maps host:5001  to  container:5000"),
    BL(),
    I("Key: Always bind to  0.0.0.0  inside the container so external traffic is accepted."),
])

# ── LAB 5: Volumes ───────────────────────────────────────────────────────────
print("Updating Lab 5...")
write(S[72], [
    KC_NOTE("day1-02-docker-storage  →  Lab 5 (Named Volumes)"),
    BL(),
    B("Create and inspect a named volume:"),
    C("docker volume create my-data"),
    C("docker volume ls"),
    C("docker volume inspect my-data"),
    BL(),
    B("Write data using one container:"),
    C("docker run -d --name app1 -v my-data:/data alpine sleep 3600"),
    C("docker exec app1 sh -c \"echo 'hello from app1' > /data/test.txt\""),
    C("docker exec app1 cat /data/test.txt"),
])

write(S[73], [
    B("Remove the container — data stays in the volume:"),
    C("docker stop app1 && docker rm app1"),
    BL(),
    B("New container reads the same volume — data still there:"),
    C("docker run -d --name app2 -v my-data:/data alpine sleep 3600"),
    C("docker exec app2 cat /data/test.txt   # prints: hello from app1"),
    BL(),
    B("Clean up:"),
    C("docker stop app2 && docker rm app2"),
    C("docker volume rm my-data"),
    BL(),
    I("Key: Named volumes outlive containers. Use for databases and persistent app data."),
])

write(S[74], [
    KC_NOTE("day1-02-docker-storage  →  Lab 5 (Bind Mounts)"),
    BL(),
    B("Map a host folder directly into a container:"),
    C("mkdir -p myfiles"),
    C("echo \"hello from host\" > myfiles/note.txt"),
    BL(),
    B("Read host file inside container:"),
    C("docker run --rm -v $(pwd)/myfiles:/data alpine cat /data/note.txt"),
    BL(),
    B("Write from container — file appears on host immediately:"),
    C("docker run --rm -v $(pwd)/myfiles:/data alpine \\"),
    C("  sh -c \"echo 'hello from container' > /data/reply.txt\""),
    C("cat myfiles/reply.txt"),
    C("rm -rf myfiles"),
])

write(S[75], [
    B("Named Volume  vs  Bind Mount — Comparison"),
    BL(),
    T("Named Volume:"),
    T("  • Managed by Docker  (/var/lib/docker/volumes/)"),
    T("  • Persists beyond container lifecycle"),
    T("  • Best for: databases, persistent app data"),
    C("  -v my-vol:/data"),
    BL(),
    T("Bind Mount:"),
    T("  • Maps an exact host path into container"),
    T("  • Changes are instant in both directions"),
    T("  • Best for: development, live code sync"),
    C("  -v $(pwd)/src:/app"),
    BL(),
    I("Use named volumes in production. Use bind mounts in development."),
])

# Lab 5-2 (slides 76, 77) — Volume in Dockerfile — keep unchanged

# ── LAB 6: Networks ──────────────────────────────────────────────────────────
print("Updating Lab 6...")
write(S[79], [
    KC_NOTE("day1-03-docker-networking  →  Lab 6"),
    BL(),
    B("Create a custom bridge network:"),
    C("docker network ls"),
    C("docker network create my-net"),
    C("docker network inspect my-net"),
    BL(),
    B("Run two containers on the same network:"),
    C("docker run -d --name app1 --network my-net busybox sleep 3600"),
    C("docker run -d --name app2 --network my-net busybox sleep 3600"),
    BL(),
    B("Containers talk by NAME — no IP address needed:"),
    C("docker exec app1 ping -c 3 app2"),
])

write(S[80], [
    B("Disconnect and reconnect:"),
    C("docker network disconnect my-net app2"),
    C("docker exec app1 ping -c 2 app2   # FAILS — not on same network"),
    BL(),
    C("docker network connect my-net app2"),
    C("docker exec app1 ping -c 2 app2   # Works again"),
    BL(),
    B("Clean up:"),
    C("docker stop app1 app2 && docker rm app1 app2"),
    C("docker network rm my-net"),
    BL(),
    I("Key: Custom bridge networks give automatic DNS by container name."),
    I("The default bridge network does NOT support DNS by name."),
])

# ── LAB 7: Port Mapping ──────────────────────────────────────────────────────
print("Updating Lab 7...")
write(S[82], [
    KC_NOTE("day1-03-docker-networking  →  Lab 7"),
    BL(),
    B("Reuse the Flask image from Lab 4:"),
    C("cd lab4 && docker build -t lab4 . && cd .."),
    BL(),
    B("Map host 5001  →  container 5000:"),
    C("docker run -d -p 5001:5000 --name web1 lab4"),
    C("curl http://localhost:5001"),
    BL(),
    B("Run second instance on a different port:"),
    C("docker run -d -p 5002:5000 --name web2 lab4"),
    C("curl http://localhost:5002"),
])

write(S[83], [
    B("Inspect port mappings:"),
    C("docker port web1"),
    C("docker port web2"),
    BL(),
    B("Let Docker pick a random host port (-P):"),
    C("docker run -d -P --name web3 lab4"),
    C("docker port web3   # shows the randomly assigned port"),
    BL(),
    B("Clean up:"),
    C("docker stop web1 web2 web3 && docker rm web1 web2 web3"),
    BL(),
    T("  -p 5001:5000   →  explicit: host:5001 routes to container:5000"),
    T("  -P             →  Docker assigns a random available host port"),
])

# ── LAB 8: Environment Variables ─────────────────────────────────────────────
print("Updating Lab 8...")
write(S[85], [
    KC_NOTE("day1-04-docker-config  →  Lab 8"),
    BL(),
    B("Create project:"),
    C("mkdir lab8 && cd lab8"),
    BL(),
    B("main.py — reads env vars with defaults:"),
    C("import os"),
    C("name = os.environ.get('MY_NAME', 'World')"),
    C("env  = os.environ.get('MY_ENV',  'development')"),
    C("print(f'Hello, {name}! Environment: {env}')"),
    BL(),
    B("Dockerfile with ENV defaults:"),
    C("FROM python:3.11-slim"),
    C("WORKDIR /app"),
    C("ENV MY_NAME=World"),
    C("ENV MY_ENV=development"),
    C("COPY main.py ."),
    C('CMD ["python", "main.py"]'),
])

write(S[86], [
    B("Build and test overrides:"),
    C("docker build -t lab8 ."),
    C("docker run lab8                              # uses Dockerfile defaults"),
    C("docker run -e MY_NAME=Alfred lab8            # override one variable"),
    C("docker run -e MY_NAME=Alfred -e MY_ENV=production lab8"),
    BL(),
    B("Use an env file:"),
    C("# .env"),
    C("MY_NAME=Student"),
    C("MY_ENV=staging"),
    C("docker run --env-file .env lab8"),
    BL(),
    B("List all env vars inside the container:"),
    C("docker run lab8 env"),
    C("cd .."),
    BL(),
    I("Key: ENV in Dockerfile sets defaults. -e or --env-file overrides at runtime."),
    I("Never hardcode secrets in images — use env vars or secret managers."),
])

# ── LAB 9: Push to Docker Hub — REFERENCE ONLY ───────────────────────────────
print("Marking Lab 9 as Reference Only...")
for idx in [88, 89, 90, 91, 92]:
    mark_ref_only(S[idx])

# ── LAB 10: Docker Compose ───────────────────────────────────────────────────
print("Updating Lab 10...")
write(S[110], [
    KC_NOTE("day1-05-docker-compose  →  Lab 10"),
    BL(),
    B("Create project:"),
    C("mkdir lab10 && cd lab10"),
    BL(),
    B("app.py:"),
    C("from flask import Flask"),
    C("app = Flask(__name__)"),
    C("@app.route('/')"),
    C("def home():"),
    C("    return '<h1>Hello from Docker Compose!</h1>'"),
    C("if __name__ == '__main__':"),
    C("    app.run(host='0.0.0.0', port=5000)"),
    BL(),
    B("requirements.txt:"),
    C("flask"),
])

write(S[111], [
    B("Dockerfile:"),
    C("FROM python:3.11-slim"),
    C("WORKDIR /app"),
    C("COPY requirements.txt ."),
    C("RUN pip install -r requirements.txt"),
    C("COPY . ."),
    C('CMD ["python", "app.py"]'),
])

write(S[112], [
    B("docker-compose.yml:"),
    C("services:"),
    C("  web:"),
    C("    build: ."),
    C("    ports:"),
    C('      - "5001:5000"'),
    BL(),
    B("Start and test:"),
    C("docker-compose up -d"),
    C("docker-compose ps"),
    C("curl http://localhost:5001"),
])

write(S[113], [
    B("Useful Compose commands:"),
    C("docker-compose logs                         # view logs"),
    C("docker-compose exec web python -c \"print('Inside container')\""),
    C("docker-compose stop                         # stop (keeps containers)"),
    C("docker-compose start                        # restart"),
    C("docker-compose down                         # stop + remove containers"),
    BL(),
    I("Note: On KillerCoda use  docker-compose  (with hyphen), not  docker compose"),
    BL(),
    I("Key: Compose manages multi-container apps from a single YAML file."),
])

# ── LAB 11: Flask + Redis ─────────────────────────────────────────────────────
print("Updating Lab 11...")
write(S[114], [
    KC_NOTE("day1-05-docker-compose  →  Lab 11"),
    BL(),
    B("Visit counter app — Flask + Redis:"),
    C("mkdir lab11 && cd lab11"),
    BL(),
    B("app.py:"),
    C("from flask import Flask"),
    C("import redis"),
    C("app = Flask(__name__)"),
    C("r = redis.Redis(host='redis', port=6379)"),
    C("@app.route('/')"),
    C("def home():"),
    C("    count = r.incr('visits')"),
    C("    return f'<h1>Visit count: {count}</h1>'"),
    C("if __name__ == '__main__': app.run(host='0.0.0.0', port=5000)"),
    BL(),
    B("requirements.txt:"),
    C("flask"),
    C("redis"),
])

write(S[115], [
    B("docker-compose.yml:"),
    C("services:"),
    C("  web:"),
    C("    build: ."),
    C("    ports:"),
    C('      - "5001:5000"'),
    C("    depends_on:"),
    C("      - redis"),
    C("  redis:"),
    C("    image: redis:7-alpine"),
])

write(S[116], [
    B("Run and test:"),
    C("docker-compose up -d"),
    C("docker-compose ps"),
    C("curl http://localhost:5001   # Visit count: 1"),
    C("curl http://localhost:5001   # Visit count: 2"),
    C("curl http://localhost:5001   # Visit count: 3"),
    BL(),
    B("Clean up:"),
    C("docker-compose down"),
    C("cd .."),
    BL(),
    I("Key: 'redis' in  host='redis'  resolves by service name via Compose networking."),
    I("depends_on ensures Redis starts before the web service."),
])

# ── LAB 12: Full Stack ───────────────────────────────────────────────────────
print("Updating Lab 12...")
write(S[117], [
    KC_NOTE("day1-05-docker-compose  →  Lab 12"),
    BL(),
    B("3-tier stack: Node.js  +  PostgreSQL  +  Redis"),
    C("mkdir lab12 && cd lab12"),
    BL(),
    B("docker-compose.yml (key sections):"),
    C("services:"),
    C("  web:"),
    C("    build: ."),
    C('    ports: ["3000:3000"]'),
    C("    environment:"),
    C("      DATABASE_URL: postgresql://user:pass@db:5432/mydb"),
    C("      REDIS_URL: redis://redis:6379"),
    C("    depends_on:"),
    C("      db: {condition: service_healthy}"),
    C("  db:"),
    C("    image: postgres:16-alpine"),
    C("    healthcheck:"),
    C('      test: ["CMD-SHELL","pg_isready -U user -d mydb"]'),
    C("      interval: 5s  retries: 5"),
    C("  redis:"),
    C("    image: redis:7-alpine"),
])

write(S[118], [
    B("Start and verify all endpoints:"),
    C("docker-compose up -d"),
    C("docker-compose ps"),
    C("curl http://localhost:3000           # visit counter"),
    C("curl http://localhost:3000/db        # current time from PostgreSQL"),
    C("curl http://localhost:3000/health    # {\"status\":\"ok\"}"),
    BL(),
    B("Clean up:"),
    C("docker-compose down -v   # -v removes named volumes too"),
    C("cd .."),
    BL(),
    I("Key: service_healthy waits for the DB healthcheck before starting web."),
    I("Named volumes (pgdata, redisdata) keep data across container restarts."),
])

# ── LAB 13: Kubernetes Pods ──────────────────────────────────────────────────
print("Updating Lab 13...")
write(S[165], [
    KC_NOTE("day2-01-k8s-pods-namespaces  →  Lab 13 (Imperative)"),
    BL(),
    B("Create and inspect a Pod:"),
    C("kubectl run my-nginx --image=nginx"),
    C("kubectl get pods"),
    C("kubectl get pods -o wide"),
    C("kubectl describe pod my-nginx"),
    BL(),
    B("Interact with the Pod:"),
    C("kubectl logs my-nginx"),
    C("kubectl exec my-nginx -- cat /etc/hostname"),
    C("kubectl exec -it my-nginx -- /bin/sh"),
    C("# Inside: ls /usr/share/nginx/html  then  exit"),
])

write(S[166], [
    B("Delete the Pod:"),
    C("kubectl delete pod my-nginx"),
    C("kubectl get pods"),
    BL(),
    I("Key: Pods are mortal — delete one and it is gone forever."),
    I("Deployments (Lab 15) automatically recreate deleted Pods."),
])

write(S[167], [
    KC_NOTE("day2-01-k8s-pods-namespaces  →  Lab 13 (Declarative YAML)"),
    BL(),
    B("pod.yaml:"),
    C("apiVersion: v1"),
    C("kind: Pod"),
    C("metadata:"),
    C("  name: my-app"),
    C("  labels:"),
    C("    app: my-app"),
    C("spec:"),
    C("  containers:"),
    C("    - name: nginx"),
    C("      image: nginx"),
    C("      ports:"),
    C("        - containerPort: 80"),
])

write(S[168], [
    B("Apply, inspect, and delete:"),
    C("kubectl apply -f pod.yaml"),
    C("kubectl get pods"),
    C("kubectl describe pod my-app"),
    BL(),
    C("kubectl delete -f pod.yaml"),
    C("kubectl get pods"),
    BL(),
    I("Key: kubectl apply -f is idempotent — safe to run multiple times."),
    I("Always prefer YAML manifests for production workloads."),
])

# ── LAB 14: Namespaces ───────────────────────────────────────────────────────
print("Updating Lab 14...")
write(S[177], [
    KC_NOTE("day2-01-k8s-pods-namespaces  →  Lab 14"),
    BL(),
    B("Imperative — create namespace and run Pod inside it:"),
    C("kubectl get namespaces"),
    C("kubectl create namespace dev"),
    C("kubectl run my-nginx --image=nginx -n dev"),
    C("kubectl get pods -n dev"),
    C("kubectl get pods --all-namespaces"),
    C("kubectl describe pod my-nginx -n dev"),
    C("kubectl delete namespace dev   # removes everything inside"),
])

write(S[178], [
    B("Declarative namespace:"),
    C("# namespace.yaml"),
    C("apiVersion: v1"),
    C("kind: Namespace"),
    C("metadata:"),
    C("  name: staging"),
    BL(),
    C("kubectl apply -f namespace.yaml"),
])

write(S[179], [
    B("Pod in a specific namespace (YAML):"),
    C("apiVersion: v1"),
    C("kind: Pod"),
    C("metadata:"),
    C("  name: my-app"),
    C("  namespace: staging"),
    C("spec:"),
    C("  containers:"),
    C("    - name: nginx"),
    C("      image: nginx"),
    BL(),
    C("kubectl apply -f pod-staging.yaml"),
    C("kubectl get pods -n staging"),
    C("kubectl delete -f namespace.yaml   # removes namespace + all Pods"),
    BL(),
    I("Deleting a namespace removes every resource inside it instantly."),
])

# ── LAB 15: Deployments ──────────────────────────────────────────────────────
print("Updating Lab 15...")
write(S[188], [
    KC_NOTE("day2-02-k8s-deployments  →  Lab 15 (Imperative)"),
    BL(),
    B("Create and inspect a Deployment:"),
    C("kubectl create deployment my-nginx --image=nginx"),
    C("kubectl get deployments"),
    C("kubectl get pods"),
    C("kubectl describe deployment my-nginx"),
    BL(),
    B("Scale to 3 replicas:"),
    C("kubectl scale deployment my-nginx --replicas=3"),
    C("kubectl get pods -w              # watch live  (Ctrl+C to stop)"),
])

write(S[189], [
    B("Self-healing — delete a Pod and watch it recreate:"),
    C("kubectl delete pod $(kubectl get pods -l app=my-nginx -o name | head -1 | cut -d/ -f2)"),
    C("kubectl get pods                 # replacement Pod appears in seconds"),
    BL(),
    B("Scale down and delete:"),
    C("kubectl scale deployment my-nginx --replicas=1"),
    C("kubectl delete deployment my-nginx"),
    BL(),
    I("Key: Deployments use a ReplicaSet to enforce the desired Pod count."),
    I("Delete a Pod → ReplicaSet detects the deficit → creates a new one automatically."),
])

write(S[190], [
    KC_NOTE("day2-02-k8s-deployments  →  Lab 15 (Declarative YAML)"),
    BL(),
    B("deployment.yaml:"),
    C("apiVersion: apps/v1"),
    C("kind: Deployment"),
    C("metadata:"),
    C("  name: my-app"),
    C("spec:"),
    C("  replicas: 3"),
    C("  selector:"),
    C("    matchLabels:"),
    C("      app: my-app"),
    C("  template:"),
    C("    metadata:"),
    C("      labels:"),
    C("        app: my-app"),
    C("    spec:"),
    C("      containers:"),
    C("        - name: nginx"),
    C("          image: nginx"),
    C("          ports:"),
    C("            - containerPort: 80"),
    BL(),
    C("kubectl apply -f deployment.yaml"),
    C("kubectl get deployments && kubectl get pods"),
    C("kubectl delete -f deployment.yaml"),
])

# ── LAB 16: Rollout & Rollback ───────────────────────────────────────────────
print("Updating Lab 16...")
write(S[197], [
    KC_NOTE("day2-03-k8s-rollouts  →  Lab 16 (Rolling Update)"),
    BL(),
    B("Deploy with a specific image version:"),
    C("kubectl create deployment my-nginx --image=nginx:1.24 --replicas=3"),
    C("kubectl get deployments"),
    C("kubectl describe deployment my-nginx | grep Image"),
    BL(),
    B("Update image — rolling update with zero downtime:"),
    C("kubectl set image deployment my-nginx nginx=nginx:1.25"),
    C("kubectl rollout status deployment my-nginx"),
    C("kubectl describe deployment my-nginx | grep Image"),
    BL(),
    B("Update again and view history:"),
    C("kubectl set image deployment my-nginx nginx=nginx:1.27"),
    C("kubectl rollout status deployment my-nginx"),
    C("kubectl rollout history deployment my-nginx"),
])

write(S[198], [
    KC_NOTE("day2-03-k8s-rollouts  →  Lab 16 (Rollback)"),
    BL(),
    B("Rollback to previous version:"),
    C("kubectl rollout undo deployment my-nginx"),
    C("kubectl describe deployment my-nginx | grep Image   # back to 1.25"),
    BL(),
    B("Rollback to a specific revision:"),
    C("kubectl rollout history deployment my-nginx"),
    C("kubectl rollout undo deployment my-nginx --to-revision=1"),
    C("kubectl describe deployment my-nginx | grep Image   # back to 1.24"),
])

write(S[199], [
    B("Declarative rollout + rollback:"),
    C("# deployment.yaml  with  image: nginx:1.24"),
    C("kubectl apply -f deployment.yaml"),
    BL(),
    B("Update by editing the YAML:"),
    C("sed -i 's/nginx:1.24/nginx:1.25/' deployment.yaml"),
    C("kubectl apply -f deployment.yaml"),
    C("kubectl rollout status deployment my-app"),
    C("kubectl rollout history deployment my-app"),
    BL(),
    B("Rollback:"),
    C("kubectl rollout undo deployment my-app"),
    BL(),
    B("Clean up:"),
    C("kubectl delete deployment my-nginx"),
    C("kubectl delete -f deployment.yaml"),
    BL(),
    I("Key: Rolling update = zero downtime. New Pods come up before old ones go down."),
])

# ── LAB 17: Services ─────────────────────────────────────────────────────────
print("Updating Lab 17...")
write(S[225], [
    KC_NOTE("day2-04-k8s-services  →  Lab 17 (ClusterIP)"),
    BL(),
    B("Create Deployment + ClusterIP Service:"),
    C("kubectl create deployment my-nginx --image=nginx --replicas=2"),
    C("kubectl expose deployment my-nginx --port=80 --target-port=80"),
    C("kubectl get services"),
    C("kubectl describe service my-nginx"),
    C("kubectl get endpoints my-nginx"),
    BL(),
    B("Test from inside the cluster:"),
    C("kubectl run test-pod --image=busybox --rm -it --restart=Never \\"),
    C("  -- wget -O- my-nginx:80"),
    BL(),
    B("Clean up:"),
    C("kubectl delete service my-nginx && kubectl delete deployment my-nginx"),
])

write(S[226], [
    KC_NOTE("day2-04-k8s-services  →  Lab 17 (NodePort)"),
    BL(),
    B("Expose a Deployment externally via NodePort:"),
    C("kubectl create deployment my-nginx --image=nginx --replicas=2"),
    C("kubectl expose deployment my-nginx --type=NodePort --port=80 --target-port=80"),
    BL(),
    C("kubectl get service my-nginx   # note the NodePort (30xxx)"),
    C("kubectl get nodes -o wide      # get node IP"),
    BL(),
    B("Test from outside the cluster:"),
    C("curl http://$(kubectl get nodes -o jsonpath='{.items[0].status.addresses[0].address}'):$(kubectl get svc my-nginx -o jsonpath='{.spec.ports[0].nodePort}')"),
    BL(),
    B("Clean up:"),
    C("kubectl delete service my-nginx && kubectl delete deployment my-nginx"),
])

write(S[227], [
    KC_NOTE("day2-04-k8s-services  →  Lab 17 (Declarative YAML)"),
    BL(),
    B("service.yaml:"),
    C("apiVersion: v1"),
    C("kind: Service"),
    C("metadata:"),
    C("  name: my-app-svc"),
    C("spec:"),
    C("  type: NodePort"),
    C("  selector:"),
    C("    app: my-app"),
    C("  ports:"),
    C("    - port: 80"),
    C("      targetPort: 80"),
    C("      nodePort: 30080"),
])

write(S[228], [
    B("Apply and test:"),
    C("kubectl apply -f deployment.yaml"),
    C("kubectl apply -f service.yaml"),
    C("kubectl get services"),
    C("kubectl get endpoints my-app-svc"),
    C("curl http://localhost:30080"),
    BL(),
    B("Clean up:"),
    C("kubectl delete -f service.yaml && kubectl delete -f deployment.yaml"),
    BL(),
    T("Service Type Summary:"),
    T("  ClusterIP    — internal only (service-to-service communication)"),
    T("  NodePort     — external via node IP:port (dev / testing)"),
    T("  LoadBalancer — cloud load balancer (production)"),
    BL(),
    I("Key: Services give Pods a stable DNS name regardless of Pod IP changes."),
])

# ── LAB 18: K8s Volumes ──────────────────────────────────────────────────────
print("Updating Lab 18...")
write(S[250], [
    KC_NOTE("day2-05-k8s-storage-jobs  →  Lab 18 (emptyDir)"),
    BL(),
    B("emptyDir — shared storage between containers in the same Pod:"),
    C("# emptydir-pod.yaml  (two containers sharing /data)"),
    C("volumes:"),
    C("  - name: shared-data"),
    C("    emptyDir: {}"),
    C("containers:"),
    C("  - name: writer"),
    C("    command: [\"sh\",\"-c\",\"echo 'hello' > /data/msg.txt && sleep 3600\"]"),
    C("    volumeMounts:"),
    C("      - name: shared-data  mountPath: /data"),
    C("  - name: reader"),
    C("    command: [\"sh\",\"-c\",\"sleep 5 && cat /data/msg.txt && sleep 3600\"]"),
    C("    volumeMounts:"),
    C("      - name: shared-data  mountPath: /data"),
    BL(),
    C("kubectl apply -f emptydir-pod.yaml"),
    C("kubectl exec shared-pod -c reader -- cat /data/msg.txt"),
    C("kubectl logs shared-pod -c reader"),
    C("kubectl delete pod shared-pod"),
    BL(),
    I("emptyDir is deleted when the Pod is removed — use for temporary shared data."),
])

write(S[251], [
    KC_NOTE("day2-05-k8s-storage-jobs  →  Lab 18 (PV / PVC)"),
    BL(),
    B("PersistentVolume (pv.yaml):"),
    C("apiVersion: v1"),
    C("kind: PersistentVolume"),
    C("metadata:"),
    C("  name: my-pv"),
    C("spec:"),
    C("  capacity:"),
    C("    storage: 1Gi"),
    C("  accessModes: [ReadWriteOnce]"),
    C("  hostPath:"),
    C("    path: /tmp/k8s-data"),
    BL(),
    B("PersistentVolumeClaim (pvc.yaml):"),
    C("apiVersion: v1"),
    C("kind: PersistentVolumeClaim"),
    C("metadata:"),
    C("  name: my-pvc"),
    C("spec:"),
    C("  accessModes: [ReadWriteOnce]"),
    C("  resources:"),
    C("    requests:"),
    C("      storage: 500Mi"),
    BL(),
    C("kubectl apply -f pv.yaml && kubectl apply -f pvc.yaml"),
    C("kubectl get pv && kubectl get pvc   # both show  Bound"),
])

write(S[252], [
    B("Pod mounting the PVC:"),
    C("volumes:"),
    C("  - name: my-storage"),
    C("    persistentVolumeClaim:"),
    C("      claimName: my-pvc"),
    C("containers:"),
    C("  - name: app  image: busybox"),
    C("    command: [\"sh\",\"-c\",\"echo 'persistent data' > /data/file.txt && sleep 3600\"]"),
    C("    volumeMounts:"),
    C("      - name: my-storage  mountPath: /data"),
    BL(),
    C("kubectl apply -f pod-with-pvc.yaml"),
    C("kubectl exec pvc-pod -- cat /data/file.txt"),
    BL(),
    B("Prove data survives Pod deletion:"),
    C("kubectl delete pod pvc-pod"),
    C("kubectl apply -f pod-with-pvc.yaml"),
    C("kubectl exec pvc-pod -- cat /data/file.txt   # data still there!"),
    BL(),
    C("kubectl delete pod pvc-pod && kubectl delete pvc my-pvc && kubectl delete pv my-pv"),
])

# ── LAB 19: Jobs & CronJobs ──────────────────────────────────────────────────
print("Updating Lab 19...")
write(S[271], [
    KC_NOTE("day2-05-k8s-storage-jobs  →  Lab 19 (Jobs)"),
    BL(),
    B("Job — Imperative:"),
    C('kubectl create job my-job --image=busybox -- echo "Hello from Job!"'),
    C("kubectl get jobs && kubectl get pods"),
    C("kubectl logs job/my-job"),
    C("kubectl delete job my-job"),
    BL(),
    B("Job YAML — 3 completions, 2 in parallel:"),
    C("spec:"),
    C("  completions: 3"),
    C("  parallelism: 2"),
    C("  template:"),
    C("    spec:"),
    C("      containers:"),
    C("        - name: counter  image: busybox"),
    C("          command: [\"sh\",\"-c\",\"echo Processing && sleep 3 && echo Done\"]"),
    C("      restartPolicy: Never"),
    BL(),
    C("kubectl apply -f job.yaml"),
    C("kubectl get pods -w"),
    C("kubectl logs job/countdown"),
    C("kubectl delete -f job.yaml"),
    BL(),
    I("completions: 3 = 3 successful runs total.  parallelism: 2 = 2 Pods at a time."),
])

write(S[272], [
    KC_NOTE("day2-05-k8s-storage-jobs  →  Lab 19 (CronJobs)"),
    BL(),
    B("CronJob — Imperative (every minute):"),
    C("kubectl create cronjob my-cron \\"),
    C('  --image=busybox --schedule="*/1 * * * *" \\'),
    C('  -- echo "Hello from CronJob!"'),
    C("kubectl get cronjobs"),
    C("# Wait 1-2 minutes, then:"),
    C("kubectl get jobs"),
    C("kubectl delete cronjob my-cron"),
    BL(),
    B("CronJob YAML (every 2 minutes):"),
    C("spec:"),
    C('  schedule: "*/2 * * * *"'),
    C("  jobTemplate:"),
    C("    spec:"),
    C("      template:"),
    C("        spec:"),
    C("          containers:"),
    C("            - name: backup  image: busybox"),
    C("              command: [\"sh\",\"-c\",\"echo Backup at $(date) && sleep 3\"]"),
    C("          restartPolicy: Never"),
])

write(S[273], [
    B("Apply and monitor CronJob:"),
    C("kubectl apply -f cronjob.yaml"),
    C("kubectl get cronjobs"),
    C("# Wait 2-4 minutes, then:"),
    C("kubectl get jobs"),
    C("kubectl delete -f cronjob.yaml"),
    BL(),
    T("Cron schedule format:"),
    T("  minute   hour   day-of-month   month   day-of-week"),
    C('  "*/1 * * * *"   →  every minute'),
    C('  "0 9 * * 1-5"   →  9:00 AM every weekday'),
    C('  "0 0 * * 0"     →  midnight every Sunday'),
    BL(),
    I("Key: CronJob  →  creates Jobs on a schedule.  Job  →  runs Pods to completion."),
])

# ── Apply lab header bars ─────────────────────────────────────────────────────
print("Adding lab header bars...")
for idx, label in LAB_BARS.items():
    if idx < len(S):
        add_lab_bar(S[idx], label)

# ── Save ─────────────────────────────────────────────────────────────────────
prs.save(DST)
print(f"\nSaved: {DST}")
print("\nLab slides updated:")
labs = [
    "Lab 1-2 (Docker basics: run, cp, ps)",
    "Lab 3-4 (Dockerfile, Flask app)",
    "Lab 5   (Named volumes + Bind mounts)",
    "Lab 6-7 (Networking + Port mapping)",
    "Lab 8   (Environment variables)",
    "Lab 9   (Reference only — banner added)",
    "Lab 10-12 (Docker Compose)",
    "Lab 13-14 (K8s Pods + Namespaces)",
    "Lab 15-16 (Deployments + Rollouts)",
    "Lab 17  (Services: ClusterIP, NodePort)",
    "Lab 18  (emptyDir + PV/PVC)",
    "Lab 19  (Jobs + CronJobs)",
]
for l in labs:
    print(f"  ✓ {l}")
